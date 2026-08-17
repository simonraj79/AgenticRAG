"""Open the file before calling it a handout.

**The measurement this module exists for**, taken 2026-08-17 against this venv's
python-pptx 1.0.2 and re-verified in a second session:

    Presentation() with ZERO slides, prs.save(...)   ->  27,387 bytes, starts PK
    open("deck.pptx","wb").write(b"PK\\x03\\x04 junk")  ->  28 bytes, starts PK

Both are `ok=True` runs. Both harvest. Both clear every assertion in the
repository -- `agentic_check.py`'s `status == "ready" and byte_size > 0` and
`sandbox_check.py` case 3's `PK` plus `>= 10_000` bytes -- and both become
downloadable handouts. The second raises `zipfile.BadZipFile` in PowerPoint.
Nothing between the model's `prs.save()` and the user's Downloads folder had
ever opened the bytes.

**This is the third rung of a ladder `jobs._problem` already has two rungs of**,
and it adds no mechanism:

    the process failed        -> _problem branch 1
    the artefact is ABSENT    -> _problem branch 2   (`loop.md` T2's own example)
    the artefact is PRESENT
      but is not the thing    -> _problem branch 3   (this module)

**NOTHING HERE MAY RAISE.** `sandbox.run` never raises, `run_handout_job` never
raises, and `_settle` never raises; this inherits that contract because it is
called from inside all three. A validator that raises takes the whole handout
job with it -- and that has happened in this repo once already, inside
`_static_refusal`, the one function whose entire job was to prevent it: it
caught `SyntaxError` and `ValueError` and a `MemoryError` escaped
(`08-streaming-and-followups.md:163-194`). So every reader here catches
`Exception`, and the top-level dispatcher catches again behind them.

**Strictness leans permissive, and that is `loop.md` T3 rather than taste.** The
string returned here feeds a RETRY, not a refusal. A false positive costs one
generation call and one subprocess; a false negative ships a file that does not
open. The two costs are nowhere near symmetric, so every threshold is set at the
forgiving end and every unexpected failure inside this module resolves to "the
artefact is fine" rather than "the artefact is bad". The precedent for getting
that backwards is `refusal_pass = 0/2`, where a measurement punished correct
behaviour and the scorecard then advised deleting the behaviour.

**`fit_text` is never imported.** It is python-pptx's only text-fitting API and
the obvious way to answer "does this bullet fit"; `pptx/text/fonts.py:41-50`
returns font directories for darwin and win32 and otherwise raises
`OSError("unsupported operating system")`. It works on a Windows dev box and
kills every deck on Render, with a message that says nothing about fonts.
Character count is not a good proxy -- it is the only honest one available, and
`scripts/deck_check.py` case 14 asserts the symbol appears nowhere under
`app/handouts`.

**Where this does NOT run.** Not in the sandbox child: validation reads bytes
the parent already holds, after the child has exited, so it needs no entry in
`ALLOWED_IMPORTS`, is not subject to `RLIMIT_AS` (absent on Windows anyway), and
cannot be lied to by a hostile program about its own output. And not in
`_harvest`: `sandbox.py`'s contract is the sandbox and its controls, while
artefact *quality* is a handout concern.

**TWO CALLERS, TWO WAYS IN, ONE SET OF READERS.** `jobs._problem` has a `Recipe`
and passes it. `tools/interpreter._run` has no recipe and never will -- a deck
asked for in the chat is decided by the MODEL calling `run_python`, so all that
exists is the filename it chose and the MIME type `HARVEST_MIME` stamped on it.
`check` therefore takes either, and the only thing that differs between them is
*which filename the message quotes*: the recipe path names the file the prompt
spelled out, and the tool path names the file the model actually wrote. Quoting
`deck.pptx` at a program that saved `ka-band.pptx` is a refusal it cannot act
on, and a wasted step out of a bounded budget (`loop.md` section 4).

**`KIND_BY_MIME` is deliberately NARROWER than `ask._handout_kind`, and this is
the one place the two must not be unified.** That function answers "how should
the panel RENDER this", so it buckets `image/svg+xml` with PNG, `application/
json` with CSV and `text/plain` with markdown -- all correct for a thumbnail and
all wrong here, because the readers below are format-exact. `_check_chart` wants
the PNG signature and would reject a perfectly good SVG; `_check_table` runs
`csv.reader` and would call a single-line JSON document "no data"; `_check_sheet`
wants a markdown heading a `.txt` file has no reason to carry. Each is a false
positive, and on the tool path a false positive costs a step of the turn the
user is waiting on. **A MIME type is validated only when a reader exists that
reads exactly that format**; everything else passes through unread, which is
today's behaviour and the permissive direction this module is built to lean in.

**`outline` IS THE SAME WALK, KEPT INSTEAD OF DISCARDED.** `_check_deck` already
opens the file and visits every slide in order to count them; the only reason a
user has to download a `.pptx` and start PowerPoint to learn that their deck is
empty is that the walk threw its titles away. So the preview costs one more pass
over bytes already in memory and no column, no route and no migration --
`handouts.preview_text` exists, is on `HandoutDetail`, and is already rendered
inside `HandoutCard`'s disclosure.

Two properties of `outline` that are NOT `check`'s and are easy to lose:

- **It is not gated on `handout_validate_artifacts`.** That flag exists so the
  regression assertion can be executed -- with it off, `_problem` returns exactly
  what it returned before this module existed. The preview is a product feature
  and the flag is a regression switch, so wiring one to the other would delete
  every preview the moment somebody reproduced a baseline. `deck_check.py` case
  55 sets the flag explicitly and asserts the outline survives.
- **It is a PREVIEW, so it is read by a person, not by a model.** `check`'s
  strings are written for the retry turn; this one goes on a card. That is why it
  is capped, why every title is ASCII-folded, and why an unopenable file gets
  `None` here rather than a second copy of the verdict `check` already produced.
"""

from __future__ import annotations

import csv
import io
import logging
import unicodedata

from app.config import settings
from app.handouts.recipes import Recipe
from app.tools import sandbox

log = logging.getLogger("uvicorn.error")

# The prompt text python-pptx's default template carries in an untouched body
# placeholder. A model that copied the sample deck rather than writing one
# produces slides whose only text is this.
#
# Matched case-insensitively on the whole stripped paragraph, never as a
# substring: a slide legitimately explaining how to edit a master would contain
# the phrase inside a real sentence, and firing on that is the false positive
# this module is built to avoid.
PLACEHOLDER_PROMPTS = frozenset(
    {
        "click to edit master text styles",
        "click to edit master title style",
        "second level",
        "third level",
        "fourth level",
        "fifth level",
        "lorem ipsum",
    }
)

# The PNG magic number, all eight bytes. `_harvest` already restricts suffixes,
# so this catches a model that wrote a JPEG or a text file under a .png name.
PNG_SIGNATURE = b"\x89PNG\r\n\x1a\n"

# What the tool path has instead of a `Recipe`. See the module docstring for why
# this is narrower than `ask._handout_kind` and must stay that way: the three
# MIME types missing from it -- `image/svg+xml`, `application/json`,
# `text/plain` -- each have a reader that would reject a correct file.
KIND_BY_MIME = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "deck",
    "image/png": "chart",
    "text/csv": "table",
    "text/markdown": "sheet",
}


def kind_for_mime(mime_type: str) -> str | None:
    """The artefact kind a MIME type may be validated as, or None for "leave it".

    None is the common answer and the safe one. It means no reader here can read
    that format, so there is no evidence to gather and the artefact passes.
    """
    return KIND_BY_MIME.get((mime_type or "").strip().lower())


# --------------------------------------------------------------------------
# Entry point
# --------------------------------------------------------------------------


def check(spec: Recipe | str, artifact: sandbox.SandboxArtifact) -> str | None:
    """A model-actionable problem string, or None if the artefact is usable.

    Pure. Reads bytes already in memory. Never raises -- see the module
    docstring for the one time that promise was broken here before.

    `spec` is a `Recipe` from the panel path or an artefact KIND from the tool
    path, which has no recipe to offer (module docstring, and
    `12-robust-handouts/06-tool-path-parity.md` question 2). The two differ in
    one place only, and it is the filename the message quotes.

    **The string is written for the model, in the register the existing
    `_problem` branches already speak**: name the defect, name the fix, and be
    specific enough that the next attempt is a correction rather than a re-roll.
    A refusal the model cannot act on wastes a step (`loop.md` section 4).
    """
    if isinstance(spec, str):
        kind = spec
        # THE MODEL CHOSE THIS NAME, so it is the name the message uses. There is
        # no `output_filename` on this path and inventing one would be worse than
        # having none: `deck.pptx` quoted at a program that saved `ka-band.pptx`
        # reads as a message about a different file.
        name = artifact.filename
    else:
        # Unchanged. `output_filename` is the name the prompt spelled out, and
        # `_primary_artifact` may have accepted a file saved under another one --
        # the repair message should name the file the prompt asked for.
        kind, name = spec.key, spec.output_filename

    reader = _READERS.get(kind)
    if reader is None:
        # A kind with no validator is not a validation failure. On the recipe
        # path this is only reachable after a new recipe is added -- and the
        # correct behaviour then is today's behaviour, not a handout that fails
        # because nobody wrote a reader yet. On the tool path it is the ORDINARY
        # case: `HARVEST_MIME` admits seven suffixes and only four of them have
        # a reader that reads exactly that format.
        return None

    try:
        return reader(artifact.content, name)
    except Exception:  # noqa: BLE001
        # THE LAST RESORT, AND IT RETURNS None ON PURPOSE.
        #
        # Reaching here means a bug in this module rather than a bad artefact --
        # the readers below already catch the failures that ARE evidence and
        # turn them into strings. A bug in the walker has produced no evidence
        # about the file, so the permissive answer is the honest one; returning
        # a string instead would fail every handout of that kind, twice each,
        # for as long as the bug lived.
        log.exception("Handout validation crashed for kind %s", kind)
        return None


# --------------------------------------------------------------------------
# Per-recipe readers
# --------------------------------------------------------------------------


def _check_deck(data: bytes, name: str) -> str | None:
    """The recipe with a measured failure. Everything else here is cheap parity.

    `name` rather than a `Recipe`, because the tool path has no recipe and the
    only thing the readers ever wanted off one was the filename to quote. See
    `check`.
    """
    try:
        from pptx import Presentation
    except Exception:  # noqa: BLE001
        # python-pptx is a declared dependency and its absence would already
        # have failed the sandbox run that produced these bytes. Returning None
        # rather than a verdict keeps the asymmetry pointing the right way: a
        # missing library is evidence about this process, not about the file.
        log.warning("python-pptx unavailable; skipping deck validation")
        return None

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        # CATCH BROADLY, and the reason is specific rather than defensive: a
        # .pptx is a zip, so junk bytes raise `zipfile.BadZipFile` from the
        # STANDARD LIBRARY one layer below python-pptx (verified 2026-08-17,
        # 28-byte probe), while a valid zip that is not a presentation raises
        # from python-pptx's own package reader. Naming the exceptions we know
        # about would let the ones we do not escape into `run_handout_job`.
        #
        # Only the exception's TYPE NAME goes into the message. Its `str()` can
        # carry a filesystem path or non-ASCII, and this string is read back on
        # a Windows console (`deck_check.py` case 27 asserts ASCII).
        return (
            f"The file {name} is not a readable PowerPoint file -- opening it "
            f"raised {type(exc).__name__}. Rewrite the program so it builds the "
            "deck with python-pptx and saves it with `prs.save('" + name + "')`; "
            "do not create the file with `open()` or write bytes into it by hand."
        )

    # From here on the file opened, so a failure in the walk below is a bug in
    # this module and must not be reported as "it did not open". The walk is
    # written not to raise -- `shapes.title` is None rather than "" on a Blank
    # layout, which is the exact shape a naive `.title.text` walk dies on
    # (`deck_check.py` case 16) -- and `check`'s outer catch is behind it.
    slides = list(prs.slides)
    floor = settings.handout_deck_min_slides
    if len(slides) < floor:
        return (
            f"The file {name} opened but has {len(slides)} slides, and the "
            f"handout needs at least {floor}. Add slides with "
            "`prs.slides.add_slide(prs.slide_layouts[1])` and put a title and "
            "bullets drawn from the MATERIAL on each one before `prs.save()`."
        )

    limit = settings.handout_deck_max_bullet_chars
    untitled: list[int] = []
    placeholder: list[int] = []
    overlong: list[tuple[int, int]] = []

    for number, slide in enumerate(slides, start=1):
        title = _slide_title(slide)
        body = _slide_paragraphs(slide, include_title=False)

        if not title:
            untitled.append(number)

        # Only when the slide has body text AND every line of it is boilerplate.
        # A slide that is deliberately title-only (a section divider) has no
        # body text at all and must not fire -- see the permissive rule above.
        if body and all(line.lower() in PLACEHOLDER_PROMPTS for line in body):
            placeholder.append(number)

        longest = max((len(line) for line in _slide_paragraphs(slide)), default=0)
        if longest > limit:
            overlong.append((number, longest))

    # PARAGRAPHS, not runs. python-pptx splits a paragraph into runs wherever
    # formatting changes, so a single bold word makes three short runs out of one
    # long bullet and a per-run measure would silently stop seeing the overflow.
    # The user reads the paragraph.

    findings: list[str] = []
    fixes: list[str] = []
    # The verb agrees with the count. Trivial, and it is in here because the
    # first draft shipped `slide 4 still carry template placeholder text` and
    # nothing in the harness could see it -- every case asserted on substrings
    # and the message was ASCII, non-empty and imperative. Found by printing one
    # real verdict and reading it, which is `build.md`'s last verification step
    # arriving in a module small enough to have felt exempt from it.
    if untitled:
        findings.append(f"{_slide_list(untitled)} {_verb(untitled)} no title text")
        fixes.append("Set `slide.shapes.title.text` on every slide.")
    if placeholder:
        findings.append(
            f"{_slide_list(placeholder)} still "
            f"{_verb(placeholder, 'carries', 'carry')} template placeholder text"
        )
        fixes.append(
            "Replace the placeholder text with real content from the MATERIAL."
        )
    if overlong:
        worst = max(overlong, key=lambda item: item[1])
        findings.append(
            f"slide {worst[0]} has a bullet of {worst[1]} characters, past the "
            f"{limit}-character limit"
        )
        fixes.append(
            f"Shorten every bullet to under {limit} characters, splitting a long "
            "point across two slides rather than shrinking the text."
        )

    if not findings:
        return None

    return (
        f"The file {name} opened with {len(slides)} slides, but "
        + "; ".join(findings)
        + ". "
        + " ".join(fixes)
        + f" Then save it again with `prs.save('{name}')`."
    )


def _check_chart(data: bytes, name: str) -> str | None:
    if not data.startswith(PNG_SIGNATURE):
        return (
            f"The file {name} is not a PNG -- its first bytes are not the PNG "
            "signature. Save the figure with `plt.savefig('" + name + "')` after "
            "plotting; do not write the file by hand."
        )

    try:
        from PIL import Image
    except Exception:  # noqa: BLE001
        log.warning("Pillow unavailable; skipping chart validation")
        return None

    try:
        with Image.open(io.BytesIO(data)) as image:
            width, height = image.size
    except Exception as exc:  # noqa: BLE001
        return (
            f"The file {name} has a PNG header but did not open as an image -- "
            f"reading it raised {type(exc).__name__}. Save the figure with "
            f"`plt.savefig('{name}')` after plotting, and let matplotlib write "
            "the whole file."
        )

    # `<= 1` rather than a real minimum size. A 200x150 chart is small and legible;
    # a 1x1 is the shape a figure written before anything was plotted takes.
    if width <= 1 or height <= 1:
        return (
            f"The image {name} is {width}x{height} pixels, which is blank rather "
            "than a chart. Plot the figures from the MATERIAL and give the figure "
            f"a real size with `plt.figure(figsize=(10, 6))` before `plt.savefig('{name}')`."
        )
    return None


def _check_table(data: bytes, name: str) -> str | None:
    text = data.decode("utf-8", errors="replace")

    try:
        rows = [row for row in csv.reader(io.StringIO(text)) if any(cell.strip() for cell in row)]
    except Exception as exc:  # noqa: BLE001
        return (
            f"The file {name} did not parse as CSV -- reading it raised "
            f"{type(exc).__name__}. Write it with `csv.writer`, one row per "
            "`writerow` call, rather than by formatting strings."
        )

    # Two rows: a header and at least one row of data. A header-only CSV is the
    # `savefig`-was-forgotten failure wearing a different extension.
    if len(rows) < 2:
        return (
            f"The file {name} has only {len(rows)} non-empty row(s), so it "
            "carries no data. Write a header row and then one row per fact drawn "
            "from the MATERIAL, each with its source."
        )
    return None


def _check_sheet(data: bytes, name: str) -> str | None:
    """The `sheet` recipe does not use the sandbox, so `_problem` never reaches
    this. It is here because `06-test-plan.md:183` promised all four artefacts
    were opened, because it costs four lines -- and, since the tool path landed,
    because a `.md` written by `run_python` DOES reach it."""
    text = data.decode("utf-8", errors="replace")

    if not text.strip():
        return (
            "The study sheet is empty. Write it as markdown, starting with a "
            "`# ` heading and covering the MATERIAL in sections."
        )
    if not any(line.lstrip().startswith("#") for line in text.splitlines()):
        return (
            "The study sheet has no markdown heading, so it will render as one "
            "wall of text. Start it with a `# ` heading and use `## ` for each "
            "section."
        )
    return None


_READERS = {
    "deck": _check_deck,
    "chart": _check_chart,
    "table": _check_table,
    "sheet": _check_sheet,
}


# --------------------------------------------------------------------------
# The preview
# --------------------------------------------------------------------------

# `preview_text` is `Text` and unbounded, so this cap is not a storage limit --
# it is about what the field costs and what it is for. It is returned on the
# detail fetch a card makes when its disclosure is opened, and a forty-slide
# outline is prompt-sized: roughly 2,300 characters, for a panel row.
#
# 2,000 fits about thirty slides, which is past anything `DECK_PROMPT` asks for
# (five to eight) by a factor of four. Whatever does not fit is COUNTED rather
# than dropped, because a preview that silently stops at slide 34 reads as a
# deck that stops at slide 34 -- the same shape as a truncated program, wrong in
# a way that looks complete.
PREVIEW_MAX_CHARS = 2_000

# Per title. A model-written slide heading is meant to be at most eight words
# (`DECK_PROMPT`), so this only ever fires on a title that has gone wrong -- and
# when one has, it must not be allowed to eat the whole preview and push the
# other thirty-nine slides into the "more slides" tail.
PREVIEW_TITLE_CHARS = 120

# What a slide with no title text is called in the list. The number still has to
# be there: dropping the row would renumber every slide after it, so the preview
# would disagree with the deck about which slide is which -- and "slide 4 is
# untitled" is exactly the finding a preview is worth having for.
UNTITLED = "(untitled)"

# Punctuation a model reaches for that HAS an honest ASCII spelling. Everything
# else non-ASCII is dropped by the fold below, which is right for an emoji and
# for a section sign and wrong for these -- deleting an em-dash joins two words
# ("Ka-bandbudget"), and deleting an apostrophe misspells a word.
#
# Applied BEFORE the NFKD fold, because NFKD leaves all of these alone: they are
# not decompositions of an ASCII character, they are different characters.
_ASCII_SUBSTITUTIONS = {
    0x2014: "--",  # em dash
    0x2013: "-",  # en dash
    0x2212: "-",  # minus sign
    0x2018: "'",  # left single quote
    0x2019: "'",  # right single quote / apostrophe
    0x201C: '"',  # left double quote
    0x201D: '"',  # right double quote
    0x2026: "...",  # horizontal ellipsis
    0x2022: "-",  # bullet
    0x00B7: "-",  # middle dot
    0x00A0: " ",  # no-break space
    0x00D7: "x",  # multiplication sign
}

# CONTROL CHARACTERS ARE ASCII, AND THAT IS THE TRAP. `"\x00".isascii()` is
# True, so an `isascii()` assertion -- the obvious one, and the one
# `deck_check.py` case 54 makes -- passes a string carrying a NUL. Found by
# probing `outline` with binary bytes under a `.csv` name, which returned
# `"0 data rows\n\x00\x01"`.
#
# It is not cosmetic. **Postgres rejects NUL in a text column**, so that string
# reaches `handout.preview_text`, the commit raises, and `run_handout_job`
# fails a handout whose artefact was fine -- the exact outcome this module's
# never-raises contract exists to prevent, arriving through the one field that
# was supposed to be free.
_CONTROL_CHARS = {code: None for code in list(range(0x20)) + [0x7F]}


def outline(spec: Recipe | str, artifact: sandbox.SandboxArtifact) -> str | None:
    """A short human-readable summary of what the artefact contains, or None.

    Pure, never raises, ASCII-safe -- `check`'s contract, for the same reason:
    this is called from inside `run_handout_job`, which never raises, and a
    handout must not be lost because its PREVIEW failed. A preview is the least
    important thing in the pipeline and it sits inside the most important one.

    `spec` is a `Recipe` or an artefact KIND, exactly as `check` takes it
    (module docstring). Unlike `check`, this one never quotes the filename --
    the user is looking at the row the filename is already printed on.

    **Two of the four recipes get None, and both are deliberate.**

    - `chart` already renders its own PNG as a thumbnail in the card
      (`HandoutCard.tsx:198-207`). A text preview beside a picture of the thing
      is noise.
    - `sheet` writes the study-sheet markdown itself into `preview_text`
      (`jobs._run_direct_recipe`). It is the only recipe whose preview IS the
      product, and replacing it with a summary would delete a feature rather
      than add one. This function is not called on that path at all, and it
      still answers None, because a dispatch table that gained a fourth entry
      "for symmetry" is how that would stop being true.

    Called after `check` has returned None, so the artefact is known-usable --
    but this does not ASSUME that: `_outline_deck` re-opens the file and answers
    None if it will not open, because the tool path and a future caller may not
    have asked `check` anything.
    """
    kind = spec if isinstance(spec, str) else spec.key

    writer = _OUTLINERS.get(kind)
    if writer is None:
        return None

    try:
        text = writer(artifact.content)
    except Exception:  # noqa: BLE001
        # Same last resort as `check`, and the same reasoning: reaching here is
        # a bug in this module rather than a fact about the file. `None` means
        # "no preview", which is what every handout had before this feature, so
        # a bug here costs the preview and nothing else.
        log.exception("Handout preview crashed for kind %s", kind)
        return None

    if not isinstance(text, str) or not text.strip():
        return None
    return text


def _outline_deck(data: bytes) -> str | None:
    """The count, then the slide titles, numbered."""
    try:
        from pptx import Presentation
    except Exception:  # noqa: BLE001
        log.warning("python-pptx unavailable; skipping deck preview")
        return None

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception:  # noqa: BLE001
        # The junk case. `check` has already said what is wrong with this file,
        # in words the model can act on; repeating it here would print a repair
        # instruction to a user as though it were the deck's contents. A row
        # that cannot be previewed shows no preview.
        return None

    # A zero-slide deck reaches here whenever validation is off, which is the
    # single most useful thing this function says: "0 slides" on the card, with
    # no download and no PowerPoint. It is the measured failure of section 1.1
    # made visible in one line.
    titles = [
        _clip(_ascii_safe(_slide_title(slide)), PREVIEW_TITLE_CHARS) or UNTITLED
        for slide in prs.slides
    ]
    return _numbered(titles, "slide")


def _outline_table(data: bytes) -> str | None:
    """The row count and the header, which is what a CSV's shape actually is.

    Not the rows themselves: a table handout is a file to open in a spreadsheet,
    and a preview that reproduces it is a worse spreadsheet. The header answers
    the question the card is asked -- "what columns did it choose" -- and the
    count answers "did it find anything".
    """
    text = data.decode("utf-8", errors="replace")
    rows = [
        row
        for row in csv.reader(io.StringIO(text))
        if any(cell.strip() for cell in row)
    ]
    if not rows:
        return None

    header = " | ".join(_ascii_safe(cell) for cell in rows[0])
    # `len(rows) - 1`, and the noun says so. "3 rows" over a header plus two
    # data rows is the kind of off-by-one a reader cannot check without opening
    # the file, which is the whole thing a preview exists to avoid.
    return _clip(
        f"{_plural(len(rows) - 1, 'data row')}\n{header}", PREVIEW_MAX_CHARS
    )


# `chart` and `sheet` are absent, not mapped to a function returning None. See
# `outline`'s docstring: the absence is the statement.
_OUTLINERS = {
    "deck": _outline_deck,
    "table": _outline_table,
}


# --------------------------------------------------------------------------
# Preview formatting
# --------------------------------------------------------------------------


def _numbered(items: list[str], noun: str) -> str:
    """`"6 slides"` and then `"1. ..."`, capped, with the remainder counted."""
    header = _plural(len(items), noun)
    lines = [header]
    used = len(header)
    # Reserved against the LONGEST tail that could be needed, so the cap holds
    # whichever line the truncation lands on. Computing it from the remainder at
    # the break point would be exact and would let the total overshoot by the
    # difference between "9 more" and "39 more".
    reserve = len(f"\n\n... ({len(items)} more {noun}s)")

    for index, item in enumerate(items, start=1):
        line = f"{index}. {item}"
        if used + 1 + len(line) + reserve > PREVIEW_MAX_CHARS:
            # THE BLANK LINE IS NOT SPACING. `HandoutCard` renders this field
            # through `Markdown`, where `1. ...` is an ordered list -- so a tail
            # placed directly under the last item is a LAZY CONTINUATION of it,
            # and the card reads "slide 33: ... (7 more slides)" as though that
            # were the slide's title. Found by rendering one, not by an
            # assertion: every string check passes either way.
            lines.append("")
            lines.append(f"... ({_plural(len(items) - index + 1, 'more ' + noun)})")
            break
        lines.append(line)
        used += 1 + len(line)

    return "\n".join(lines)


def _plural(count: int, noun: str) -> str:
    """`"1 slide"` / `"6 slides"`. `_check_deck`'s verb agreement, for nouns.

    Here for the reason case 28 of `deck_check.py` exists: the first build of
    `_check_deck` shipped `slide 4 still carry ...`, which is ASCII, non-empty
    and imperative, and no assertion in the harness could see it. A preview is
    read by a person, so "1 slides" on a card is the same defect on the same
    day it would be noticed.
    """
    return f"{count} {noun}" if count == 1 else f"{count} {noun}s"


def _clip(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    return text[: max(0, limit - 3)].rstrip() + "..."


def _ascii_safe(text: str) -> str:
    """Model-written text, made safe for a Windows console and a `String` field.

    **This is the class of text that has broken three throwaway scripts in this
    repo** -- an emoji out of `agent_templates.icon`, a section sign, and a
    box-drawing character copied out of a Markdown tree. `preview_text` is read
    back by `deck_check.py`, by `agentic_check.py` and by whatever one-off
    script is written next, and the Windows console codepage turns any of them
    into `UnicodeEncodeError` at the `print`, nowhere near the cause.

    Three steps, and the order is the point:

    1. **Collapse whitespace first.** A python-pptx title with two paragraphs
       comes back with a `\\n` in it, and one newline inside a numbered list
       makes the preview claim a slide that does not exist.
    2. **Substitute the punctuation that has an ASCII spelling**, before the
       fold, because NFKD does not touch it -- an em-dash is a character, not a
       decomposition.
    3. **Fold, then drop.** NFKD turns an accented letter into a letter plus a
       combining mark, so `"Ka-band\u00e9"` survives as `"Ka-bande"` rather than
       losing the vowel; everything with no ASCII form at all -- emoji, a
       section sign, CJK -- is dropped, and the whitespace is collapsed a second
       time because dropping a character between two spaces leaves two.
    4. **Drop the control characters, which the first three steps do not.**
       `str.split()` eats the whitespace ones and `isascii()` is True for all of
       them, so a NUL out of a binary file read as a CSV survives every check
       here and then fails the Postgres write. See `_CONTROL_CHARS`.

    `provisional_filename` solves the same problem by slugging to `[a-z0-9-]`,
    which is right for a filename and would be unreadable for a slide title.
    This is the same care, one register up.
    """
    flattened = " ".join(str(text).split())
    swapped = flattened.translate(_ASCII_SUBSTITUTIONS)
    folded = unicodedata.normalize("NFKD", swapped)
    plain = folded.encode("ascii", "ignore").decode("ascii")
    return " ".join(plain.translate(_CONTROL_CHARS).split())


# --------------------------------------------------------------------------
# Slide walking. Defensive by contract, not by habit -- see `check`.
# --------------------------------------------------------------------------


def _slide_title(slide) -> str:
    """The slide's title text, or "" -- including when there is no placeholder.

    `shapes.title` is **None** on a layout with no title placeholder (Blank),
    which is a different code path from a title whose text is empty, and is the
    one a naive `.title.text` walk raises on. `deck_check.py` case 16 pins it.
    """
    holder = getattr(slide.shapes, "title", None)
    if holder is None or not getattr(holder, "has_text_frame", False):
        return ""
    return holder.text_frame.text.strip()


def _slide_paragraphs(slide, *, include_title: bool = True) -> list[str]:
    """Every non-empty paragraph of text on the slide."""
    title_shape = getattr(slide.shapes, "title", None)
    out: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if not include_title and title_shape is not None and shape is title_shape:
            continue
        for paragraph in shape.text_frame.paragraphs:
            line = paragraph.text.strip()
            if line:
                out.append(line)
    return out


def _slide_list(numbers: list[int]) -> str:
    """`slide 3` / `slides 1 and 4` / `slides 1, 4 and 6`. Capped at three.

    The model needs enough to locate the defect, not an inventory -- and an
    unbounded list of slide numbers in a repair prompt is how a message stops
    being read.
    """
    shown = [str(n) for n in numbers[:3]]
    more = len(numbers) - len(shown)
    if len(shown) == 1:
        body = shown[0]
    else:
        body = ", ".join(shown[:-1]) + " and " + shown[-1]
    tail = f" (and {more} more)" if more > 0 else ""
    return f"slide{'' if len(numbers) == 1 else 's'} {body}{tail}"


def _verb(numbers: list[int], singular: str = "has", plural: str = "have") -> str:
    """Agree with the subject `_slide_list` just produced."""
    return singular if len(numbers) == 1 else plural
