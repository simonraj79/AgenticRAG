"""Layer 1 harness for `app/handouts/`. No DB, no API, no model, no subprocess.

WHY THIS FILE EXISTS.

Nothing under `scripts/` imported anything from `app.handouts` until 2026-08-17.
Ten pure functions had no coverage at all, and the module's own docstrings say
which branches are easy to leave out -- `_problem`'s second one in as many words.
Meanwhile the only assertions about a produced slide deck lived in
`agentic_check.py`, which needs a database, a live model, Pinecone, Cohere and
twenty minutes, and whose `--only` flag skips the HTTP block that holds them
entirely (`agentic_check.py:1836`). So the cheap checks were unrunnable and the
expensive ones were `byte_size > 0`.

Measured 2026-08-17 against this venv's python-pptx 1.0.2:

    Presentation() with zero slides  ->  27,387 bytes, starts with PK
    b"PK\\x03\\x04" + junk           ->  a `ready` handout that will not open

Both cleared every assertion in the repository. `06-test-plan.md:183` and
`04-handouts-panel.md:404` have both promised since they were written that a deck
is opened with python-pptx; neither was ever executed. That is `build.md`'s own
rule -- a criterion names a harness file and a case id or it is a wish -- failing
for the third document, and this file is the first half of the correction.

WHAT IS DELIBERATELY NOT HERE.

No validator. `app/handouts/validate.py` is feature 02, and the point of shipping
this file first is to have somewhere to watch it fail. Cases 11-13 assert what the
fixtures ARE, so that feature 02's cases can assert what a validator SAYS about
them.

THE FIXTURES ARE BUILT, NOT COMMITTED.

`12-robust-handouts/01-deck-harness-floor.md` said to commit four .pptx files.
Building them here instead: large binaries in git are permanent, a committed
empty.pptx would freeze one python-pptx version's output while the interesting
fact is what THIS library produces, and a reviewer can read `_thin_honest()` and
see a three-slide deck where a binary tells them nothing. The byte count is
asserted as a range rather than an equality for the same reason.

    backend/.venv/Scripts/python.exe scripts/deck_check.py

ASCII in print(), exit 1 on any failure -- `new features/06-test-plan.md`.
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "backend"))

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from app.handouts.jobs import (  # noqa: E402
    ATTEMPT_SEPARATOR,
    _join_attempts,
    _primary_artifact,
    _problem,
    _repair_message,
    _strip_fence,
)
from app.handouts.recipes import (  # noqa: E402
    RECIPES,
    Material,
    derive_title,
    provisional_filename,
    render,
)
from app.tools.sandbox import SandboxArtifact, SandboxResult  # noqa: E402

failures: list[str] = []

DECK = RECIPES["deck"]
CHART = RECIPES["chart"]


def check(name: str, condition: bool, detail: str = "") -> None:
    status = "[ok]  " if condition else "[FAIL]"
    print(f"{status} {name}" + (f" -- {detail}" if detail else ""))
    if not condition:
        failures.append(name)


def short(text: str, limit: int = 120) -> str:
    """ASCII-safe, single-line, bounded."""
    flat = " ".join(str(text).split())
    if len(flat) > limit:
        flat = flat[:limit] + "..."
    return ascii(flat)


# ---------------------------------------------------------------------------
# Builders. Hand-constructed inputs for functions that are pure but not simple.
# ---------------------------------------------------------------------------


def artifact(filename: str, content: bytes = b"x" * 32, mime: str = "") -> SandboxArtifact:
    return SandboxArtifact(filename=filename, mime_type=mime or "text/plain", content=content)


def ok_result(*artifacts: SandboxArtifact) -> SandboxResult:
    return SandboxResult(
        ok=True, exit_code=0, stdout="", stderr="", artifacts=list(artifacts)
    )


def failed_result(error: str, kind: str = "runtime") -> SandboxResult:
    return SandboxResult(
        ok=False,
        exit_code=1,
        stdout="",
        stderr="Traceback (most recent call last):\n  ZeroDivisionError",
        artifacts=[],
        error=error,
        error_kind=kind,
    )


# ---------------------------------------------------------------------------
# Fixture decks
# ---------------------------------------------------------------------------

# 240 is `handout_deck_max_bullet_chars`'s planned default (PLAN.md 3.1). The
# overflow fixture must clear it by enough that a later re-tuning from the
# measured distribution does not quietly make this fixture valid.
OVERFLOW_BULLET = (
    "The Ka-band downlink budget assumes a clear-sky margin that the material "
    "does not restate for rain-fade conditions, and the handover window is "
    "quoted only for the primary ground station, so the figure below should be "
    "read as an upper bound rather than a commitment, and the reader should "
    "check the modulation and coding scheme against the link table before "
    "using any of it for planning purposes at all, twice over. "
) * 2


def _save(prs) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def deck_empty() -> bytes:
    """Zero slides. A valid zip, a valid pptx, and a useless handout."""
    return _save(Presentation())


def deck_junk() -> bytes:
    """The 28-byte probe. Passes `agentic_check.py` S8's `byte_size > 0` today."""
    return b"PK\x03\x04 this is not a real pptx"


def _titled(prs, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullets[0]
    for extra in bullets[1:]:
        body.add_paragraph().text = extra
    return slide


def deck_thin_honest() -> bytes:
    """THE CONTROL. Three slides, titled, short bullets -- a correct honest shrink.

    `DECK_PROMPT` tells the model to use only what the material supports, so a
    thin corpus SHOULD produce a short deck. A validator that fires on this is
    the `refusal_pass = 0/2` defect again: a measurement punishing the behaviour
    the prompt exists to produce, and then recommending its removal.

    A DELETED validator passes every case that only checks bad decks are
    rejected. This is the only fixture that can tell a working detector from an
    absent one.
    """
    prs = Presentation()
    _titled(prs, "Ka-band downlink", ["Downlink runs at 26 GHz [comms.md]"])
    _titled(prs, "Link margin", ["Clear-sky margin is 3 dB [comms.md]"])
    _titled(prs, "Handover", ["Primary station hands over hourly [comms.md]"])
    return _save(prs)


def deck_overflow() -> bytes:
    """Six slides, one bullet far past any plausible placeholder."""
    prs = Presentation()
    for n in range(1, 6):
        _titled(prs, f"Subsystem {n}", [f"A short and entirely reasonable bullet {n}"])
    _titled(prs, "Caveats", [OVERFLOW_BULLET])
    return _save(prs)


def deck_untitled() -> bytes:
    """Two slides on a layout with no title placeholder at all.

    `slide_layouts[6]` is Blank. `shapes.title` is None there rather than empty,
    which is a different code path from a title whose text is "" and is the one
    a naive `.title.text` walk raises on.
    """
    prs = Presentation()
    for _ in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = "text but no title placeholder"
        box.text_frame.paragraphs[0].font.size = Pt(18)
    return _save(prs)


def open_deck(data: bytes):
    """(Presentation, None) or (None, reason). Never raises -- see sandbox_check."""
    try:
        return Presentation(io.BytesIO(data)), None
    except Exception as exc:  # noqa: BLE001
        return None, f"{type(exc).__name__}: {exc}"


def slide_titles(prs) -> list[str]:
    out = []
    for slide in prs.slides:
        holder = getattr(slide.shapes, "title", None)
        if holder is None or not getattr(holder, "has_text_frame", False):
            out.append("")
        else:
            out.append(holder.text_frame.text.strip())
    return out


def longest_run(prs) -> int:
    longest = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                longest = max(longest, len(para.text))
    return longest


print("=" * 74)
print("deck_check -- app/handouts pure functions, and what a .pptx actually is")
print("=" * 74)

# ---------------------------------------------------------------------------
# 1-10. The ten uncovered pure functions (01-deck-harness-floor.md, A7).
# ---------------------------------------------------------------------------
print("\n-- pure functions: ten with no coverage at all until now --")

# 1. _strip_fence: the ordinary case, and the language tag.
check(
    "1. _strip_fence unwraps a whole-reply fence and drops the language tag",
    _strip_fence('```python\nprs.save("deck.pptx")\n```') == 'prs.save("deck.pptx")',
    short(_strip_fence('```python\nprs.save("deck.pptx")\n```')),
)

# 2. _strip_fence: the guard that stops a study sheet being silently corrupted.
#    Its docstring calls this out; nothing tested it. A reply containing its own
#    fenced block has four markers, not two, and must come back untouched.
sheet = "# Notes\n\nRun this:\n\n```python\nx = 1\n```\n\nThat is all."
check(
    "2. _strip_fence leaves an inner fenced block alone (count != 2 guard)",
    _strip_fence(sheet) == sheet,
    "a study sheet's own code block must survive",
)

# 3. _primary_artifact: all three tiers, including the deliberate forgiveness.
tier1 = _primary_artifact(CHART, [artifact("notes.md"), artifact("chart.png")])
tier2 = _primary_artifact(CHART, [artifact("figure.png")])
tier3 = _primary_artifact(CHART, [artifact("data.csv")])
check(
    "3. _primary_artifact: exact name, then extension, then None",
    tier1 is not None
    and tier1.filename == "chart.png"
    and tier2 is not None
    and tier2.filename == "figure.png"
    and tier3 is None,
    f"tier1={tier1 and tier1.filename} tier2={tier2 and tier2.filename} tier3={tier3}",
)

# 4. _problem branch 1: the process failed.
crash = _problem(DECK, failed_result("The code raised an error: boom"), None)
check(
    "4. _problem branch 1 returns the sandbox's own error on a crash",
    crash is not None and "boom" in crash,
    short(crash),
)

# 5. _problem branch 2: THE one its docstring says is easy to leave out.
missed = _problem(DECK, ok_result(artifact("notes.md")), None)
check(
    "5. _problem branch 2 fires on ran-fine-but-wrote-no-file",
    missed is not None
    and "deck.pptx" in missed
    and "notes.md" in missed,
    short(missed),
)

# 6. _problem with an artefact present. THIS IS THE BASELINE feature 02's
#    regression assertion (PLAN.md 3.6 R-a) is measured against: it is None for
#    ANY bytes, including a zero-slide deck.
#
#    AMENDED BY FEATURE 02, and the amendment is the point rather than a repair.
#    This case was written to record pre-validation behaviour and it read the
#    AMBIENT flag, which now defaults to True -- so it asserted the flag-OFF
#    contract while running flag-ON, and feature 02 could not both satisfy its
#    own A5 (an empty deck IS a problem) and leave this green. Exactly one of
#    the two had to move. The baseline itself is untouched: the same pairs, the
#    same expected `None`, now measured under the condition its own detail line
#    always named. Case 25 asserts the whole table.
from app.config import settings as _settings  # noqa: E402

_prior = _settings.handout_validate_artifacts
try:
    _settings.handout_validate_artifacts = False
    present = _problem(DECK, ok_result(artifact("deck.pptx")), artifact("deck.pptx"))
    baseline_empty = _problem(
        DECK,
        ok_result(artifact("deck.pptx", deck_empty())),
        artifact("deck.pptx", deck_empty()),
    )
finally:
    _settings.handout_validate_artifacts = _prior
check(
    "6. _problem returns None whenever a file exists -- with validation OFF, even an empty deck",
    present is None and baseline_empty is None,
    "R-a baseline: feature 02 must keep this true with the flag OFF",
)

# 7. _join_attempts: one bare, two bannered.
one = _join_attempts([("print(1)", "first")])
two = _join_attempts([("print(1)", "first"), ("print(2)", "second")])
check(
    "7. _join_attempts stores a single attempt bare and two with banners",
    one == "print(1)"
    and two.count("ATTEMPT") == 2
    and "print(1)" in two
    and "print(2)" in two,
    f"one={short(one)} two_has_banners={two.count('ATTEMPT')}",
)

# 8. render: substitution, and the brace trap str.format would spring.
braced = Material(context="ctx {not_a_field}", conversation="")
rendered = render(DECK, brief="a brief with {braces}", material=braced)
check(
    "8. render substitutes without parsing braces in prompt, brief or material",
    "a brief with {braces}" in rendered and "ctx {not_a_field}" in rendered,
    "str.format here would raise KeyError at generation time",
)

# 9. Material: the refuse-rather-than-invent gate, both directions.
check(
    "9. Material.is_empty and conversation_block, both directions",
    Material().is_empty
    and not Material(context="x").is_empty
    and not Material(conversation="y").is_empty
    and Material(context="x").conversation_block == ""
    and "RECENT CONVERSATION" in Material(conversation="y").conversation_block,
    "an empty MATERIAL must refuse, not generate from parametric memory",
)

# 10. Filenames and titles, including the hostile brief 04-handouts-panel.md:405
#     names in prose and nothing ever ran.
hostile = provisional_filename(DECK, 'a"; rm -rf /')
long_title = derive_title("T" * 500)
check(
    "10. provisional_filename slugs a hostile brief; derive_title caps at 200",
    hostile.endswith(".pptx")
    and '"' not in hostile
    and ";" not in hostile
    and "/" not in hostile
    and " " not in hostile
    and len(long_title) == 200,
    f"hostile={short(hostile)} title_len={len(long_title)}",
)

# ---------------------------------------------------------------------------
# 11-13. What the fixture decks ARE. Feature 02 asserts what a validator SAYS.
# ---------------------------------------------------------------------------
print("\n-- the fixtures: what a .pptx actually contains --")

empty_bytes = deck_empty()
prs_empty, why_empty = open_deck(empty_bytes)
check(
    "11. empty.pptx opens, has ZERO slides, and clears every old assertion",
    prs_empty is not None
    and len(prs_empty.slides) == 0
    and empty_bytes.startswith(b"PK")
    and 20_000 < len(empty_bytes) < 40_000,
    f"{len(empty_bytes)} bytes, starts PK, {0 if prs_empty is None else len(prs_empty.slides)} slides",
)

junk_bytes = deck_junk()
prs_junk, why_junk = open_deck(junk_bytes)
check(
    "12. junk.pptx does NOT open, and the raise is caught rather than propagated",
    prs_junk is None and why_junk is not None and "BadZipFile" in why_junk,
    f"{len(junk_bytes)} bytes -> {short(why_junk)}",
)

thin_bytes = deck_thin_honest()
prs_thin, _ = open_deck(thin_bytes)
titles_thin = slide_titles(prs_thin) if prs_thin else []
check(
    "13. thin-honest.pptx is a GOOD deck: 3 slides, all titled, bullets under 240",
    prs_thin is not None
    and len(prs_thin.slides) == 3
    and all(titles_thin)
    and longest_run(prs_thin) < 240,
    f"titles={titles_thin} longest_run={longest_run(prs_thin) if prs_thin else '-'}",
)

# ---------------------------------------------------------------------------
# 14. R1, asserted structurally rather than remembered.
# ---------------------------------------------------------------------------
print("\n-- R1: the fix that is green here and dead on Render --")

# python-pptx's only text-fitting API is `fit_text`, and it is the obvious way to
# answer "does this bullet fit". `pptx/text/fonts.py` returns font directories
# for darwin and win32 and otherwise raises OSError("unsupported operating
# system"). It works on this Windows box and fails on every Linux deploy, with a
# message that says nothing about fonts. Character count is the honest proxy.
SEARCHED = [ROOT / "backend" / "app" / "handouts", ROOT / "scripts" / "deck_check.py"]

# Assembled from parts on purpose. This file is one of the files it scans, so a
# literal needle here matches its own source and the check fails on itself --
# which it did, first run. A self-scanning check has to be written so that
# looking for the thing is not doing the thing.
NEEDLE = "." + "fit" + "_text("

offenders: list[str] = []
for target in SEARCHED:
    files = sorted(target.rglob("*.py")) if target.is_dir() else [target]
    for path in files:
        body = path.read_text(encoding="utf-8")
        for number, line in enumerate(body.splitlines(), 1):
            if NEEDLE in line:
                offenders.append(f"{path.relative_to(ROOT).as_posix()}:{number}")
check(
    "14. fit_text is called nowhere in app/handouts or this harness",
    not offenders,
    f"offenders={offenders}" if offenders else "R1 cannot regress silently",
)

# ---------------------------------------------------------------------------
# 15-17. Extras beyond the acceptance criteria, cheap and load-bearing.
# ---------------------------------------------------------------------------
print("\n-- extras --")

# 15. _repair_message: the traceback goes in VERBATIM. A paraphrase throws away
#     the line number and the exception type, which are the whole of the signal.
repair = _repair_message("prs.save()", failed_result("boom"), "The code raised: boom")
body = repair.content
check(
    "15. _repair_message carries the code, the problem and stderr verbatim",
    "prs.save()" in body
    and "The code raised: boom" in body
    and "ZeroDivisionError" in body,
    "the model reads its own traceback; that is why one retry is worth having",
)

# 16. A layout with no title placeholder returns None from `shapes.title`, not
#     an empty string. A naive `.title.text` walk raises here, inside a validator
#     whose whole contract is that it never raises.
prs_untitled, _ = open_deck(deck_untitled())
titles_untitled = slide_titles(prs_untitled) if prs_untitled else ["x"]
check(
    "16. a Blank-layout slide reports no title without raising",
    prs_untitled is not None
    and len(prs_untitled.slides) == 2
    and titles_untitled == ["", ""],
    f"titles={titles_untitled}",
)

# 17. The overflow fixture must actually overflow, by a margin, so that
#     re-tuning the threshold from the measured distribution does not quietly
#     make this fixture valid.
prs_over, _ = open_deck(deck_overflow())
check(
    "17. overflow.pptx has 6 slides and a run far past 240 characters",
    prs_over is not None
    and len(prs_over.slides) == 6
    and longest_run(prs_over) > 480,
    f"slides={len(prs_over.slides) if prs_over else '-'} longest_run={longest_run(prs_over) if prs_over else '-'}",
)

# ---------------------------------------------------------------------------
# 20-27. Feature 02 -- artefact validation (12-robust-handouts/02, A1-A8).
#
# WRITTEN AND RUN RED BEFORE `app/handouts/validate.py` EXISTED. That order is
# the repo rule rather than a preference: `agentic_check.py` S3 went green twice
# while proving nothing, because its case was added after the code it tests. A
# case written to pass is a case that measures nothing.
#
# A9 and A10 are NOT here. They need a live model and a database
# (`agentic_check.py` S8 and a new S28), and this file is layer 1 by
# construction. A10 in particular -- raise `handout_deck_min_slides` above what
# the fixture corpus can support and watch a real job end `failed` -- is the one
# that proves the branch is reachable in production rather than merely reached
# by this harness. IT IS OWED, and until it is written the eight cases below
# pass over a validator that nothing in production may ever call.
# ---------------------------------------------------------------------------
print("\n-- feature 02: validate.check, and the branch that calls it --")

from app.config import settings  # noqa: E402

# Imported defensively, and this is scaffolding with a purpose. On the run that
# watches feature 02 fail, `app/handouts/validate.py` does not exist yet, and a
# bare top-level import would abort the whole file with a traceback -- so the
# nineteen cases above would go unreported too, and "watch it fail" would show
# no cases failing at all. A harness reports; it does not crash. The sentinel is
# neither None nor str, so every case below fails rather than passing by luck on
# a `is None` comparison.
try:
    from app.handouts import validate  # noqa: E402

    _validate_why = ""
except Exception as exc:  # noqa: BLE001
    validate = None  # type: ignore[assignment]
    _validate_why = f"{type(exc).__name__}: {exc}"


class _Missing:
    def __repr__(self) -> str:
        return f"<app.handouts.validate absent: {_validate_why}>"


MISSING = _Missing()


def verdict(recipe, data: bytes, filename: str):
    """`validate.check` over raw bytes. MISSING until feature 02 ships."""
    if validate is None:
        return MISSING
    return validate.check(recipe, artifact(filename, data))


TABLE = RECIPES["table"]
SHEET = RECIPES["sheet"]


def deck_missing_one_title() -> bytes:
    """Three slides, one of them untitled. Clears the slide floor on purpose.

    `deck_untitled()` above has two slides, so on a deck with
    `handout_deck_min_slides = 3` the slide-count branch fires first and the
    title branch is never reached. A fixture that cannot reach the branch it
    is named for is `loop.md` section 5's "test that cannot fail" in miniature.
    """
    prs = Presentation()
    _titled(prs, "Ka-band downlink", ["Downlink runs at 26 GHz [comms.md]"])
    _titled(prs, "Link margin", ["Clear-sky margin is 3 dB [comms.md]"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text_frame.text = "Handover happens hourly [comms.md]"
    return _save(prs)


def png_bytes(width: int, height: int) -> bytes:
    """A real PNG, written by Pillow -- the library the validator reads it with."""
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (width, height), (12, 74, 110)).save(buf, format="PNG")
    return buf.getvalue()


def zip_not_pptx() -> bytes:
    """A VALID zip that is not a presentation.

    The junk fixture fails at the zip layer (`BadZipFile`). This one gets past
    it and fails one layer up, inside python-pptx's own package reader -- a
    different exception type from a different library, which is the whole
    reason `check` has to catch broadly rather than name the errors it knows.
    """
    import zipfile

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("hello.txt", "not a deck")
    return buf.getvalue()


GOOD_CSV = b"Subsystem,Value,Source\nKa-band,26 GHz,comms.md\nMargin,3 dB,comms.md\n"
HEADER_ONLY_CSV = b"Subsystem,Value,Source\n"
GOOD_SHEET = b"# Ka-band revision\n\n- Downlink runs at 26 GHz [comms.md]\n\n## Test yourself\n\n1. What is the margin?\n"
NO_HEADING_SHEET = b"Downlink runs at 26 GHz. Clear-sky margin is 3 dB.\n"

# Collected as the cases run, then asserted over as a set by case 27. Gathering
# them rather than grepping the source keeps the assertion about what the
# function RETURNS, which is what the model actually reads.
messages: list[tuple[str, object]] = []


def record(label: str, value) -> object:
    if isinstance(value, str):
        messages.append((label, value))
    return value


# 20. A1 -- the zero-slide deck. 27,387 bytes, opens fine, teaches nothing.
empty_msg = record("deck empty", verdict(DECK, deck_empty(), "deck.pptx"))
check(
    "20. A1 validate.check rejects a zero-slide deck and says so in slides",
    isinstance(empty_msg, str) and "0 slides" in empty_msg,
    short(empty_msg),
)

# 21. A2 -- junk, and three other shapes of not-a-deck. The point of this case
#     is `does not raise`: a validator that raises takes `run_handout_job` with
#     it, which has happened in this repo once already, inside the one function
#     whose entire job was to stop it.
hostile = [
    ("junk", deck_junk()),
    ("empty bytes", b""),
    ("valid zip, not a deck", zip_not_pptx()),
    ("random", bytes(range(256)) * 4),
]
hostile_out = []
hostile_raised = ""
for label, data in hostile:
    try:
        hostile_out.append((label, record(f"deck {label}", verdict(DECK, data, "deck.pptx"))))
    except BaseException as exc:  # noqa: BLE001
        hostile_raised = f"{label}: {type(exc).__name__}: {exc}"
        hostile_out.append((label, exc))
check(
    "21. A2 four kinds of not-a-deck each return a string, and NOTHING raises",
    not hostile_raised and all(isinstance(v, str) and v.strip() for _, v in hostile_out),
    hostile_raised or f"{[(l, short(v, 40)) for l, v in hostile_out]}",
)

# 22. A3 -- THE CONTROL, and the only case that separates a working validator
#     from a deleted one. `DECK_PROMPT` carries an honest-shrink rule, so a thin
#     corpus SHOULD produce a short deck; firing here is the `refusal_pass = 0/2`
#     defect, where a measurement punished correct behaviour and then advised
#     deleting the thing that produced it.
thin_verdict = verdict(DECK, deck_thin_honest(), "deck.pptx")
check(
    "22. A3 a THREE-slide honest deck is accepted -- the R5 control",
    thin_verdict is None,
    "a deleted validator passes 20 and 21; only this case sees it",
)

# 23. A4 -- the over-long bullet. Character count is the only honest proxy;
#     case 14 above asserts `fit_text` is never reached for.
over_msg = record("deck overflow", verdict(DECK, deck_overflow(), "deck.pptx"))
check(
    "23. A4 an over-long bullet is named, with the limit that rejected it",
    isinstance(over_msg, str)
    and "character" in over_msg.lower()
    and str(getattr(settings, "handout_deck_max_bullet_chars", "?")) in over_msg,
    short(over_msg),
)

# 24. A5 -- THE BRANCH IS WIRED. Cases 20-23 test a function; this tests that
#     `_problem` reaches it. A validator nothing calls passes every one of them.
empty_art = artifact("deck.pptx", deck_empty())
wired = _problem(DECK, ok_result(empty_art), empty_art)
check(
    "24. A5 _problem returns non-None for a present-but-empty deck",
    wired is not None,
    short(wired) if wired else "branch 3 is not reached",
)

# 25. A6 -- PLAN.md 3.6 R-a. With the flag off, `_problem` is what case 6
#     recorded as today's behaviour, for every pair in a table. And with it on,
#     at least one row must differ -- an off switch that changes nothing is
#     gating nothing, and this case would otherwise pass over a deleted branch.
junk_art = artifact("deck.pptx", deck_junk())
thin_art = artifact("deck.pptx", deck_thin_honest())
tiny_art = artifact("chart.png", png_bytes(1, 1))

PAIRS = [
    ("crash", DECK, failed_result("The code raised an error: boom"), None),
    ("ran, no file", DECK, ok_result(artifact("notes.md")), None),
    ("empty deck", DECK, ok_result(empty_art), empty_art),
    ("junk deck", DECK, ok_result(junk_art), junk_art),
    ("thin deck", DECK, ok_result(thin_art), thin_art),
    ("1x1 chart", CHART, ok_result(tiny_art), tiny_art),
]

# `getattr` with a default, and the branch below, exist because pydantic raises
# `ValueError: "Settings" object has no field` on an unknown assignment -- which
# on the red run aborts the file at this line instead of failing this case.
prior_flag = getattr(settings, "handout_validate_artifacts", None)
if prior_flag is None:
    off = on = [MISSING] * len(PAIRS)
else:
    try:
        settings.handout_validate_artifacts = False
        off = [_problem(r, res, art) for _, r, res, art in PAIRS]
        settings.handout_validate_artifacts = True
        on = [_problem(r, res, art) for _, r, res, art in PAIRS]
    finally:
        # Restored in a `finally` -- `loop.md` section 5: a scenario owns the
        # conditions it needs, and hands them back.
        settings.handout_validate_artifacts = prior_flag

off_matches_today = (
    isinstance(off[0], str)
    and "boom" in off[0]
    and isinstance(off[1], str)
    and "deck.pptx" in off[1]
    and "notes.md" in off[1]
    and off[2] is None
    and off[3] is None
    and off[4] is None
    and off[5] is None
)
check(
    "25. A6 flag OFF reproduces case 6 exactly; flag ON changes at least one row",
    prior_flag is not None and off_matches_today and off != on,
    f"off={[short(v, 24) if v else None for v in off]} differs_on={sum(a != b for a, b in zip(off, on))}",
)

# 26. A7 -- the other three recipes. `06-test-plan.md:183` promised all four
#     were opened, in prose, and none of them ever was.
good_chart = verdict(CHART, png_bytes(640, 480), "chart.png")
bad_chart_tiny = record("chart 1x1", verdict(CHART, png_bytes(1, 1), "chart.png"))
bad_chart_notpng = record("chart not a png", verdict(CHART, b"\x89 nope", "chart.png"))
good_table = verdict(TABLE, GOOD_CSV, "table.csv")
bad_table = record("table header only", verdict(TABLE, HEADER_ONLY_CSV, "table.csv"))
good_sheet = verdict(SHEET, GOOD_SHEET, "sheet.md")
bad_sheet_empty = record("sheet empty", verdict(SHEET, b"   \n\n", "sheet.md"))
bad_sheet_flat = record("sheet no heading", verdict(SHEET, NO_HEADING_SHEET, "sheet.md"))
check(
    "26. A7 chart / table / sheet: the good ones pass, the thin ones are named",
    good_chart is None
    and good_table is None
    and good_sheet is None
    and all(
        isinstance(v, str) and v.strip()
        for v in (bad_chart_tiny, bad_chart_notpng, bad_table, bad_sheet_empty, bad_sheet_flat)
    ),
    f"good=({good_chart}, {good_table}, {good_sheet}) tiny={short(bad_chart_tiny, 40)}",
)

# 27. A8 -- every message is one the model can act on. A refusal it cannot act
#     on wastes a step (`loop.md:236-238`), and a non-ASCII one is the failure
#     three throwaway scripts in this repo have already hit on the Windows
#     console. The missing-title branch is driven here because no earlier case
#     reaches it: `deck_untitled()` has two slides, so the slide floor fires
#     first and the title check is never consulted.
record("deck missing a title", verdict(DECK, deck_missing_one_title(), "deck.pptx"))

IMPERATIVES = (
    "Add ",
    "Rewrite ",
    "Shorten ",
    "Give ",
    "Write ",
    "Save ",
    "Set ",
    "Replace ",
    "Start ",
    "Check ",
    "Plot ",
)
bad_messages = [
    (label, short(text, 60))
    for label, text in messages
    if not text.strip()
    or not text.isascii()
    or not any(verb in text for verb in IMPERATIVES)
]
check(
    "27. A8 every verdict is non-empty, ASCII, and carries an imperative",
    len(messages) >= 9 and not bad_messages,
    # The count is printed unconditionally. The first draft printed "all
    # actionable" whenever `bad_messages` was empty, which on the red run --
    # zero messages collected, so nothing to be bad -- read as a reassurance
    # beside a [FAIL]. A detail line that describes a passing run while the
    # case fails is the same defect the whole file is about, in miniature.
    f"{len(messages)} verdicts collected (>= 9 wanted); bad={bad_messages}",
)

# 28. Beyond the acceptance criteria, and it is here because case 27 could not
#     see the bug that produced it. The first build of `validate.py` emitted
#
#         slide 4 still carry template placeholder text
#
#     which is ASCII, non-empty, imperative and contains every substring cases
#     20-27 assert on. It was found by printing one verdict and READING it --
#     `build.md`'s last verification step, in a module small enough to have felt
#     exempt from it. This is the assertion that would have caught it.


def deck_two_untitled() -> bytes:
    prs = Presentation()
    _titled(prs, "Ka-band downlink", ["Downlink runs at 26 GHz [comms.md]"])
    _titled(prs, "Link margin", ["Clear-sky margin is 3 dB [comms.md]"])
    for _ in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        box.text_frame.text = "body text with no title placeholder"
    return _save(prs)


one_missing = verdict(DECK, deck_missing_one_title(), "deck.pptx")
two_missing = verdict(DECK, deck_two_untitled(), "deck.pptx")
check(
    "28. the verdict's verb agrees with the number of slides it names",
    isinstance(one_missing, str)
    and "slide 3 has no title text" in one_missing
    and isinstance(two_missing, str)
    and "slides 3 and 4 have no title text" in two_missing,
    f"one={short(one_missing, 60)} two={short(two_missing, 60)}",
)

# ---------------------------------------------------------------------------
# 30-32. Feature 03 -- the deck material budget (12-robust-handouts/03, A1-A3).
#
# WRITTEN AND RUN RED BEFORE `Recipe` HAD THE TWO FIELDS, for the same reason
# cases 20-28 were: a case added after the code it tests is a case written to
# pass. Case 30 was red on `NO_FIELD != None`; case 32 was red because the deck
# carried no budget at all, and it went red a SECOND time -- the useful time --
# the moment the budget was widened underneath a cap that could not carry it.
#
# A4, A5 and A6 ARE NOT HERE AND ARE OWED. They are two `agentic_check.py`
# scenarios and a live measurement, and this file is layer 1 by construction:
#
#   A4 / S29  force the fixture agent's `rerank_top_n` to 1 in the scenario
#             itself, restore in a `finally`, and require a deck that either
#             clears the slide floor or ends `failed` -- never a silently thin
#             `ready` deck. It has to starve retrieval ITSELF; relying on the
#             fixture's defaults being hostile enough is how S3 passed twice
#             while proving nothing.
#   A5 / S30  read the effective k/top_n off the RETRIEVE trace payload and
#             assert they are the RECIPE's, not the agent's.
#   A6        slide-count distribution at top_n=3 versus the new value, n >= 6
#             per arm, ARM ORDER ALTERNATED (PLAN.md R8 -- this repo has one
#             recorded case of a measurement reporting confidently about its
#             own loop order).
#
# Until S30 exists, NOTHING IN THIS REPOSITORY PROVES THE OVERRIDE REACHES
# PINECONE. The three cases below prove the recipe CARRIES a budget and that
# the cap can carry it. They cannot see an argument that is dropped on the way
# to `aretrieve`, and A4 alone would pass over exactly that -- a deck clears a
# slide floor because the corpus was easy, and the override never took effect.
# ---------------------------------------------------------------------------
print("\n-- feature 03: the deck material budget --")

from langchain_core.documents import Document  # noqa: E402

from app.handouts.recipes import MAX_CONTEXT_CHARS, _truncate  # noqa: E402
from app.rag.pipeline import format_context  # noqa: E402
from app.rag.retriever import META_FILENAME  # noqa: E402


# The sentinel must NOT be None, and that is the whole design of case 30.
# `getattr(recipe, "retrieve_k", None)` reads exactly like the line below and
# passes on the red run, because "the field does not exist" and "the field is
# None" collapse to one value -- a case that cannot fail, in the one file whose
# entire subject is cases that cannot fail.
class _NoField:
    def __repr__(self) -> str:
        return "<no such field on Recipe>"


NO_FIELD = _NoField()


def budget_of(key: str) -> tuple[object, object]:
    recipe = RECIPES[key]
    return (
        getattr(recipe, "retrieve_k", NO_FIELD),
        getattr(recipe, "rerank_top_n", NO_FIELD),
    )


# 30. A1 -- the identity case, ASSERTED rather than assumed.
#
#     `None` on both fields is what leaves the other three recipes untouched BY
#     CONSTRUCTION rather than by care, and an assertion is the only thing that
#     keeps "by construction" true after somebody sets a table budget in a
#     hurry. The deck is in the same case deliberately: three recipes at
#     `(None, None)` is also what a DELETED feature looks like, and only the
#     pair separates the two -- the same shape as case 22 beside case 20.
budgets = {key: budget_of(key) for key in ("chart", "table", "sheet", "deck")}
untouched = all(budgets[key] == (None, None) for key in ("chart", "table", "sheet"))
check(
    "30. A1 chart/table/sheet carry NO budget; the deck carries one",
    untouched
    and isinstance(budgets["deck"][0], int)
    and isinstance(budgets["deck"][1], int),
    f"{budgets}",
)


# 31. A2 -- PLAN.md 3.6 R-b. THIS FEATURE CHANGES MATERIAL, NEVER WORDING.
#
#     The fixture below is the whole rendered prompt, committed verbatim, and
#     the duplication is the MECHANISM rather than a lapse. This repo's standing
#     rule is that a contract stated twice drifts and the copy that drifted is
#     never the one you are reading -- which is precisely why a golden string
#     works here: drift IS the failure being detected, so the second copy has to
#     be unable to follow the first.
#
#     The fixture material carries a non-empty conversation on purpose. With an
#     empty one the prompt simply ends at the corpus, and `conversation_block`'s
#     "RECENT CONVERSATION ..." preamble -- wording that reaches the model on
#     every handout made from a live thread -- would go unmeasured.
#
#     A NOTE FOR WHOEVER SEES THIS FAIL. The question is never "how do I make it
#     pass". Pasting the new output in is a case written to pass, and this file
#     already carries two comments about what that has cost. The question is
#     whether DECK_PROMPT was MEANT to change: if it was, the fixture moves in
#     the same commit as the prompt; if it was not, the prompt regressed.
FIXTURE_BRIEF = "The Ka-band downlink budget"
FIXTURE_MATERIAL = Material(
    context="[comms.md]\nDownlink runs at 26 GHz.",
    conversation="Q: What band?\nA: Ka-band [comms.md]",
    chunk_ids=["c1"],
    turn_count=1,
)

DECK_RENDER_FIXTURE = """\
GROUNDING COMES FIRST. It outranks every instruction below, and no amount of polish is worth breaking it:
- Use only what is in the MATERIAL. Do not use prior knowledge and do not complete a half-covered idea from what you know about the subject generally.
- NEVER INVENT A NUMBER. Not a figure, not a percentage, not a date, not a count, not a unit. If the MATERIAL does not state it, it does not go in.
- A well-formatted artefact reads as though somebody checked it. A chart with labelled axes, a deck with clean bullets and a table with a header row all carry authority the MATERIAL may not have earned, which is exactly why fabricating here is worse than fabricating in prose.
- If the MATERIAL supports only part of the brief, make the smaller honest version and say what is missing. That is a correct outcome, not a failure.

HOW YOUR CODE RUNS. Read this before writing a line:
- Your entire reply is saved as one Python file and executed. Reply with code and nothing else: no prose, no explanation, no markdown fence.
- THERE IS NO FILESYSTEM AND NO NETWORK. You cannot read a CSV, open a path, download anything or query anything. Every number, label and string you use must be written into the code as a LITERAL, copied out of the MATERIAL below.
- Imports are restricted to: matplotlib, numpy, pandas, pptx, math, statistics, datetime, decimal, json, csv, io, re, textwrap, string, itertools, functools, collections, pathlib, base64. Anything else -- os, sys, requests, subprocess -- is refused before the code runs.
- eval, exec, compile, getattr, setattr and any attribute starting and ending with double underscores are refused for the same reason.
- Keep it short and straight-line. There is a wall-clock limit and a memory limit, and a clever solution that times out produces nothing at all.

- Write EXACTLY ONE file, into the current directory, named `deck.pptx`, using python-pptx (`from pptx import Presentation`). Do not write any other file.
- Do not use images, icons, charts, custom fonts or template files. There are no files to load them from and no network to fetch them over; each one is a crash, not a downgrade.
- Set 16:9 with `prs.slide_width` and `prs.slide_height` in `Inches`. Use `prs.slide_layouts[0]` for the title slide and `prs.slide_layouts[1]` for content slides -- those two exist in the default template; higher indices vary.
- Finish with `prs.save("deck.pptx")`.
- You may `print()` one short line describing the deck. It becomes the caption in the panel.

WHAT TO MAKE. A short deck answering the brief:
- Five to eight slides. A title slide, then ONE idea per slide.
- Each content slide gets a heading of at most eight words and three to five bullets. Each bullet is a complete claim taken from the MATERIAL, not a fragment and not a topic label.
- End a bullet with the source filename in square brackets when it came from one, exactly as an answer would cite it.
- If the MATERIAL only supports four slides, make four. A deck padded to a round number is padded with invention.
- Do not write a "Conclusions", "Next steps" or "Recommendations" slide unless the MATERIAL contains conclusions, next steps or recommendations. Those are the slides a model writes from nowhere.

BRIEF: The Ka-band downlink budget

MATERIAL (retrieved from this agent's corpus):
[comms.md]
Downlink runs at 26 GHz.


RECENT CONVERSATION (the user's own thread with this agent; treat the answers below as material too):
Q: What band?
A: Ka-band [comms.md]"""


def first_difference(actual: str, expected: str) -> str:
    """Where two strings part company. "3483 != 3502" is not actionable."""
    for offset, (got, want) in enumerate(zip(actual, expected)):
        if got != want:
            return (
                f"offset {offset}: got {short(actual[offset : offset + 48], 48)} "
                f"want {short(expected[offset : offset + 48], 48)}"
            )
    shared = min(len(actual), len(expected))
    if len(actual) == len(expected):
        return f"identical, {shared} chars"
    name, tail = (
        ("actual", actual[shared:]) if len(actual) > shared else ("expected", expected[shared:])
    )
    return f"agree for {shared} chars, then {name} continues {short(tail, 48)}"


rendered_deck = render(DECK, brief=FIXTURE_BRIEF, material=FIXTURE_MATERIAL)
check(
    "31. A2 render(deck) is byte-identical to the committed fixture (R-b)",
    rendered_deck == DECK_RENDER_FIXTURE,
    first_difference(rendered_deck, DECK_RENDER_FIXTURE),
)


# 32. A3 -- MAX_CONTEXT_CHARS measured against the NEW width, which is the half
#     of this feature that would otherwise ship as nothing.
#
#     Widening `rerank_top_n` while the cap stays where it was buys the model
#     one or two extra chunks and bills for ten: the rerank call is paid in
#     full, the prompt is assembled in full, and `_truncate` then deletes the
#     back half with no error, no warning and no log line. That is the exact
#     shape of the defect this whole change set exists to correct -- a green
#     pipeline over a product that is not there -- reproduced inside its own fix.
#
#     MEASURED 2026-08-17, this repo's own markdown through
#     `app.rag.ingest._prepare_chunks` at the PRODUCTION default chunk_size=800
#     TOKENS (the splitter is `from_tiktoken_encoder`, so 800 is not 800
#     characters, and reading it as characters is how 12,000 ever looked roomy):
#
#         PRD.md   37 chunks   p50 2,380   p90 3,201   max 3,316 chars
#         EVAL.md  17 chunks   p50 2,428   p90 3,099   max 3,113 chars
#
#     So a retrieved chunk is ~2,400 characters and the 90th percentile is
#     ~3,200. p90 is the figure to size a cap against, not the median: a cap
#     that clears the typical chunk and binds on the ordinary heavy one is not
#     a backstop, it is a coin flip about which half of the corpus the model
#     reads.
#
#     The width is read OFF THE RECIPE rather than hard-coded, so re-tuning the
#     budget re-runs this arithmetic automatically. Raise `rerank_top_n` to 20
#     without touching the cap and this case goes red, which is the service it
#     is here to provide.
TYPICAL_CHUNK_CHARS = 3_200
TRUNCATION_SUFFIX = "\n... [truncated]"


def context_at(width: int, chunk_chars: int) -> str:
    """The corpus half exactly as `gather_material` composes it.

    `format_context` then `_truncate`, both imported rather than re-implemented.
    A second copy of that composition here would drift from production and this
    case would go on passing about a prompt nobody builds.
    """
    filler = "Ka-band downlink material for the deck. "
    docs = [
        Document(
            page_content=(filler * (chunk_chars // len(filler) + 1))[:chunk_chars],
            metadata={META_FILENAME: f"comms-{index}.md"},
        )
        for index in range(width)
    ]
    return _truncate(format_context(docs), MAX_CONTEXT_CHARS)


deck_width = budgets["deck"][1]
if isinstance(deck_width, int):
    at_width = context_at(deck_width, TYPICAL_CHUNK_CHARS)
    on_long_chunks = context_at(deck_width, 8_000)
    # Provenance is gathered before the cap is applied, so a truncated context
    # must NOT shrink `chunk_ids`. `meta["chunk_ids"]` records what the model
    # was ALLOWED to see, explicitly not what it used -- and a wider budget
    # makes that distinction matter more, because the gap between allowed and
    # used is now large enough to mislead somebody reading the provenance.
    capped = Material(
        context=on_long_chunks, chunk_ids=[f"c{i}" for i in range(deck_width)]
    )
else:
    at_width = on_long_chunks = ""
    capped = Material()

check(
    "32. A3 the cap carries the deck's own width, and still binds on long chunks",
    isinstance(deck_width, int)
    # (a) at the recipe's configured width, p90 chunks arrive WHOLE.
    and TRUNCATION_SUFFIX not in at_width
    and len(at_width) >= deck_width * TYPICAL_CHUNK_CHARS
    # (b) the cap is still a real backstop: truncated AT it, never over-running.
    and TRUNCATION_SUFFIX in on_long_chunks
    and len(on_long_chunks) == MAX_CONTEXT_CHARS + len(TRUNCATION_SUFFIX)
    # (c) truncating the context does not truncate the provenance, and an
    #     empty MATERIAL still refuses (`is_empty`, case 9's gate, at the new
    #     width -- a wider budget must not turn "no material" into "a little").
    and len(capped.chunk_ids) == deck_width
    and not capped.is_empty
    and Material().is_empty,
    f"top_n={deck_width!r} cap={MAX_CONTEXT_CHARS} "
    f"needs={0 if not isinstance(deck_width, int) else deck_width * TYPICAL_CHUNK_CHARS} "
    f"at_width={len(at_width)} long={len(on_long_chunks)}",
)

# ---------------------------------------------------------------------------
# 40-44. Feature 04 -- failure legibility, sections D and E
#        (12-robust-handouts/04: A5, A9, A10, plus the source_code defect that
#        was measured after that file was written).
#
# WRITTEN AND RUN RED BEFORE ANY OF THE THREE DEFECTS WAS FIXED, for the fourth
# time in this file and for the same reason: `agentic_check.py` S3 went green
# twice while proving nothing, because its case was added after the code it
# tests.
#
# THE THREE DEFECTS, ALL THREE MEASURED RATHER THAN REASONED ABOUT.
#
#   D  A truncated program retries at the SAME cap and fails identically.
#      Observed live, twice (`agentic_check.py:2041-2050`): with the slide floor
#      forced to 40 the model inflated the deck until it ran out of tokens and
#      came back "Python syntax error on line 322: unterminated string literal".
#      Nothing in the backend reads `finish_reason`, so the model is told only
#      that its code raised -- and it goes and debugs a line that is fine.
#
#   E  An EMPTY generation is run in a subprocess and the retry is misdirected.
#      Read back out of `handouts.source_code` on a real run: attempt 1 was the
#      empty string, `static_check("")` accepted it (an empty program is a valid
#      empty AST), a subprocess was spawned to run nothing, `_problem` branch 2
#      fired CORRECTLY, and `_repair_message` then told the model to "check that
#      the save call is actually reached and that the filename matches" -- about
#      a program it never wrote. Trigger right, diagnosis wrong. It recovered on
#      attempt 2, so the row ended `ready` and nothing recorded it.
#
#   --  A `failed` handout stores NO source_code at all. Measured `len == 0` on
#      two failed rows. `_run_sandbox_recipe` raises before the caller assigns
#      it, so both attempts are lost exactly when somebody needs to read them.
#      `05-handout-*`'s claim that "both attempts are joined by
#      ATTEMPT_SEPARATOR -- that is what a user reads to see the retry" is true
#      only of SUCCESSFUL rows.
#
# A NUMBERING COLLISION, RECORDED RATHER THAN SILENTLY RESOLVED.
# `04-failure-legibility.md` gives case 42 to its A6 (every `error_kind` value
# is <= 16 characters and has a frontend label) and cases 43/44 to A9/A10. The
# brief that built this block assigns 42/43 to A9/A10 and 44 to the source_code
# defect, so the numbers here follow the brief and **A6 IS NOT IN THIS FILE AND
# IS OWED**. Nothing anywhere asserts that `ERROR_KIND_INVALID` plus the
# sandbox's five kinds all fit `String(16)`, or that the card can label them.
#
# WHY THESE ARE LAYER-1 DECIDABLE AT ALL. All three defects live between
# `_generate` and `_attempt`, which is the one stretch of this job with no I/O
# in it: a stub model handing back a hand-built `AIMessage` and a patched
# `_attempt` reach every branch with no model, no database and no subprocess.
# That is precisely why the cases could be written before the fix.
#
# WHAT IS STILL OWED BESIDES A6. Nothing here proves a LIVE model ever
# truncates; the two live observations cited above are the evidence for that,
# and they are an observation rather than an assertion. And case 44 drives
# `_settle` DIRECTLY -- the join from `run_handout_job`'s `except` down to it is
# asserted only by reading the source, which catches a deletion and not a
# behaviour. An `agentic_check.py` scenario reading `source_code` off a
# genuinely failed row is the assertion that would close both.
# ---------------------------------------------------------------------------
print("\n-- feature 04: truncation, an empty generation, and a failed row's code --")

import asyncio  # noqa: E402
import inspect  # noqa: E402
import uuid  # noqa: E402

from langchain_core.messages import AIMessage  # noqa: E402

# The MODULE, not the names. `_attempt` and `SessionLocal` are patched below and
# `_run_sandbox_recipe`/`_settle` resolve both as module globals at call time --
# a `from app.handouts.jobs import _attempt` on either side would make the patch
# invisible and this harness would quietly spawn subprocesses and open sockets
# instead of measuring anything. `interpreter.sandbox.run` in case 60 carries the
# same reasoning.
from app.handouts import jobs as handout_jobs  # noqa: E402
from app.handouts.jobs import HandoutFailure  # noqa: E402

_ABSENT = object()


class _NoFlag:
    """The pre-change shape of `_generate`: a bare `str`, with no flag on it.

    NOT `None` and NOT `False`, deliberately. `getattr(out, "truncated", False)`
    reads exactly like the line below and passes the red run, because "the
    function does not report truncation" and "the function reports no
    truncation" collapse to one value -- a case that cannot fail, in the file
    whose whole subject is cases that cannot fail. Case 30 above carries the
    identical sentinel for the identical reason.
    """

    def __repr__(self) -> str:
        return "<no truncation flag: _generate returned a bare str>"


NO_FLAG = _NoFlag()


def reply(text: str, finish: object = _ABSENT) -> AIMessage:
    """One hand-built model reply. `finish` omitted means the key is ABSENT.

    Three different facts, and A4 needs all three separated: `"length"` (cut
    off), `"stop"` (finished), and no `finish_reason` key at all (a provider
    that does not report one). A default of `None` would collapse the last two.
    """
    if finish is _ABSENT:
        return AIMessage(content=text)
    return AIMessage(content=text, response_metadata={"finish_reason": finish})


class _Model:
    """A chat model that hands back prepared replies in order. Never calls out.

    `ainvoke` is the only method `_generate` uses. The messages it is handed are
    RECORDED, because the retry turn's wording is half of what sections D and E
    are about and it is visible only from the model's side -- section E's whole
    defect was a repair message that was well-formed, actionable, ASCII, and
    about the wrong thing.
    """

    def __init__(self, replies: list[AIMessage]) -> None:
        self.replies = replies
        self.seen: list[list] = []

    async def ainvoke(self, messages):
        self.seen.append(list(messages))
        return self.replies[min(len(self.seen) - 1, len(self.replies) - 1)]


def generated(message: AIMessage):
    """`(text, truncated, raised)` from `_generate` over one hand-built reply."""
    try:
        out = asyncio.run(handout_jobs._generate(_Model([message]), []))
    except BaseException as exc:  # noqa: BLE001
        return "", NO_FLAG, f"{type(exc).__name__}: {exc}"
    if isinstance(out, str):
        return out, NO_FLAG, ""
    return getattr(out, "text", NO_FLAG), getattr(out, "truncated", NO_FLAG), ""


def run_recipe(replies: list[AIMessage], results: list[SandboxResult], recipe=DECK):
    """Drive `_run_sandbox_recipe` with a stub model FACTORY and no subprocess.

    The factory is the point rather than a convenience. "The retry requests a
    higher cap" is only assertable if every cap this job asks for passes through
    one observable place; handing `_run_sandbox_recipe` a pre-built model would
    make attempt 1's budget an assumption of this harness rather than a
    measurement of the code.

    Returns `(caps, ran, seen, failure, raised)`. `raised` is REPORTED, never
    swallowed: on the red run the first thing that happens is an AttributeError
    from a factory being invoked as a model, and a helper that hid it would make
    five cases fail for an unstated reason.
    """
    caps: list[object] = []
    ran: list[str] = []
    model = _Model(replies)

    def build(max_tokens=None, *args, **kwargs):
        caps.append(max_tokens)
        return model

    async def fake_attempt(code):
        ran.append(code)
        return results[min(len(ran) - 1, len(results) - 1)]

    prior_attempt = handout_jobs._attempt
    handout_jobs._attempt = fake_attempt
    failure = None
    raised = ""
    try:
        asyncio.run(handout_jobs._run_sandbox_recipe(build, [], recipe))
    except HandoutFailure as exc:
        failure = exc
    except BaseException as exc:  # noqa: BLE001
        raised = f"{type(exc).__name__}: {exc}"
    finally:
        handout_jobs._attempt = prior_attempt
    return caps, ran, model.seen, failure, raised


def last_turn(seen: list[list]) -> str:
    """Everything the model was handed on its SECOND call, flattened."""
    if len(seen) < 2:
        return ""
    return " ".join(str(getattr(m, "content", "")) for m in seen[1])


# 40. A4 -- `finish_reason`, read defensively off a hand-built message.
#
#     `04-failure-legibility.md`: "langchain-openai puts it in
#     `response_metadata["finish_reason"]`, and OpenRouter may report `"length"`
#     or a provider-specific string. Read defensively, treat an unrecognised
#     value as 'not truncated', and pin the shape with a hand-constructed
#     AIMessage rather than a live call."
#
#     The three rows the criterion names are `length`, `stop` and ABSENT. Three
#     more are here because the defensive half is the half that can be wrong in
#     silence: an uppercase spelling (the Google-native shape is `MAX_TOKENS`),
#     an unknown provider string, and an explicit `None`. Only the first of
#     those may read as truncation.
#
#     UNRECOGNISED MUST MEAN NOT-TRUNCATED, and the asymmetry is deliberate
#     (`loop.md` T3). A false positive here raises the cap and rewords the
#     repair turn on a run that was never cut off -- one wasted retry at worst.
#     A false negative is today's behaviour, which is the defect.
PROGRAM = 'from pptx import Presentation\nprs = Presentation()\nprs.save("deck.pptx")'

t_len, f_len, r_len = generated(reply(PROGRAM, "length"))
t_stop, f_stop, r_stop = generated(reply(PROGRAM, "stop"))
t_absent, f_absent, r_absent = generated(reply(PROGRAM))
_t_upper, f_upper, _r_upper = generated(reply(PROGRAM, "MAX_TOKENS"))
_t_odd, f_odd, _r_odd = generated(reply(PROGRAM, "eos_token_reached"))
_t_none, f_none, _r_none = generated(reply(PROGRAM, None))

check(
    "40. A4 _generate reports finish_reason=length as truncation; stop, absent and unknown are not",
    not (r_len or r_stop or r_absent)
    and f_len is True
    and f_upper is True
    and f_stop is False
    and f_absent is False
    and f_odd is False
    and f_none is False
    # The code still comes back. A truncation flag bolted on at the cost of the
    # payload would pass every line above it.
    and t_len == PROGRAM
    and t_stop == PROGRAM
    and t_absent == PROGRAM,
    r_len
    or r_stop
    or r_absent
    or f"length={f_len!r} MAX_TOKENS={f_upper!r} stop={f_stop!r} absent={f_absent!r} "
    f"unknown={f_odd!r} none={f_none!r}",
)

# 41. A5 -- THE RETRY THAT CAN ACTUALLY SUCCEED, and it is two halves.
#
#     Half one is the cap: a retry at the same 4,096 tokens produces the same
#     length and fails identically, which is `PLAN.md` R7 and is what was
#     observed live twice.
#
#     Half two is the sentence, and without it the raised cap buys a longer
#     WRONG program. The model is told only that its code raised, so it goes
#     looking for a bug on the line where the string ended -- a line that is
#     fine. `"cut off"` in the retry turn is the assertion that the diagnosis
#     travelled with the budget.
#
#     The fixture program stops mid-string on purpose: that is the shape the
#     live failure had ("unterminated string literal"), and it is what makes the
#     sandbox's own message plausible enough to mislead.
TRUNCATED_PROGRAM = (
    "from pptx import Presentation\n"
    "prs = Presentation()\n"
    "for n in range(40):\n"
    "    s = prs.slides.add_slide(prs.slide_layouts[1])\n"
    '    s.shapes.title.text = "Slide %d of the Ka-band downlink budget, marg'
)
SYNTAX_FAIL = failed_result(
    "Python syntax error on line 322: unterminated string literal", kind="syntax"
)

caps41, ran41, seen41, fail41, raised41 = run_recipe(
    [reply(TRUNCATED_PROGRAM, "length"), reply(TRUNCATED_PROGRAM, "length")],
    [SYNTAX_FAIL],
)
retry41 = last_turn(seen41).lower()
cap_setting = getattr(settings, "handout_code_max_tokens", None)
check(
    "41. A5 a truncation-triggered retry asks for a HIGHER cap, and tells the model it was cut off",
    not raised41
    # The cap is a SETTING now (PLAN.md 3.1), not a module constant, so it can
    # be moved without an edit to the job.
    and isinstance(cap_setting, int)
    and len(caps41) == 2
    and caps41[0] == cap_setting
    and isinstance(caps41[1], int)
    and caps41[1] > caps41[0]
    and len(ran41) == 2
    and "cut off" in retry41,
    raised41
    or f"caps={caps41} setting={cap_setting!r} attempts={len(ran41)} "
    f"retry_says_cut_off={'cut off' in retry41}",
)

# 42. A9 -- THE SUBPROCESS THAT RAN NOTHING.
#
#     "Assert the subprocess is never reached, not merely that the run failed."
#     `_attempt` is counted rather than `sandbox.run`, because `_attempt` is the
#     boundary the check has to sit above: it performs the static check, and
#     `static_check("")` ACCEPTS an empty program -- an empty AST is a valid one
#     -- so every guard below this line lets nothing through to a real process
#     while still paying for one.
#
#     Both shapes in one run: `""` on attempt 1 and whitespace-only on attempt
#     2. Whitespace-only is the one a `if not text:` guard would let through, and
#     `_strip_fence` already `.strip()`s, so the two arrive here identically --
#     which is worth asserting rather than assuming.
#
#     THE CAP MUST NOT MOVE HERE, and that conjunct is the pair to case 41. An
#     empty reply with `finish_reason="stop"` is a model that returned nothing,
#     not a model that ran out of room; raising the budget for it would make
#     "the cap went up" mean nothing.
NO_FILES = ok_result()  # ran clean, wrote nothing -- `_problem` branch 2

caps42, ran42, seen42, fail42, raised42 = run_recipe(
    [reply("", "stop"), reply("   \n\t  ", "stop")],
    [NO_FILES],
)
check(
    "42. A9 an empty and a whitespace-only generation never reach _attempt at all",
    not raised42
    and ran42 == []
    and fail42 is not None
    and len(caps42) == 2
    and caps42[0] == caps42[1],
    raised42
    or f"attempt_calls={len(ran42)} caps={caps42} failure={short(str(fail42), 60)}",
)

# 43. A10 -- THE SENTENCE, which is the whole of section E.
#
#     The trigger was already right on the measured run; the DIAGNOSIS was
#     wrong. `_problem` branch 2 said "check that the save call is actually
#     reached and that the filename matches" about a program that was never
#     written, and every error-shaped check in the repository was green for it.
#
#     Both the failure message and THE TURN THE MODEL ACTUALLY READ are
#     asserted. Only the second is the defect: a correct `error` string beside a
#     misdirecting repair message is exactly what shipped.
msg43 = str(fail42) if fail42 is not None else ""
retry43 = last_turn(seen42)
# ".pptx" rather than "deck.pptx": the recipe path always asks for one name, so
# a message naming any filename at all is the failure being tested for.
BANNED43 = ("save", "filename", ".pptx")
bad43 = sorted(
    {w for w in BANNED43 if w in msg43.lower() or w in retry43.lower()}
)
check(
    "43. A10 the no-code problem names the missing CODE, and mentions no save call and no filename",
    bool(msg43.strip())
    and msg43.isascii()
    and "no code" in msg43.lower()
    and bool(retry43.strip())
    and retry43.isascii()
    and "no code" in retry43.lower()
    and not bad43,
    f"banned={bad43} problem={short(msg43, 80)}",
)

# 44. A `failed` handout must KEEP THE CODE IT TRIED.
#
#     Measured `len(source_code) == 0` on two failed rows. `_run_sandbox_recipe`
#     raises, and the caller assigns `source_code` only on the success path, so
#     both attempts are discarded at exactly the moment somebody needs to read
#     them -- and `04-handouts-panel.md`'s "both attempts are joined by
#     ATTEMPT_SEPARATOR; that is what a user reads to see the retry" is true
#     only of rows that did not need the retry.
#
#     Two halves, because they can fail independently: the attempts have to
#     survive the raise, and `_settle` -- which writes the failed row from its
#     own second session -- has to write them.
ATTEMPT_ONE = "from pptx import Presentation\nprs = Presentation()  # attempt one"
ATTEMPT_TWO = "from pptx import Presentation\nprs = Presentation()  # attempt two"

caps44, ran44, seen44, fail44, raised44 = run_recipe(
    [reply(ATTEMPT_ONE, "stop"), reply(ATTEMPT_TWO, "stop")],
    [failed_result("The code raised an error: boom")],
)
kept44 = getattr(fail44, "source_code", None)


class _Row:
    """The columns `_settle` writes, and the state it must refuse to touch."""

    def __init__(self, status: str = "pending") -> None:
        self.status = status
        self.error = None
        self.meta = None
        self.source_code = None


class _FakeSession:
    """Enough of `AsyncSession` for `_settle`. No engine, no socket, no row.

    `_settle`'s subject is what it WRITES, and every branch of that is decided
    before SQLAlchemy is involved -- so a fake session measures the function
    rather than the database, and this file stays layer 1.
    """

    def __init__(self, row) -> None:
        self.row = row
        self.commits = 0

    async def __aenter__(self):
        return self

    async def __aexit__(self, *exc):
        return False

    async def scalar(self, statement):
        return self.row

    async def commit(self):
        self.commits += 1


def settle_onto(row, **kwargs):
    """`_settle` against a fake session. Returns `(session, raised)`.

    An unknown keyword raises at the CALL, outside `_settle`'s own try/except,
    so on the red run this reports the missing parameter by name rather than
    reporting a swallowed nothing.
    """
    session = _FakeSession(row)
    prior = handout_jobs.SessionLocal
    handout_jobs.SessionLocal = lambda: session
    try:
        asyncio.run(
            handout_jobs._settle(
                uuid.uuid4(), uuid.uuid4(), "both attempts failed", **kwargs
            )
        )
        return session, ""
    except BaseException as exc:  # noqa: BLE001
        return session, f"{type(exc).__name__}: {exc}"
    finally:
        handout_jobs.SessionLocal = prior


settled_code = kept44 if isinstance(kept44, str) else "<no source_code on the failure>"
pending_row = _Row("pending")
_sess44, settle_raised = settle_onto(
    pending_row, error_kind="syntax", attempts=2, source_code=settled_code
)
# R9, re-pinned here because this feature adds a THIRD field to the write and a
# field written unconditionally would resurrect a terminal row's code.
ready_row = _Row("ready")
settle_onto(ready_row, error_kind="syntax", attempts=2, source_code="must not land")

meta44 = pending_row.meta if isinstance(pending_row.meta, dict) else {}

# THE JOIN BETWEEN THE TWO HALVES, and it is asserted STRUCTURALLY because
# nothing cheaper reaches it. The two halves above can both be correct while the
# wire between them is missing: `_run_sandbox_recipe` puts the code on the
# exception, `_settle` writes whatever it is handed, and `run_handout_job` --
# which catches the one and calls the other -- quietly passes nothing. Driving
# `run_handout_job` itself needs an agent, a handout, a recipe and a corpus, so
# it belongs to `agentic_check.py`; until that scenario exists this reads the
# source, in the manner of case 14, and catches the deletion rather than the
# behaviour. It is the weakest assertion in this block and it is here because
# the alternative was none.
job_source = inspect.getsource(handout_jobs.run_handout_job)
settle_call = job_source.split("await _settle(", 1)[-1] if "await _settle(" in job_source else ""
wired44 = "exc.source_code" in job_source and "source_code=" in settle_call

check(
    "44. a failed run keeps BOTH attempts, and _settle writes them onto the failed row",
    not raised44
    and isinstance(kept44, str)
    and bool(kept44.strip())
    and ATTEMPT_ONE in kept44
    and ATTEMPT_TWO in kept44
    and kept44.count("ATTEMPT") == 2
    and not settle_raised
    and pending_row.status == "failed"
    and pending_row.source_code == settled_code
    and meta44.get("error_kind") == "syntax"
    and meta44.get("attempts") == 2
    # The `pending` guard survives: a row that already settled keeps its own.
    and ready_row.status == "ready"
    and ready_row.source_code is None
    and wired44,
    raised44
    or settle_raised
    or f"source_code={short(str(kept44), 60)} row_status={pending_row.status} "
    f"row_code={short(str(pending_row.source_code), 40)} meta={meta44} wired={wired44}",
)

# ---------------------------------------------------------------------------
# 50-56. Feature 05 -- the deck outline preview (12-robust-handouts/05, A1-A7).
#
# WRITTEN AND RUN RED BEFORE `validate.outline` AND `jobs._preview_for` EXISTED,
# for the third time in this file and for the same reason: `agentic_check.py` S3
# went green twice while proving nothing because its case was added after the
# code it tests.
#
# WHAT THIS FEATURE IS. `handouts.preview_text` already exists, is already on
# `HandoutDetail`, and is already rendered inside `HandoutCard`'s `Reveal`. For a
# sandbox recipe it holds the model's own single `print()` line -- a caption,
# where an outline would fit. Feature 02's validator ALREADY opens the deck and
# walks its slides; this is that walk returning its titles instead of discarding
# them. So there is no column, no route, no migration and no new component, and
# the whole feature is two pure functions -- which is exactly the shape this file
# exists to measure.
#
# A8 is the frontend half (`HandoutCard.test.tsx`, `npm test`): jsdom computes no
# layout, and whether the outline's lines survive into the DOM is a React fact
# rather than a Python one.
#
# A9 IS OWED. It is an `agentic_check.py` S8 assertion -- the STORED
# `preview_text` names a slide count matching the deck a live model actually
# produced -- and it needs a database, Pinecone, Cohere and a model. Until it is
# written, the seven cases below prove the outline is CORRECT and prove nothing
# about it ever reaching a row. That is the same gap cases 20-28 and 60-64 each
# record for themselves, and it is worth restating rather than assuming somebody
# reads the other two notes.
# ---------------------------------------------------------------------------
print("\n-- feature 05: the deck outline preview --")

# Imported defensively for the reason the `validate` import above is, and with
# one addition: on the red run `app/handouts/validate.py` DOES exist -- feature
# 02 shipped it -- and simply has no `outline`. A bare call would then raise
# `AttributeError` from inside the helper's own `try`, and be recorded as "the
# function raised" when the truth is "the function is not built yet". The two
# are different facts and only one of them is a defect, so the presence check is
# separate from the raise check.
try:
    from app.handouts.jobs import _preview_for  # noqa: E402

    _preview_why = ""
except Exception as exc:  # noqa: BLE001
    _preview_for = None  # type: ignore[assignment]
    _preview_why = f"{type(exc).__name__}: {exc}"

HAS_OUTLINE = validate is not None and hasattr(validate, "outline")


class _NoOutline:
    """A sentinel of this block's OWN, rather than `MISSING` above.

    Reusing `MISSING` prints `<app.handouts.validate absent: >` on the red run,
    which is a detail line describing a module that imported perfectly. This
    file already carries two comments about a detail line that describes the
    wrong thing; that is the third.
    """

    def __repr__(self) -> str:
        return (
            "<validate.outline not built>"
            if validate is not None
            else f"<app.handouts.validate absent: {_validate_why}>"
        )


NO_OUTLINE = _NoOutline()

# 2,000 characters, from `05-deck-outline-preview.md` section C. Written here as
# a literal rather than imported from `validate`, deliberately: importing the
# cap would assert that the module agrees with itself, and the fact worth
# pinning is that a detail fetch does not carry a prompt-sized field.
PREVIEW_CAP = 2_000


def deck_forty() -> bytes:
    """Forty slides with titles long enough that the whole outline cannot fit.

    ~57 characters per rendered line against a 2,000-character cap, so roughly
    two thirds of the deck lists and the rest has to be counted. A fixture that
    only just overflows would stop overflowing the first time somebody shortens
    the title format, and the case would go green while measuring nothing --
    `deck_overflow()` above carries the same reasoning about the bullet limit.
    """
    prs = Presentation()
    for n in range(1, 41):
        _titled(
            prs,
            f"Slide {n}: Ka-band downlink budget, margin and handover",
            [f"A short and entirely reasonable bullet {n} [comms.md]"],
        )
    return _save(prs)


# An em-dash, a section sign and an emoji -- the three shapes that have actually
# broken a throwaway script in this repo, plus the one a model reaches for most.
# Written as escapes rather than as literals so that this file stays ASCII and
# cannot itself become the thing it is testing for.
UNICODE_TITLES = [
    "Ka-band \u2014 downlink budget",
    "Margin \u00a7 3.2 clear sky",
    "Handover \U0001f680 window",
]


def deck_unicode() -> bytes:
    """Three slides whose titles are model-written text of the worst kind."""
    prs = Presentation()
    for title in UNICODE_TITLES:
        _titled(prs, title, ["Downlink runs at 26 GHz [comms.md]"])
    return _save(prs)


def ran_and_printed(text: str, *artifacts: SandboxArtifact) -> SandboxResult:
    """`ok_result` with a caption. The stdout half is the whole of case 56."""
    return SandboxResult(
        ok=True, exit_code=0, stdout=text, stderr="", artifacts=list(artifacts)
    )


def outline_of(spec, data: bytes, filename: str = "deck.pptx"):
    """`(value, raised)` from `validate.outline`. MISSING until this ships.

    A raise is RETURNED, never propagated and never swallowed. `outline` carries
    `check`'s contract -- it is called from inside `run_handout_job`, which never
    raises -- so "did it raise" is the assertion A2 and A3 are actually about,
    and a helper that caught quietly would turn it into a green row.
    """
    if not HAS_OUTLINE:
        return NO_OUTLINE, ""
    try:
        return validate.outline(spec, artifact(filename, data)), ""
    except BaseException as exc:  # noqa: BLE001 - cases 51 and 52 measure this
        return NO_OUTLINE, f"{filename}: {type(exc).__name__}: {exc}"


def preview_of(recipe, stdout: str, art: SandboxArtifact):
    """What `jobs` would write into `handouts.preview_text` for one attempt."""
    if _preview_for is None:
        return NO_OUTLINE
    try:
        return _preview_for(recipe, ran_and_printed(stdout, art), art)
    except BaseException as exc:  # noqa: BLE001
        return f"<raised {type(exc).__name__}: {exc}>"


# 50. A1 -- the honest three-slide deck. Titles AND their order: an outline that
#     lists the right slides in the wrong order is a preview of a deck nobody
#     has, and `"contains all three"` alone cannot see it.
THIN_TITLES = ["Ka-band downlink", "Link margin", "Handover"]
thin_outline, thin_raised = outline_of(DECK, deck_thin_honest())
thin_at = (
    [thin_outline.find(t) for t in THIN_TITLES] if isinstance(thin_outline, str) else []
)
check(
    "50. A1 outline(deck, thin-honest) counts 3 slides and numbers all three titles in order",
    isinstance(thin_outline, str)
    and "3 slides" in thin_outline
    and all(
        f"{n}. {title}" in thin_outline for n, title in enumerate(THIN_TITLES, start=1)
    )
    and all(at >= 0 for at in thin_at)
    and thin_at == sorted(thin_at),
    thin_raised or short(thin_outline),
)

# 51. A2 -- THE ARTEFACT THIS FEATURE EXISTS TO WARN ABOUT. A preview that
#     crashes on the empty deck is worse than no preview at all: the user gets a
#     `failed` row for a job that produced a file, and the one card that would
#     have said "this deck is empty" is the card that took the job down.
empty_outline, empty_raised = outline_of(DECK, deck_empty())
check(
    "51. A2 outline(deck, empty) says zero slides and does NOT raise",
    not empty_raised and isinstance(empty_outline, str) and "0 slides" in empty_outline,
    empty_raised or short(empty_outline),
)

# 52. A3 -- None, three ways, and one string.
#
#     The junk deck is A3 itself: `check` has already said what is wrong with it
#     in words the model can act on, and a preview has nothing to add.
#
#     The other three conjuncts are the DISPATCH, folded in here rather than
#     given a case number because 50-56 are the numbers this feature owns. They
#     are not decoration. `chart` already renders its own PNG as a thumbnail, so
#     a text preview beside it is noise; and `sheet`'s `preview_text` IS the
#     study sheet -- the only recipe whose preview is the product -- so an
#     outline there would delete the feature rather than add one. A dispatch
#     table that gained a fourth entry by symmetry would pass every other case
#     in this block.
junk_outline, junk_raised = outline_of(DECK, deck_junk())
chart_outline, chart_raised = outline_of(CHART, png_bytes(640, 480), "chart.png")
sheet_outline, sheet_raised = outline_of(SHEET, GOOD_SHEET, "sheet.md")
table_outline, table_raised = outline_of(TABLE, GOOD_CSV, "table.csv")
check(
    "52. A3 junk -> None without raising; chart and sheet get NO outline; the table does",
    not (junk_raised or chart_raised or sheet_raised or table_raised)
    and junk_outline is None
    and chart_outline is None
    and sheet_outline is None
    and isinstance(table_outline, str)
    and "Subsystem" in table_outline
    and "2" in table_outline,
    junk_raised
    or chart_raised
    or sheet_raised
    or table_raised
    or f"junk={junk_outline} chart={chart_outline} sheet={sheet_outline} table={short(table_outline, 60)}",
)

# 53. A4 -- the cap, and the tail that makes the truncation honest.
#
#     `preview_text` is `Text` and unbounded, but it is returned on a detail
#     fetch and a forty-slide outline is prompt-sized. The tail is the half that
#     matters: a preview silently cut at 2,000 characters looks like a deck that
#     stops at slide 34, which is the same defect as a truncated deck program --
#     plausible right up to the point it ends.
forty_outline, forty_raised = outline_of(DECK, deck_forty())
forty_lines = forty_outline.splitlines() if isinstance(forty_outline, str) else []
# `[1:]` -- the header is "40 slides", which also starts with a digit. Counting
# it made the detail line read "34 listed" over an outline that listed 33, which
# is a detail line describing the wrong thing beside a passing case. Two such
# comments already exist in this file; this is the third.
forty_listed = len([line for line in forty_lines[1:] if line[:1].isdigit()])
check(
    "53. A4 a 40-slide outline stops at 2000 chars and counts the slides it did not list",
    isinstance(forty_outline, str)
    and forty_outline.startswith("40 slides")
    and len(forty_outline) <= PREVIEW_CAP
    and len(forty_lines) > 2
    and forty_lines[-1].startswith("... (")
    and forty_lines[-1].endswith("more slides)")
    # The blank line before it, and it is not spacing. `HandoutCard` renders
    # this field through `Markdown`, where `1. ...` is an ordered list, so a
    # tail sitting directly under the last item is a LAZY CONTINUATION of it --
    # the card would read "slide 33: ... (7 more slides)" as that slide's title.
    # Found by rendering one and looking at it; every string assertion above
    # passes either way.
    and forty_lines[-2] == ""
    # It truncated, and it truncated something: a tail on an outline that listed
    # every slide, or on one that listed none, would both satisfy the two lines
    # above.
    and 0 < forty_listed < 40,
    forty_raised
    or f"{len(forty_outline) if isinstance(forty_outline, str) else '-'} chars, "
    f"{forty_listed} listed, tail={short(forty_lines[-1] if forty_lines else '', 40)}",
)

# 54. A5 -- ASCII, and the words still there.
#
#     `isascii()` ALONE IS PASSED BY THE EMPTY STRING, and by an outline that
#     dropped every title it could not encode. The word assertions are what
#     separate "made safe" from "made empty" -- the same pairing as case 22
#     beside case 20. Slide titles are model-written text going into a field a
#     Python script may print, and the Windows console codepage has broken three
#     throwaway scripts in this repo on exactly this class of character.
uni_outline, uni_raised = outline_of(DECK, deck_unicode())
UNICODE_SURVIVORS = (
    "Ka-band",
    "downlink budget",
    "Margin",
    "3.2 clear sky",
    "Handover",
    "window",
)

# AND THE HALF `isascii()` CANNOT SEE. `"\x00".isascii()` is True, so the
# assertion above passes a string carrying a NUL -- and Postgres rejects NUL in
# a text column, so that string fails the commit and takes a finished handout
# with it. The deck path cannot produce one (XML forbids it, so python-pptx
# raises long before this), which is exactly why the probe is a CSV: binary
# bytes under a `.csv` name are reachable, and were returning
# `"0 data rows\n\x00\x01"` when this was written.
CONTROL_CSV = b"Sub\x00system,Va\x01lue,Sou\x7frce\nKa-band,26 GHz,comms.md\n"
ctrl_outline, ctrl_raised = outline_of(TABLE, CONTROL_CSV, "table.csv")
printable = isinstance(ctrl_outline, str) and all(
    ch == "\n" or 0x20 <= ord(ch) <= 0x7E for ch in ctrl_outline
)
check(
    "54. A5 an em-dash, a section sign and an emoji come back as pure ASCII, words intact",
    isinstance(uni_outline, str)
    and uni_outline.isascii()
    and "3 slides" in uni_outline
    and all(word in uni_outline for word in UNICODE_SURVIVORS)
    and not ctrl_raised
    and printable
    and "Subsystem" in ctrl_outline,
    uni_raised or ctrl_raised or f"{short(uni_outline)} ctrl={short(ctrl_outline, 40)}",
)

# 55. A6 -- THE PREVIEW IS A PRODUCT FEATURE; THE FLAG IS A REGRESSION SWITCH.
#
#     `handout_validate_artifacts` exists so that `_problem` can be shown to be
#     byte-identical to its pre-validation self (PLAN.md 3.6 R-a). It is not a
#     product option, and wiring the outline behind it -- which is the obvious
#     thing to do, since the outline reuses the validator's own walk -- would
#     delete every preview the moment somebody flipped the switch to reproduce a
#     baseline.
#
#     The flag is set EXPLICITLY in a try/finally rather than read from the
#     ambient value. Case 6 above was written the other way and asserted the
#     flag-OFF contract while running flag-ON; case 64 below records the same
#     correction. Two agents have now had to make it.
_prior55 = getattr(settings, "handout_validate_artifacts", None)
off_outline: object = NO_OUTLINE
off_preview: object = NO_OUTLINE
off_raised55 = ""
if _prior55 is not None:
    try:
        settings.handout_validate_artifacts = False
        off_outline, off_raised55 = outline_of(DECK, deck_thin_honest())
        off_preview = preview_of(
            DECK, "deck written with 3 slides", artifact("deck.pptx", deck_thin_honest())
        )
    finally:
        settings.handout_validate_artifacts = _prior55
check(
    "55. A6 the outline and the preview are NOT gated on handout_validate_artifacts",
    _prior55 is not None
    and not off_raised55
    and isinstance(off_outline, str)
    and "3 slides" in off_outline
    and isinstance(off_preview, str)
    and "3 slides" in off_preview
    and "1. Ka-band downlink" in off_preview,
    off_raised55 or f"flag was {_prior55}; preview={short(off_preview, 60)}",
)

# 56. A7 -- PLAN.md 3.6 R-a, applied to this field.
#
#     A recipe with no outline function behaves exactly as it did before this
#     feature: `result.stdout.strip() or None`, with `.strip()` and with `None`
#     rather than `""` -- the panel renders no caption at all for an empty one,
#     and `""` would render an empty block instead.
#
#     The last conjunct is the other direction, and without it this case passes
#     over a feature that was never wired: an implementation that always returned
#     the caption would satisfy every line above it.
CAPTION = "chart written with 4 bars"
chart_art = artifact("chart.png", png_bytes(640, 480))
deck_art56 = artifact("deck.pptx", deck_thin_honest())
kept = preview_of(CHART, CAPTION, chart_art)
kept_padded = preview_of(CHART, "  " + CAPTION + "  \n", chart_art)
kept_blank = preview_of(CHART, "   \n  ", chart_art)
deck_preview = preview_of(DECK, CAPTION, deck_art56)
check(
    "56. A7 a recipe with no outline still gets the stdout caption, byte for byte",
    kept == CAPTION
    and kept_padded == CAPTION
    and kept_blank is None
    and isinstance(deck_preview, str)
    and "3 slides" in deck_preview,
    _preview_why
    or f"kept={short(kept, 40)} padded={short(kept_padded, 40)} blank={kept_blank!r} "
    f"deck={short(deck_preview, 40)}",
)

# ---------------------------------------------------------------------------
# 60-64. Feature 06 -- tool-path parity (12-robust-handouts/06, A1-A5).
#
# THE OTHER DOOR. Everything above measures the panel button's path. A deck asked
# for IN CHAT goes through `run_python` instead, which shares no prompt, no
# grounding rules and no validation with it -- so feature 02 landing alone means
# the identical 27,387-byte empty deck ships through the door the workshop
# actually demonstrates.
#
# WRITTEN AND RUN RED BEFORE `interpreter._run` VALIDATED ANYTHING, for the same
# reason cases 20-28 were: `agentic_check.py` S3 went green twice while proving
# nothing because its case was added after the code it tests.
#
# A6 is `llm_check.py` case 30 -- the request must not widen -- because it is a
# fact about `build_chat_model`, not about handouts.
#
# A7, A8 and A9 are `agentic_check.py` scenarios (S32, S33, and S16 unchanged).
# THEY ARE OWED. Until they exist, everything below passes over a validator that
# nothing in a real turn may ever reach: these cases prove the wiring, not that
# a live model is ever told about it.
# ---------------------------------------------------------------------------
print("\n-- feature 06: the other door -- run_python's artefacts --")

import asyncio  # noqa: E402

from app.tools import interpreter  # noqa: E402
from app.tools.registry import ToolContext  # noqa: E402

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def run_tool(result, *, code="prs.save('deck.pptx')", purpose="a deck"):
    """`interpreter._run` over a hand-built SandboxResult. The sandbox never runs.

    `_run` is a closure inside `build_python_tool`, and `StructuredTool` keeps it
    on `.coroutine` -- so this reaches the real function rather than a copy of it.
    `sandbox.run` is patched on the MODULE object, because `interpreter` calls it
    as `sandbox.run(...)` and resolves that attribute per call; a
    `from ... import run` on either side would make the patch invisible and this
    harness would quietly spawn subprocesses instead.

    Returns `(text, outcome, persisted_filenames, raised)`. **It does not swallow
    exceptions** -- `raised` is the whole of case 63, and a helper that caught
    them would turn the one assertion a code interpreter cannot afford to lose
    into a green row.
    """
    # Neither field is touched by `_run`; the artifact list is the only part of
    # `ToolContext` this path writes.
    ctx = ToolContext(agent=None, ledger=None)  # type: ignore[arg-type]
    tool = interpreter.build_python_tool(ctx)

    async def _fake_run(source, *, timeout_s=None):
        return result

    prior_run = interpreter.sandbox.run
    interpreter.sandbox.run = _fake_run
    try:
        text, outcome = asyncio.run(tool.coroutine(code=code, purpose=purpose))
    except BaseException as exc:  # noqa: BLE001 - case 63 measures exactly this
        return "", None, [], f"{type(exc).__name__}: {exc}"
    finally:
        interpreter.sandbox.run = prior_run
    return text, outcome, [a.filename for a in ctx.artifacts], ""


# 60. A1 -- `check` reachable BY KIND. The tool path has no `Recipe`: the model
#     chose the filename and the sandbox stamped a MIME type, and that is all
#     there is. The second half of this case is the part a by-kind entry point is
#     easy to get wrong -- with a Recipe the name in the message is always
#     `deck.pptx`, and quoting that back at a model which wrote `ka-band.pptx`
#     is a refusal it cannot act on.
def by_kind(kind: str, filename: str, data: bytes, mime: str):
    """`validate.check` reached by KIND, reported rather than raised.

    Defensive for the same reason the `validate` import above is: on the run that
    watches this feature fail, `check` still takes a `Recipe` and dies on
    `recipe.key` with an `AttributeError` raised OUTSIDE its own try block. A
    bare call there aborts the whole file, so the twenty-odd cases above go
    unreported and "watch it fail" shows no failures at all.
    """
    if validate is None:
        return MISSING
    try:
        return validate.check(kind, artifact(filename, data, mime))
    except Exception as exc:  # noqa: BLE001
        return f"<raised {type(exc).__name__}: {exc}>"


kind_empty = by_kind("deck", "deck.pptx", deck_empty(), PPTX_MIME)
kind_named = by_kind("deck", "ka-band.pptx", deck_empty(), PPTX_MIME)
kind_thin = by_kind("deck", "deck.pptx", deck_thin_honest(), PPTX_MIME)
kind_unknown = by_kind("file", "notes.txt", b"hello", "text/plain")
check(
    "60. A1 validate.check takes a KIND, quotes the model's own filename, and still passes a good deck",
    isinstance(kind_empty, str)
    and "0 slides" in kind_empty
    and isinstance(kind_named, str)
    and "ka-band.pptx" in kind_named
    and kind_thin is None
    and kind_unknown is None,
    f"empty={short(kind_empty, 40)} named={short(kind_named, 40)} thin={kind_thin} unknown={kind_unknown}",
)

# 61. A2 -- the rejection reaches the model as a message it can act on, and the
#     outcome is `ok=False` so it lands in TOOL_ERROR (an EXISTING trace type --
#     PLAN.md 3.5 forbids minting one here). "Not a bare 'invalid'" is the point:
#     a refusal the model cannot act on wastes a step, and steps are bounded by
#     `max_tool_steps` and MAX_CONSECUTIVE_FAILED_STEPS.
#
#     The last two conjuncts are agreement, folded in here rather than given a
#     case of their own because 60-64 are the case numbers this feature owns.
#     They are not decoration: the first build of the rejection block said
#     "These files ... They were NOT saved" over ONE deck, which is ASCII,
#     non-empty, imperative and contains every other substring this case
#     asserts. `validate.py` shipped the identical defect the day before and its
#     case 28 exists because of it. Found by printing one message and reading
#     it, which is `build.md`'s last verification step -- twice now, in two
#     modules, both small enough to have felt exempt from it.
text61, out61, kept61, raised61 = run_tool(
    ok_result(artifact("deck.pptx", deck_empty(), PPTX_MIME))
)
text61b, _o61b, _k61b, _r61b = run_tool(
    ok_result(
        artifact("deck.pptx", deck_empty(), PPTX_MIME),
        artifact("second.pptx", deck_junk(), PPTX_MIME),
    )
)
check(
    "61. A2 a rejected artefact is ToolOutcome(ok=False) naming the defect AND an action",
    not raised61
    and out61 is not None
    and out61.ok is False
    and isinstance(text61, str)
    and "deck.pptx" in text61
    and "0 slides" in text61
    and text61.isascii()
    and any(verb in text61 for verb in IMPERATIVES)
    and bool(out61.error)
    and kept61 == []
    and "This file was opened and checked" in text61
    and "These files were opened and checked" in text61b,
    raised61 or f"ok={out61 and out61.ok} kept={kept61} text={short(text61, 90)}",
)

# 62. A3 -- PER ARTEFACT. The tool path keeps every file
#     (`interpreter.py:183-191`) where the recipe path keeps one
#     (`jobs.py:268-274`), so a run can persist the good ones and reject the bad
#     one. That is a genuine difference from the recipe path and it is correct:
#     throwing away a valid CSV because the deck beside it was empty punishes the
#     half that worked.
text62, out62, kept62, raised62 = run_tool(
    ok_result(
        artifact("good.csv", GOOD_CSV, "text/csv"),
        artifact("empty.pptx", deck_empty(), PPTX_MIME),
    ),
    purpose="a table and a deck",
)
check(
    "62. A3 good.csv + empty.pptx persists ONLY good.csv, and says so about both",
    not raised62
    and kept62 == ["good.csv"]
    and "good.csv" in text62
    and "empty.pptx" in text62
    and "NOT saved" in text62
    and out62 is not None
    and out62.ok is False,
    raised62 or f"kept={kept62} text={short(text62, 90)}",
)

# 63. A4 -- `_run` NEVER RAISES, for any validator outcome INCLUDING a validator
#     that itself throws. `loop.md` section 4: a tool that fails is a message,
#     never an exception, and it is the single most valuable behaviour a code
#     interpreter has. `validate.check` already promises never to raise; this
#     case assumes that promise will one day be broken -- it was broken once
#     before in this repo, inside `_static_refusal`, the one function whose whole
#     job was to prevent it.
#
#     Every probe also asserts the artefact SURVIVES. A bug in validation is
#     evidence about the validator, not about the file, so the permissive
#     resolution is the honest one -- and the alternative is a validator bug that
#     silently deletes every handout the chat path produces.


def _boom(spec, art):
    raise RuntimeError("validator bug")


PROBES = [
    ("raises", _boom),
    ("returns a non-string", lambda spec, art: object()),
    ("returns an empty string", lambda spec, art: "   "),
    ("returns None", lambda spec, art: None),
]
probe_rows = []
if validate is None:
    probe_rows = [("validate absent", "import failed", [])]
else:
    _prior_check = validate.check
    try:
        for label, fake in PROBES:
            validate.check = fake
            _t, _o, _kept, _raised = run_tool(
                ok_result(artifact("deck.pptx", deck_empty(), PPTX_MIME))
            )
            probe_rows.append((label, _raised, _kept))
    finally:
        validate.check = _prior_check
check(
    "63. A4 _run never raises for any validator outcome, and keeps the file when the validator is the bug",
    len(probe_rows) == len(PROBES)
    and all(not raised for _, raised, _k in probe_rows)
    and all(kept == ["deck.pptx"] for _, _r, kept in probe_rows),
    f"{[(l, r or 'no raise', k) for l, r, k in probe_rows]}",
)

# 64. A5 -- PLAN.md 3.6 R-a, applied to this path. "Byte-identical to today"
#     cannot be asserted against code that no longer exists, so it is asserted
#     three ways that together mean it:
#
#       a. with the flag OFF, `validate.check` is NEVER CALLED -- so the executed
#          path is the pre-change one modulo an `if`;
#       b. flag OFF is byte-identical -- rendered text, outcome and persisted
#          filenames -- to flag ON with a validator that finds nothing, which is
#          the pre-change path by construction;
#       c. flag ON with the REAL validator differs on at least one fixture. An
#          off switch that changes nothing is gating nothing, and without (c)
#          this case passes over a deleted feature.
#
#     The flag is set EXPLICITLY inside a try/finally rather than read from the
#     ambient value. Case 6 above was written the other way and asserted the
#     flag-OFF contract while running flag-ON, which took a repair to notice.
F64 = [
    ("bad deck", ok_result(artifact("deck.pptx", deck_empty(), PPTX_MIME))),
    (
        "mixed",
        ok_result(
            artifact("good.csv", GOOD_CSV, "text/csv"),
            artifact("empty.pptx", deck_empty(), PPTX_MIME),
        ),
    ),
    ("no files at all", ok_result()),
]
consulted: list[str] = []


def _counting(spec, art):
    consulted.append(getattr(art, "filename", "?"))
    return None


off_rows: list[tuple] = []
noop_rows: list[tuple] = []
real_rows: list[tuple] = []
_flag64 = getattr(settings, "handout_validate_artifacts", None)
if validate is not None and _flag64 is not None:
    _prior_check = validate.check
    try:
        settings.handout_validate_artifacts = False
        off_rows = [run_tool(res)[:3] for _, res in F64]
        validate.check = _counting
        run_tool(F64[0][1])
        validate.check = lambda spec, art: None
        settings.handout_validate_artifacts = True
        noop_rows = [run_tool(res)[:3] for _, res in F64]
        validate.check = _prior_check
        real_rows = [run_tool(res)[:3] for _, res in F64]
    finally:
        validate.check = _prior_check
        settings.handout_validate_artifacts = _flag64
check(
    "64. A5 flag OFF never consults the validator and is byte-identical to a no-op one; flag ON differs",
    len(off_rows) == len(F64)
    and not consulted
    and off_rows == noop_rows
    and off_rows != real_rows,
    f"consulted={consulted} identical_to_noop={off_rows == noop_rows} "
    f"differs_on={sum(a != b for a, b in zip(off_rows, real_rows))}",
)

# ---------------------------------------------------------------------------

print()
print("=" * 74)
if failures:
    print(f"{len(failures)} FAILED:")
    for name in failures:
        print(f"  - {name}")
    sys.exit(1)
print("all deck_check cases passed")
