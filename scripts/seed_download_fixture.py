"""Seed the fixture `download_ui_check.py` needs. Idempotent.

Run:  backend/.venv/Scripts/python.exe scripts/seed_download_fixture.py
      backend/.venv/Scripts/python.exe scripts/seed_download_fixture.py --cleanup

**Why this exists as a separate script rather than inside the browser harness.**
`download_ui_check.py` runs on the GLOBAL interpreter, because Playwright is
deliberately not a backend dependency -- so it has no database session, no
models and no storage client. It can drive a browser and it cannot make a
handout. This script has the opposite capabilities, which is why the split falls
here.

**Why it does not just run a recipe.** Making a real handout through the API
costs a corpus, a retrieval and a model call -- 30-60 seconds, a provider
dependency, and a non-deterministic result. The browser proof is about the
download path, not about generation, so its fixture should be the cheapest thing
that is genuinely a deck: `python-pptx` writes three real slides here, and the
same bytes go to Postgres and to R2 through the ordinary seam.

The identity is `ui-check@groundwork.local`, the dev-login shim's user. That is a
SEPARATE row from the same address signed in through Google -- `google_sub` is
`dev|<email>` -- which is the property that makes the shim safe and the reason
this harness cannot assume it will find anything.
"""

from __future__ import annotations

import asyncio
import io
import sys
import uuid
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "backend"))

from sqlalchemy import delete, select  # noqa: E402

from app import storage  # noqa: E402
from app.db.models import (  # noqa: E402
    Agent,
    Conversation,
    Document,
    Handout,
    User,
)
from app.rag.ingest import ingest_file  # noqa: E402
from app.db.session import SessionLocal  # noqa: E402

DEV_EMAIL = "ui-check@groundwork.local"
DEV_SUB = f"dev|{DEV_EMAIL}"
AGENT_NAME = "Download check agent"

PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
PNG_MIME = "image/png"


def build_deck() -> bytes:
    """Three real slides, each with a title. Not a stub.

    A zero-slide `Presentation()` is 27,387 bytes and opens fine, which is
    exactly the artefact the robust-handouts change set was built to reject --
    so a fixture that used one would let the browser proof pass against a file
    that proves nothing.
    """
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    for index, (title, body) in enumerate(
        [
            ("Ka-band link budget", "Downlink margin is 3.2 dB at 10 degrees."),
            ("Power allocation", "The array delivers 12.4 kW at end of life."),
            ("Handover weeks", "Crew expands from eleven to nineteen."),
        ],
        start=1,
    ):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"{index}. {title}"
        frame = slide.placeholders[1].text_frame
        frame.text = body
        frame.paragraphs[0].font.size = Pt(18)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def build_chart() -> bytes:
    """A real PNG, so D3's thumbnail assertion has something to render."""
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(2.2, 2.2), dpi=100)
    ax.bar(["TT&C", "Payload", "Thermal"], [2.1, 7.8, 2.5], color="#34d399")
    ax.set_title("Power (kW)", fontsize=9)
    fig.tight_layout()
    buffer = io.BytesIO()
    fig.savefig(buffer, format="png")
    plt.close(fig)
    return buffer.getvalue()


async def seed() -> int:
    async with SessionLocal() as db:
        user = await db.scalar(select(User).where(User.google_sub == DEV_SUB))
        if user is None:
            print(
                f"  no dev user yet. Sign in once at http://localhost:5173 as "
                f"{DEV_EMAIL}, then re-run."
            )
            return 1

        agent = await db.scalar(
            select(Agent).where(
                Agent.owner_user_id == user.id, Agent.name == AGENT_NAME
            )
        )
        if agent is None:
            agent = Agent(
                id=uuid.uuid4(),
                owner_user_id=user.id,
                name=AGENT_NAME,
                description="Fixture for scripts/download_ui_check.py. Safe to delete.",
            )
            db.add(agent)
            await db.commit()
            await db.refresh(agent)
            print(f"  created agent {agent.id}")
        else:
            print(f"  reusing agent {agent.id}")

        # A corpus, because an agent with zero documents renders
        # `EmptyAgentWorkspace` -- no composer, no handout dock, and therefore no
        # download anchor for the browser proof to click. That empty state is
        # correct product behaviour and it makes D1 and D3 unmeasurable, so the
        # fixture has to clear it. One small file is enough; nothing here asks a
        # question, so retrieval quality is irrelevant.
        documents = (
            await db.scalars(select(Document).where(Document.agent_id == agent.id))
        ).all()
        if not documents:
            fixture = ROOT / "scripts" / "fixtures" / "comms-subsystem.md"
            print(f"  ingesting {fixture.name} so the workspace is not the empty state")
            await ingest_file(db, agent, fixture, uploaded_by_user_id=user.id)
            print("  ingested")

        existing = (
            await db.scalars(select(Handout).where(Handout.agent_id == agent.id))
        ).all()
        if existing:
            print(f"  {len(existing)} handout(s) already present; nothing to do")
            return 0

        # `user_id` is NOT NULL on `conversations` -- a thread belongs to the
        # person who opened it, not only to the agent. Omitting it raised an
        # asyncpg NotNullViolation naming the column, which is the readable end
        # of the error spectrum and still worth pinning here so the next reader
        # does not rediscover it.
        conversation = Conversation(
            id=uuid.uuid4(),
            agent_id=agent.id,
            user_id=user.id,
            title="Download fixture",
        )
        db.add(conversation)
        await db.commit()

        for kind, filename, mime, content in [
            ("deck", "ka-band-brief.pptx", PPTX_MIME, build_deck()),
            ("chart", "power-allocation.png", PNG_MIME, build_chart()),
        ]:
            handout_id = uuid.uuid4()
            key = None
            if storage.enabled():
                key = storage.handout_key(agent.id, handout_id, mime)
                await asyncio.to_thread(storage.put_object, key, content, mime)

            db.add(
                Handout(
                    id=handout_id,
                    agent_id=agent.id,
                    conversation_id=conversation.id,
                    created_by_user_id=user.id,
                    kind=kind,
                    title=f"Fixture {kind}",
                    filename=filename,
                    mime_type=mime,
                    byte_size=len(content),
                    content=content,
                    storage_key=key,
                    origin="recipe",
                    status="ready",
                    source_code="# fixture, written by scripts/seed_download_fixture.py",
                )
            )
            print(f"  seeded {kind}: {filename} ({len(content):,} bytes) key={key}")

        await db.commit()

    print("\n  Now run:  python scripts/download_ui_check.py")
    return 0


async def cleanup() -> int:
    async with SessionLocal() as db:
        user = await db.scalar(select(User).where(User.google_sub == DEV_SUB))
        if user is None:
            print("  no dev user; nothing to clean")
            return 0
        agents = (
            await db.scalars(
                select(Agent).where(
                    Agent.owner_user_id == user.id, Agent.name == AGENT_NAME
                )
            )
        ).all()
        for agent in agents:
            # Objects first -- the same ordering the application uses, and the
            # reason this helper does it explicitly is that the Core DELETE below
            # cascades in Postgres with no Python seeing the handout rows.
            if storage.enabled():
                removed = await asyncio.to_thread(
                    storage.delete_prefix, storage.agent_prefix(agent.id)
                )
                print(f"  deleted {removed} object(s) for {agent.id}")
            await db.execute(delete(Handout).where(Handout.agent_id == agent.id))
            await db.execute(delete(Conversation).where(Conversation.agent_id == agent.id))
            await db.delete(agent)
        await db.commit()
        print(f"  removed {len(agents)} fixture agent(s)")
    return 0


if __name__ == "__main__":
    raise SystemExit(
        asyncio.run(cleanup() if "--cleanup" in sys.argv else seed())
    )
