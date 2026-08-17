"""Standalone check of the code-interpreter sandbox. No database, no API.

The twelve cases from `new features/02-code-interpreter.md` section 8, three more
found by probing afterwards, plus one assertion about the child's environment
that cannot be made from inside the sandbox (the test would have to print
`os.environ`, which is refused).

Cases 1-7 and 15-17 are behaviour. **Cases 8-14 are the security contract**: a
change to `sandbox.py` that leaves the behaviour cases green and turns any of
8-14 red has removed a control, and this script says so as loudly as a terminal
allows. New behaviour cases are numbered from 16 so the security block never
moves.

Cases 16-17 and the deeper half of case 3 arrived on 2026-08-17, when the whole
suite was green and the sandbox was still handing out PowerPoint files nobody had
opened. The old case 3 asserted that a deck started with `PK` and was over 10,000
bytes; a `Presentation()` with **zero slides** is 27,387 bytes and starts with
`PK`. Case 3's own source writes three slides and prints "deck written with 3
slides", and neither fact was read.

Cases 13-15 exist because the first twelve all passed and the sandbox was still
porous. `import matplotlib` is allowed; `matplotlib.os` is then a public
attribute of an allowed module, and the import allowlist, the dunder rule and
the name denylist all saw nothing wrong with it. `numpy.ctypeslib.ctypes` is the
same shape and reaches native code rather than the filesystem. Neither leaked
anything -- `_minimal_env()` had already removed what the first one went looking
for -- which is the layered model behaving as designed and not a reason to leave
the outer layer open. Case 15 guards the other direction: the denylist that
blocks `matplotlib.os` must not block `matplotlib.pyplot`.

Cases 18-20 are `new features/12-robust-handouts/04-failure-legibility.md`
sections A and B -- three signals the sandbox already had and threw away:

  18  a crash AFTER a successful save must still return the file. It used to
      return `artifacts=[]`, so a deck that saved and then raised on its
      `print()` line came back as "you wrote nothing", which is both false and
      the least useful thing to tell a model that is about to retry.
  19  the timeout path must KEEP returning `artifacts=[]`. The asymmetry with 18
      is deliberate: the workdir is torn down under a process that was killed
      mid-write, and a half-written file is not a file. This case exists so that
      nobody "tidies" the two paths into agreement.
  20  stderr on a SUCCESSFUL run must reach the model, and the child's own
      `[sandbox]` housekeeping note must not. Case 20 asserts both halves,
      because a filter that drops everything passes the second alone.

Usage:
    backend/.venv/Scripts/python.exe scripts/sandbox_check.py

Exits 1 if anything fails.

ASCII only in print(). The Windows console codepage mangles em-dashes, section
signs and emoji, and this has broken three throwaway scripts in this repo
already. Anything echoed back from the sandbox goes through ascii().
"""

from __future__ import annotations

import asyncio
import io
import os
import re
import sys
import time
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "backend"))

from app.tools.interpreter import (  # noqa: E402
    TRACEBACK_TAIL_LINES,
    _render_success,
)
from app.tools.sandbox import (  # noqa: E402
    SandboxResult,
    _minimal_env,
    run,
    static_check,
)

# A fake secret planted in the parent BEFORE the child is ever spawned. If any
# of this leaks into _minimal_env(), the environment assertion at the end fails.
os.environ["OPENROUTER_API_KEY"] = "sk-or-v1-FAKE-CANARY-FOR-SANDBOX-CHECK"
os.environ["DATABASE_URL"] = "postgresql://canary:canary@example.invalid/canary"

SECRET_PATTERN = re.compile(r"KEY|SECRET|TOKEN|PASSWORD|DATABASE_URL", re.IGNORECASE)

SECURITY_CASES = {8, 9, 10, 11, 12, 13, 14}


def short(text: str, limit: int = 160) -> str:
    """ASCII-safe, single-line, bounded. Everything from the sandbox goes here."""
    flat = " ".join(text.split())
    if len(flat) > limit:
        flat = flat[:limit] + "..."
    # ascii() quotes the result and escapes anything the console cannot draw.
    return ascii(flat)


# --------------------------------------------------------------------------
# The cases
# --------------------------------------------------------------------------

CASE_1 = 'print("hello")\n'

CASE_2 = """
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

labels = ["alpha", "beta", "gamma", "delta"]
values = [3, 1, 4, 1]
fig, ax = plt.subplots(figsize=(4, 3))
ax.bar(labels, values)
ax.set_title("Sandbox chart")
fig.tight_layout()
fig.savefig("chart.png", dpi=100)
print("chart written with", len(labels), "bars")
"""

CASE_3 = """
from pptx import Presentation
from pptx.util import Pt

prs = Presentation()
layout = prs.slide_layouts[5]
count = 0
for heading in ("One", "Two", "Three"):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = heading
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    count += 1
prs.save("deck.pptx")
print("deck written with", count, "slides")
"""

CASE_4 = """
import pandas as pd

frame = pd.DataFrame({"n": [1, 2, 3, 4], "square": [1, 4, 9, 16]})
frame.to_csv("data.csv", index=False)
print("rows:", len(frame))
"""

CASE_5 = "value = 1 / 0\n"

CASE_6 = "def f(:\n    pass\n"

CASE_7 = "while True:\n    pass\n"

CASE_8 = "import os\nprint(os.environ)\n"

CASE_9 = "import socket\ns = socket.socket()\n"

CASE_10 = '__import__("os").system("echo pwned")\n'

CASE_11 = "print(().__class__.__base__.__subclasses__())\n"

CASE_12 = """
with open("big.txt", "w", encoding="utf-8") as handle:
    handle.write("x" * 9_000_000)
print("wrote a big file")
"""


def check_1(result: SandboxResult) -> str | None:
    if not result.ok:
        return f"expected ok, got error_kind={result.error_kind} error={short(result.error or '')}"
    if "hello" not in result.stdout:
        return f"stdout did not contain hello: {short(result.stdout)}"
    if result.artifacts:
        return f"expected no artifacts, got {len(result.artifacts)}"
    return None


def check_2(result: SandboxResult) -> str | None:
    if not result.ok:
        return f"expected ok, got error_kind={result.error_kind} error={short(result.error or '')} stderr={short(result.stderr)}"
    if len(result.artifacts) != 1:
        return f"expected 1 artifact, got {[a.filename for a in result.artifacts]}"
    art = result.artifacts[0]
    if art.mime_type != "image/png":
        return f"expected image/png, got {art.mime_type}"
    if art.byte_size < 2000:
        return f"png is implausibly small: {art.byte_size} bytes"
    if not art.content.startswith(b"\x89PNG"):
        return "artifact does not carry a PNG signature"
    return None


def slide_title(slide) -> str:
    """The slide's title text, or "" -- including when the layout has no title.

    `shapes.title` is None on a layout without a title placeholder, which is a
    real thing a model can produce and not an error.
    """
    holder = getattr(slide.shapes, "title", None)
    if holder is None or not getattr(holder, "has_text_frame", False):
        return ""
    return holder.text_frame.text.strip()


def open_deck(art) -> tuple[object | None, str | None]:
    """Open a harvested .pptx. Returns (Presentation, None) or (None, reason).

    Catch broadly. A .pptx is a zip, so junk raises `zipfile.BadZipFile` from the
    standard library one layer BELOW python-pptx -- catching only pptx's own
    exceptions lets it through (verified 2026-08-17). This function must never
    raise: it is the thing that decides whether a raise happened.
    """
    try:
        from pptx import Presentation

        return Presentation(io.BytesIO(art.content)), None
    except Exception as exc:  # noqa: BLE001 - see the docstring
        return None, f"does not open as a deck: {type(exc).__name__}: {exc}"


def deck_problem(art, *, min_slides: int = 1) -> str | None:
    """Is this artifact a usable slide deck? Returns a reason, or None.

    The first three assertions are what `check_3` said before 2026-08-17. They
    are kept because they are cheap and they localise the failure -- a wrong mime
    is a different bug from a corrupt file. Everything after them is the part
    that was missing: **the file is opened**.

    Deliberately independent of `app/handouts/validate.py`. This file's contract
    is the sandbox, so a regression in the application's validator must not be
    able to make the sandbox's own deck case pass. Double-entry, not duplication.
    """
    if not art.mime_type.endswith("presentationml.presentation"):
        return f"unexpected mime {art.mime_type}"
    if not art.content.startswith(b"PK"):
        return "pptx is not a zip container"
    if art.byte_size < 10_000:
        return f"pptx is implausibly small: {art.byte_size} bytes"

    prs, why = open_deck(art)
    if prs is None:
        return why
    slides = list(prs.slides)
    if len(slides) < min_slides:
        return f"deck has {len(slides)} slide(s), fewer than the {min_slides} required"
    untitled = [n for n, slide in enumerate(slides, 1) if not slide_title(slide)]
    if untitled:
        return f"slide(s) {untitled} have no title text"
    return None


def check_3(result: SandboxResult) -> str | None:
    if not result.ok:
        return f"expected ok, got error_kind={result.error_kind} error={short(result.error or '')} stderr={short(result.stderr)}"
    if len(result.artifacts) != 1:
        return f"expected 1 artifact, got {[a.filename for a in result.artifacts]}"
    art = result.artifacts[0]
    problem = deck_problem(art, min_slides=3)
    if problem is not None:
        return problem
    # CASE_3's source writes three slides and says so on stdout. Both facts were
    # sitting there unread while the case asserted only that the file was large.
    prs, _ = open_deck(art)
    slides = list(prs.slides)  # type: ignore[union-attr]
    if len(slides) != 3:
        return f"CASE_3 writes 3 slides; the deck has {len(slides)}"
    titles = [slide_title(s) for s in slides]
    if titles != ["One", "Two", "Three"]:
        return f"slide titles are {titles}, not the ones CASE_3 wrote"
    if "deck written with 3 slides" not in result.stdout:
        return f"the code's own count is missing from stdout: {short(result.stdout)}"
    return None


def check_4(result: SandboxResult) -> str | None:
    if not result.ok:
        return f"expected ok, got error_kind={result.error_kind} error={short(result.error or '')} stderr={short(result.stderr)}"
    if len(result.artifacts) != 1:
        return f"expected 1 artifact, got {[a.filename for a in result.artifacts]}"
    art = result.artifacts[0]
    if art.mime_type != "text/csv":
        return f"expected text/csv, got {art.mime_type}"
    if b"square" not in art.content:
        return f"csv body looks wrong: {short(art.content.decode('utf-8', 'replace'))}"
    return None


def check_5(result: SandboxResult) -> str | None:
    if result.ok:
        return "expected failure, got ok"
    if result.error_kind != "runtime":
        return f"expected error_kind=runtime, got {result.error_kind}"
    if "ZeroDivisionError" not in result.stderr:
        return f"ZeroDivisionError missing from stderr: {short(result.stderr)}"
    return None


def check_6(result: SandboxResult) -> str | None:
    if result.ok:
        return "expected failure, got ok"
    if result.error_kind != "syntax":
        return f"expected error_kind=syntax, got {result.error_kind}"
    if static_check(CASE_6) is None:
        return "static_check accepted a syntax error, so a process would be spawned"
    # A refusal costs a parse, not a process. 500 ms is generous for one.
    if result.duration_ms > 500:
        return f"took {result.duration_ms} ms, which suggests a process was spawned"
    if result.exit_code != -1:
        return f"expected the no-process sentinel exit_code=-1, got {result.exit_code}"
    return None


def check_7(result: SandboxResult) -> str | None:
    if result.ok:
        return "expected failure, got ok"
    if result.error_kind != "timeout":
        return f"expected error_kind=timeout, got {result.error_kind}"
    return None


def refusal_check(code: str, must_mention: tuple[str, ...]):
    """Cases 8-11: refused by the static check, before anything is spawned."""

    def check(result: SandboxResult) -> str | None:
        if result.ok:
            return "THE CODE RAN. The static check let it through."
        if result.error_kind != "import":
            return f"expected error_kind=import, got {result.error_kind}"
        message = static_check(code)
        if message is None:
            return "static_check returned None, so this reached a subprocess"
        for token in must_mention:
            if token not in message:
                return f"refusal does not name {token!r}: {short(message)}"
        if "Allowed" not in message and "Blocked" not in message:
            return f"refusal lists no alternatives, so the model cannot retry: {short(message)}"
        if result.duration_ms > 500:
            return f"took {result.duration_ms} ms, which suggests a process was spawned"
        return None

    return check


# Cases 13-15 were added after the first twelve passed, because probing found a
# hole reasoning had missed. `import matplotlib` is allowed and legitimate;
# `matplotlib.os` is then a public attribute of an allowed module, so the import
# allowlist, the dunder rule and the name denylist all saw nothing wrong.
# `numpy.ctypeslib.ctypes` is the same shape and worse -- ctypes reaches native
# code rather than merely the filesystem.
#
# Neither probe leaked anything, because `_minimal_env()` had already stripped
# the credentials the first one went looking for. That is the layered model
# working. It is still worth a regression test, because the fix (attribute names
# in `DENIED_MODULE_ATTRS`) is one frozenset that a later edit could trim without
# failing any of cases 1-12.
#
# Case 15 is the other half and matters just as much: a denylist wide enough to
# block `matplotlib.os` must not block `matplotlib.pyplot`. Without it, the safe
# response to any future scare is to add names until the tool stops working.
CASE_13 = 'import matplotlib\nprint(matplotlib.os.environ.get("PATH"))\n'
CASE_14 = "import numpy\nprint(numpy.ctypeslib.ctypes)\n"
CASE_15 = (
    "import matplotlib\n"
    "matplotlib.use('Agg')\n"
    "import matplotlib.pyplot as plt\n"
    "import pathlib\n"
    "fig, ax = plt.subplots()\n"
    "ax.bar(['a', 'b'], [1, 2])\n"
    "fig.savefig('ok.png')\n"
    "pathlib.Path('ok.csv').write_text('a,b\\n1,2\\n')\n"
    "print('both wrote')\n"
)


def check_15(result: SandboxResult) -> str | None:
    if not result.ok:
        return (
            "legitimate matplotlib + pathlib code was refused; the attribute "
            f"denylist is too wide: {short(result.error or result.stderr)}"
        )
    names = sorted(a.filename for a in result.artifacts)
    if names != ["ok.csv", "ok.png"]:
        return f"expected ok.csv and ok.png, got {names}"
    return None


def check_12(result: SandboxResult) -> str | None:
    if result.ok:
        return "a 9 MB file was accepted; the per-file cap is not enforced"
    if result.error_kind != "output":
        return f"expected error_kind=output, got {result.error_kind} ({short(result.error or '')})"
    if result.artifacts:
        return (
            f"returned {len(result.artifacts)} artifact(s) alongside the failure; "
            "a partial set would be reported to the model as success"
        )
    if "big.txt" not in (result.error or ""):
        return f"the error does not name the offending file: {short(result.error or '')}"
    return None


# Cases 16-17 were added after the whole suite was green and the sandbox was
# still handing PowerPoint files it had never opened to users.
#
# Measured 2026-08-17 against this venv's python-pptx 1.0.2: `Presentation()`
# with zero slides saves to 27,387 bytes and starts with `PK`, so it cleared
# every assertion this file made about a deck -- and `agentic_check.py`'s
# `byte_size > 0` too. `06-test-plan.md:183` has promised since it was written
# that S8 opens the .pptx with python-pptx; nothing ever did.
#
# Case 17 is the corrupt-file half. Note it must be LARGER than 10,000 bytes to
# be interesting: the 28-byte probe that started this is caught here by the size
# floor and accepted by `agentic_check.py`, so the small version proves nothing
# about *this* file. A truncated or half-written zip is the realistic shape and
# it defeats all three of the old assertions at once.
CASE_16 = """
from pptx import Presentation

prs = Presentation()
prs.save("deck.pptx")
print("deck written with", len(prs.slides), "slides")
"""

CASE_17 = """
with open("deck.pptx", "wb") as handle:
    handle.write(b"PK\\x03\\x04")
    handle.write(b"x" * 20_000)
print("wrote something shaped like a zip")
"""


def check_16(result: SandboxResult) -> str | None:
    """A deck with no slides must be rejected. It is a valid zip and a useless file."""
    if not result.ok:
        return f"expected the code to RUN, got error_kind={result.error_kind} error={short(result.error or '')}"
    if len(result.artifacts) != 1:
        return f"expected 1 artifact, got {[a.filename for a in result.artifacts]}"
    art = result.artifacts[0]
    if deck_problem(art) is None:
        return (
            f"a zero-slide deck of {art.byte_size} bytes was accepted as a usable "
            "deck; nothing in this file opens the .pptx"
        )
    return None


def check_17(result: SandboxResult) -> str | None:
    """Bytes that open as neither a zip nor a deck must be rejected."""
    if not result.ok:
        return f"expected the code to RUN, got error_kind={result.error_kind} error={short(result.error or '')}"
    if len(result.artifacts) != 1:
        return f"expected 1 artifact, got {[a.filename for a in result.artifacts]}"
    art = result.artifacts[0]
    if deck_problem(art) is None:
        return (
            f"{art.byte_size} bytes of padding behind a PK header was accepted as a "
            "usable deck; the check is reading the first two bytes, not the file"
        )
    return None


# Cases 18-20: failure legibility.
# `new features/12-robust-handouts/04-failure-legibility.md` sections A and B.
#
# All three are about signals the sandbox computes and then discards, so each one
# reads a field the suite has never looked at. 18 and 20 fail before that feature
# is built; 19 passes before and after, and its whole job is to fail if the
# asymmetry it pins is ever "corrected".

# 18 -- the file is saved, and THEN the program falls over. Measured 2026-08-17:
# `artifacts=[]`, because the non-zero-exit branch returned before `_harvest` was
# ever called. `_problem` (jobs.py:300) then fired branch 1 with a runtime error,
# and the repair turn told a model that had written a perfectly good deck that its
# code did not run. The deck is deliberately valid and three slides long, so the
# assertion below is "you got as far as this", not merely "a file came back".
CASE_18 = """
from pptx import Presentation

prs = Presentation()
layout = prs.slide_layouts[5]
for heading in ("One", "Two", "Three"):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = heading
prs.save("deck.pptx")
raise ValueError("saved the deck, then fell over")
"""

# 19 -- the same shape, except the program never stops. The file exists on disk at
# the moment the kill lands, and it must still NOT come back.
CASE_19 = """
import pathlib

pathlib.Path("partial.csv").write_text("a,b\\n1,2\\n", encoding="utf-8")
while True:
    pass
"""

# 20 -- a program that reports its own trouble on stderr and exits 0.
#
# The odd-looking `open(2, ...)` is not a trick to get around anything: `sys` is
# on DENIED_MODULE_ATTRS and `import sys` is refused (correctly), so a descriptor
# is the only route a sandboxed program has to stderr, and fd 2 is a pipe the
# parent already owns and already captures. The alternative -- provoking a numpy
# RuntimeWarning -- was measured working too and rejected: its wording moves
# between numpy versions, and this case asserts on an exact string.
#
# The message is the one 04-failure-legibility.md section A names, because that is
# the real shape: a model that wrapped a slide in try/except, lost slide 4, and
# saved the other three anyway. Exit code 0, stdout claiming success.
CASE_20 = """
handle = open(2, "w", closefd=False)
print("could not add slide 4", file=handle)
handle.flush()
print("deck written with 3 slides")
"""

# 21 -- crash AND cap breach, which case 18 made reachable for the first time.
#
# Harvesting before the exit code is read means `_harvest` can now return its
# "you blew a cap" error on the crash path, where previously only a clean run
# could get there. Two rules meet here and the case exists to pin which wins:
#
#   all-or-nothing (`sandbox.py`)  a partial set is never reported as success
#   "you got as far as this" (18)  a crash should return what was saved
#
# All-or-nothing wins, and must. A deck alongside a file that blew the run's
# budget is exactly the "half a deck" outcome the cap exists to refuse, and the
# fact that the program also crashed does not make a truncated set safe to keep.
# So `artifacts` is empty again -- but the cap message survives into the error,
# because otherwise the deck vanishes with no explanation at all, which is the
# defect 04-failure-legibility.md was written to fix arriving one layer along.
CASE_21 = """
from pptx import Presentation

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "One"
prs.save("deck.pptx")
with open("big.txt", "w", encoding="utf-8") as handle:
    handle.write("x" * 9_000_000)
raise ValueError("saved a deck and a 9 MB file, then fell over")
"""


def check_21(result: SandboxResult) -> str | None:
    if result.ok:
        return "expected failure: the program raises"
    if result.error_kind != "runtime":
        return (
            f"expected error_kind=runtime (the crash is the cause), got "
            f"{result.error_kind}"
        )
    if result.artifacts:
        return (
            f"returned {[a.filename for a in result.artifacts]} despite a cap "
            "breach; all-or-nothing must still hold on the crash path"
        )
    if "big.txt" not in (result.error or ""):
        return (
            "the error does not name the file that blew the cap, so the deck "
            f"vanished with no explanation: {short(result.error or '')}"
        )
    return None


# What the child prints to stderr about itself, before the user's code runs at
# all (`_sandbox_child.py:104-132`). On Windows it is on EVERY run.
SANDBOX_NOTE_PREFIX = "[sandbox]"


def check_18(result: SandboxResult) -> str | None:
    """A crash after a good save must return the file alongside ok=False."""
    if result.ok:
        return "expected failure: the program raises after saving"
    if result.error_kind != "runtime":
        return f"expected error_kind=runtime, got {result.error_kind}"
    if len(result.artifacts) != 1:
        return (
            f"expected the saved deck back alongside the failure, got "
            f"{[a.filename for a in result.artifacts]}; a model told it wrote "
            "nothing will rewrite code that was already correct"
        )
    art = result.artifacts[0]
    if art.filename != "deck.pptx":
        return f"expected deck.pptx, got {art.filename!r}"
    # "You got as far as this" is only a useful signal if the file is real. A
    # harvest that returned a truncated deck would be worse than returning none.
    problem = deck_problem(art, min_slides=3)
    if problem is not None:
        return f"the deck kept from the crashed run is not usable: {problem}"
    if "ValueError" not in result.stderr:
        return f"the traceback is missing from stderr: {short(result.stderr)}"
    return None


def check_19(result: SandboxResult) -> str | None:
    """The timeout path keeps returning nothing, and that asymmetry is deliberate.

    Case 18 says a crashed run gives its files back; this says a KILLED run does
    not. The difference is that the child was stopped at an arbitrary
    instruction, so a file on disk may be half a file -- and the workdir is torn
    down at `sandbox.py`'s `finally` regardless. Keeping the two paths different
    is the point, so this case is here to go red if they are ever unified.
    """
    if result.ok:
        return "expected failure, got ok"
    if result.error_kind != "timeout":
        return f"expected error_kind=timeout, got {result.error_kind}"
    if result.artifacts:
        return (
            f"the timeout path returned {[a.filename for a in result.artifacts]}; "
            "a file written by a process that was killed mid-write is not a file, "
            "and this asymmetry with case 18 is deliberate"
        )
    if "No files were kept" not in (result.error or ""):
        return f"the error no longer says the files were dropped: {short(result.error or '')}"
    return None


def check_20(result: SandboxResult) -> str | None:
    """stderr reaches the model on a SUCCESSFUL run; the child's note does not.

    Both halves are asserted, and the pair is the case. Asserting only that
    `[sandbox]` is absent would pass against a `_render_success` that reads no
    stderr at all -- which is exactly the code this case was written against.
    """
    if not result.ok:
        return f"expected ok, got error_kind={result.error_kind} error={short(result.error or '')}"

    rendered = _render_success(result)

    if "could not add slide 4" not in rendered:
        return (
            "the program's own stderr is absent from the success render: a caught "
            f"exception on an exit-0 run is seen by nobody. render={short(rendered, 300)}"
        )
    # The success render still has to be a success render.
    if "deck written with 3 slides" not in rendered:
        return f"stdout went missing from the success render: {short(rendered, 300)}"

    # Platform note: on Linux the child imports `resource` cleanly and writes no
    # note at all, so the real subprocess cannot prove the filter there. The
    # synthetic result below makes the assertion bite on every platform -- a case
    # that can only fail on one OS is a case that reports success on the other.
    if SANDBOX_NOTE_PREFIX in rendered:
        return (
            "the child's own [sandbox] housekeeping note reached the model; every "
            "local run would gain that noise and the real signal is buried again"
        )

    synthetic = SandboxResult(
        ok=True,
        exit_code=0,
        stdout="done\n",
        stderr=(
            "[sandbox] resource limits are unavailable on this platform (Windows); "
            "running with process isolation and the import hook only\n"
            "MATPLOTLIB CACHE WARNING GOES HERE\n"
        ),
        artifacts=[],
        duration_ms=10,
    )
    synth_render = _render_success(synthetic)
    if "MATPLOTLIB CACHE WARNING GOES HERE" not in synth_render:
        return "a real stderr line was dropped alongside the [sandbox] note"
    if SANDBOX_NOTE_PREFIX in synth_render:
        return "the [sandbox] note survived the filter in the synthetic case"

    # The caps. `sandbox_max_output_chars` already bounds stderr in the parent;
    # this is the wrapper's own bound, and without it a chatty-but-successful
    # program spends the next turn's prompt budget on warnings.
    noisy = SandboxResult(
        ok=True,
        exit_code=0,
        stdout="",
        stderr="\n".join(f"warning line {n}" for n in range(200)) + "\n",
        artifacts=[],
        duration_ms=10,
    )
    noisy_render = _render_success(noisy)
    if "warning line 0" in noisy_render:
        return (
            f"200 stderr lines were rendered whole; TRACEBACK_TAIL_LINES "
            f"({TRACEBACK_TAIL_LINES}) is not being applied to the success render"
        )
    if "warning line 199" not in noisy_render:
        return "the TAIL of stderr was dropped; the last lines are the useful ones"

    return None


CASES = [
    (1, 'print("hello")', CASE_1, None, check_1),
    (2, "matplotlib bar chart -> chart.png", CASE_2, None, check_2),
    (3, "python-pptx 3-slide deck -> deck.pptx", CASE_3, None, check_3),
    (4, "pandas -> data.csv", CASE_4, None, check_4),
    (5, "1/0", CASE_5, None, check_5),
    (6, "def f(:", CASE_6, None, check_6),
    (7, "while True: pass", CASE_7, 5.0, check_7),
    (8, "import os; os.environ", CASE_8, None, refusal_check(CASE_8, ("'os'",))),
    (9, "import socket", CASE_9, None, refusal_check(CASE_9, ("'socket'",))),
    (
        10,
        '__import__("os").system(...)',
        CASE_10,
        None,
        refusal_check(CASE_10, ("__import__",)),
    ),
    (
        11,
        "().__class__.__base__.__subclasses__()",
        CASE_11,
        None,
        refusal_check(CASE_11, ("__subclasses__",)),
    ),
    (12, "writes a 9 MB file", CASE_12, None, check_12),
    (
        13,
        "matplotlib.os.environ (module reached via an ALLOWED library)",
        CASE_13,
        None,
        refusal_check(CASE_13, ("'os'",)),
    ),
    (
        14,
        "numpy.ctypeslib.ctypes (native code via an ALLOWED library)",
        CASE_14,
        None,
        refusal_check(CASE_14, ("'ctypes'",)),
    ),
    (15, "legitimate code still runs after the attribute rules", CASE_15, None, check_15),
    (16, "python-pptx deck with ZERO slides", CASE_16, None, check_16),
    (17, "PK header + 20 KB of padding, named deck.pptx", CASE_17, None, check_17),
    (18, "saves deck.pptx, THEN raises", CASE_18, None, check_18),
    (19, "writes a file, then never stops (timeout keeps nothing)", CASE_19, 5.0, check_19),
    (20, "prints to stderr and exits 0", CASE_20, None, check_20),
    (21, "crash AND cap breach -- all-or-nothing still wins", CASE_21, None, check_21),
]


# --------------------------------------------------------------------------
# The environment assertion
# --------------------------------------------------------------------------


def check_env() -> str | None:
    """The child must inherit no credential-shaped variable. Asserted here.

    It cannot be asserted from inside the sandbox: the test would have to print
    `os.environ`, and `import os` is refused (case 8). So the parent checks the
    environment it is about to hand over instead.
    """
    env = _minimal_env()
    offenders = sorted(name for name in env if SECRET_PATTERN.search(name))
    if offenders:
        return f"credential-shaped variables would reach the child: {offenders}"
    if env.get("MPLBACKEND") != "Agg":
        return f"MPLBACKEND is {env.get('MPLBACKEND')!r}, not 'Agg'"
    unexpected = sorted(
        set(env)
        - {"PATH", "LANG", "LC_ALL", "HOME", "SYSTEMROOT", "TEMP", "TMP", "MPLBACKEND"}
    )
    if unexpected:
        return f"unexpected variables in the child environment: {unexpected}"
    return None


# --------------------------------------------------------------------------


async def main() -> int:
    print("=" * 78)
    print("sandbox_check -- code interpreter, 20 cases plus the environment assertion")
    print("=" * 78)
    print(f"python:  {sys.executable}")
    print(f"platform:{os.name}")
    print()

    rows: list[tuple[str, str, str, str]] = []
    failures: list[tuple[str, str]] = []

    for number, label, code, timeout_s, check in CASES:
        started = time.perf_counter()
        result = await run(code, timeout_s=timeout_s)
        wall = time.perf_counter() - started

        problem = check(result)
        # The two timeout cases: the kill has to actually land. Without this a
        # `timeout` verdict arriving thirty seconds late still reads as green.
        if problem is None and number in (7, 19):
            budget = (timeout_s or 30.0) + 2.0
            if wall > budget:
                problem = f"returned after {wall:.1f}s, past the {budget:.0f}s budget"

        status = "[ok]" if problem is None else "[FAIL]"
        detail = "" if problem is None else problem
        rows.append((str(number), label, status, f"{wall:.1f}s"))

        marker = " " if problem is None else "!"
        print(f"{marker} {status:6} {number:>2}. {label}")
        print(
            f"           ok={result.ok} exit={result.exit_code} kind={result.error_kind} "
            f"artifacts={[a.filename for a in result.artifacts]} {wall:.1f}s"
        )
        if result.error:
            print(f"           error: {short(result.error)}")
        if result.stdout.strip():
            print(f"           stdout: {short(result.stdout)}")
        if problem is not None:
            failures.append((f"case {number} ({label})", problem))
            if number in SECURITY_CASES:
                print()
                print("  " + "!" * 74)
                print("  !!! SECURITY CASE FAILED -- A CONTROL HAS BEEN REMOVED !!!")
                print(f"  !!! case {number}: {label}")
                print(f"  !!! {problem}")
                print("  !!! Read new features/02-code-interpreter.md section 5 before")
                print("  !!! changing anything else in app/tools/.")
                print("  " + "!" * 74)
            else:
                print(f"           FAILED: {problem}")
        print()

    env_problem = check_env()
    env_status = "[ok]" if env_problem is None else "[FAIL]"
    marker = " " if env_problem is None else "!"
    print(f"{marker} {env_status:6} env. _minimal_env() carries no credentials")
    print(f"           keys: {sorted(_minimal_env())}")
    if env_problem is not None:
        failures.append(("environment assertion", env_problem))
        print()
        print("  " + "!" * 74)
        print("  !!! SECURITY CASE FAILED -- THE CHILD WOULD INHERIT A SECRET !!!")
        print(f"  !!! {env_problem}")
        print("  !!! _minimal_env() is the strongest control in the design. Only ever")
        print("  !!! remove names from it.")
        print("  " + "!" * 74)
    rows.append(("env", "_minimal_env() carries no credentials", env_status, "-"))
    print()

    print("=" * 78)
    print("summary")
    print("=" * 78)
    print(f"{'#':>4}  {'result':<7} {'time':>6}  case")
    print(f"{'-' * 4}  {'-' * 7} {'-' * 6}  {'-' * 46}")
    for number, label, status, wall in rows:
        print(f"{number:>4}  {status:<7} {wall:>6}  {label}")
    print()

    if failures:
        print(f"{len(failures)} FAILED:")
        for name, problem in failures:
            print(f"  - {name}: {problem}")
        return 1

    print(f"all {len(rows)} checks passed")
    return 0


if __name__ == "__main__":
    sys.exit(asyncio.run(main()))
