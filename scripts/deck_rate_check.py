"""MEASUREMENT, not an assertion harness. What does a model-written deck look like?

`new features/12-robust-handouts/01-deck-harness-floor.md` section G, criterion A11.

WHY THIS EXISTS.

There is no recorded first-attempt success rate for the `deck` recipe anywhere in
this repository. The number quoted repo-wide -- "handouts drop to 5/6
first-attempt files at 3.7x the latency" (`CLAUDE.md:100`) -- was measured on
**6 chart recipes per arm**. Charts only. The deck is the longest and most
structured of the four programs, and `PLAN.md` R3 flags inheriting the chart
conclusion as a live risk.

Two things come out of one run:

  1. The slide-count and bullet-length distributions that set
     `handout_deck_min_slides` and `handout_deck_max_bullet_chars`. `PLAN.md`
     3.1 says both are "to be re-set from the distribution measured in feature
     01, not from instinct", and R5 says why: a validator tuned by feel is the
     `refusal_pass = 0/2` defect, a measurement punishing the behaviour the
     prompt exists to produce and then recommending its removal.
  2. R3 settled for decks specifically -- is `reasoning=False` right here, or was
     that a chart-only finding?

WHAT IS SCORED, AND WHY IT IS NOT "DID IT ERROR".

`meta["attempts"] == 1` -- the artefact was right FIRST TIME. A first-attempt
miss is invisible to an error-shaped check, because `_problem` catches it and the
retry usually succeeds, so the row still ends `ready`. It costs a whole extra
model call and sandbox run. That is `loop.md` T2 applied to a measurement rather
than to a feature, and it is the same thing `09-deepseek-agentic.md` scored.

ARM ORDER IS BALANCED, DELIBERATELY.

`PLAN.md` R8: this repo has one recorded measurement that reported confidently
about loop position rather than about its variable -- `asyncio.gather` ran third
in every trial and looked slower than its parts until it was run first. So the
pairs alternate (off,on), (on,off), (off,on), ... and every arm spends an equal
number of turns in each position.

THE BUDGET IS THE FIXTURE'S, AND IT IS HOSTILE. READ THE NUMBERS ACCORDINGLY.

`agentic_check.py`'s fixture agent runs `retrieve_k=3, rerank_top_n=2` over a
seven-chunk corpus, so a deck here is written from TWO chunks against a prompt
asking for five to eight slides. That is the most starved configuration in the
repo and it is the RIGHT place to calibrate a floor: the floor must not fire on
an honest shrink, so it should be set from the worst honest shrink there is.
It is the wrong place to ask "how good is a normal deck" -- that is feature 03's
measurement, at a real retrieval budget.

    backend/.venv/Scripts/python.exe scripts/deck_rate_check.py [--n 6]

Needs the `agentic_check.py --setup` fixture, a live database and OpenRouter,
Pinecone and Cohere. Roughly ten minutes. Writes every deck it makes to the
scratchpad so one can be opened by eye -- which is the step no harness replaces.

ASCII in print().
"""

from __future__ import annotations

import argparse
import asyncio
import io
import re
import sys
import time
import uuid
from dataclasses import dataclass, field
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "backend"))

import httpx  # noqa: E402
from pptx import Presentation  # noqa: E402
from sqlalchemy import select  # noqa: E402

from app.config import settings  # noqa: E402
from app.db.models import Agent, Handout, User  # noqa: E402
from app.db.session import SessionLocal, engine  # noqa: E402

AGENT_NAME = "Agentic Check"
BRIEF = "A short deck on the three downlink paths and what each one carries."
OUT_DIR = Path(
    r"C:\Users\Simon\AppData\Local\Temp\claude\D--vscode-ClaudeRAGAgent"
    r"\e21300d6-d9ae-41c8-94c4-a9eb62173fb1\scratchpad\decks"
)


# A bullet's citation, e.g. "Downlink runs at 26 GHz [comms-subsystem.md]".
# `DECK_PROMPT` says "End a bullet with the source filename in square brackets",
# so the share of bullets carrying one is a GROUNDING proxy that costs nothing --
# no judge, no Ragas, computable straight from the bytes.
#
# It is a proxy and not a measure: a citation can be present and wrong. What it
# does catch is the failure that matters here, a slide the model wrote about
# something it was never given. Measured 2026-08-17 at rerank_top_n=2, the share
# is BIMODAL -- eight decks at 70-76% and two at 15-19% -- so the mean alone
# hides the thing worth seeing. Report the per-deck distribution, never just the
# pooled number.
CITATION = re.compile(r"\[([^\]\[]{1,80})\]")


@dataclass
class Run:
    index: int
    arm: str  # "off" | "on"
    position: int  # 1-based position within the whole sequence
    status: str = "?"
    attempts: int | None = None
    slides: int | None = None
    titles: list[str] = field(default_factory=list)
    longest_run: int | None = None
    bullet_lengths: list[int] = field(default_factory=list)
    bullets: int = 0
    cited: int = 0
    sources: set[str] = field(default_factory=set)
    byte_size: int = 0
    wall_s: float = 0.0
    error: str = ""
    saved_to: str = ""

    @property
    def first_try(self) -> bool:
        return self.status == "ready" and self.attempts == 1

    @property
    def cite_pct(self) -> float:
        return 100.0 * self.cited / self.bullets if self.bullets else 0.0


def ascii_safe(text: str, limit: int = 90) -> str:
    return ascii(" ".join(str(text).split())[:limit])


def analyse(body: bytes, run: Run) -> None:
    """Open the deck and record its shape. Never raises."""
    run.byte_size = len(body)
    try:
        prs = Presentation(io.BytesIO(body))
    except Exception as exc:  # noqa: BLE001
        run.error = run.error or f"will not open: {type(exc).__name__}"
        run.slides = 0
        return
    slides = list(prs.slides)
    run.slides = len(slides)
    for slide in slides:
        holder = getattr(slide.shapes, "title", None)
        run.titles.append(
            holder.text_frame.text.strip()
            if holder is not None and getattr(holder, "has_text_frame", False)
            else ""
        )
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if holder is not None and shape is holder:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                run.bullet_lengths.append(len(text))
                run.bullets += 1
                found = CITATION.findall(text)
                if found:
                    run.cited += 1
                    run.sources.update(f.strip() for f in found)
    run.longest_run = max(run.bullet_lengths) if run.bullet_lengths else 0


async def one_run(client, db, base: str, agent_id: uuid.UUID, run: Run) -> None:
    started = time.perf_counter()

    # The arm. `_model_for` reads this at job time, and runs are strictly
    # sequential here, so setting it per POST is safe.
    settings.generation_reasoning = run.arm == "on"

    r = await client.post(base, json={"recipe": "deck", "brief": BRIEF})
    if r.status_code != 202:
        run.status = "post-failed"
        run.error = f"{r.status_code}: {ascii_safe(r.text)}"
        run.wall_s = time.perf_counter() - started
        return
    handout_id = r.json()["id"]

    deadline = time.monotonic() + 300
    row = None
    while time.monotonic() < deadline:
        await asyncio.sleep(3)
        rr = await client.get(f"{base}/{handout_id}")
        if rr.status_code != 200:
            continue
        row = rr.json()
        if row["status"] in ("ready", "failed"):
            break
    run.wall_s = time.perf_counter() - started

    if row is None or row["status"] not in ("ready", "failed"):
        run.status = "timeout"
        return
    run.status = row["status"]
    run.error = ascii_safe(row.get("error") or "")

    # `meta` is on no response model, by design -- read it from the row.
    async with SessionLocal() as fresh:
        h = await fresh.get(Handout, uuid.UUID(handout_id))
        if h is not None and isinstance(h.meta, dict):
            run.attempts = h.meta.get("attempts")

    if run.status != "ready":
        return
    body_r = await client.get(f"{base}/{handout_id}/download")
    if body_r.status_code != 200:
        run.error = f"download {body_r.status_code}"
        return
    analyse(body_r.content, run)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / f"deck-{run.index:02d}-reasoning-{run.arm}.pptx"
    path.write_bytes(body_r.content)
    run.saved_to = path.name


def pct(part: int, whole: int) -> str:
    return f"{part}/{whole}" + (f" ({100 * part / whole:.0f}%)" if whole else "")


def summarise(runs: list[Run], arm: str) -> str:
    arm_runs = [r for r in runs if r.arm == arm]
    if not arm_runs:
        return f"  {arm:3}  (no runs)"
    ready = [r for r in arm_runs if r.status == "ready"]
    first = [r for r in arm_runs if r.first_try]
    slides = sorted(r.slides for r in ready if r.slides is not None)
    walls = sorted(r.wall_s for r in arm_runs)
    p50 = walls[len(walls) // 2] if walls else 0.0
    bullets = sorted(b for r in ready for b in r.bullet_lengths)
    return (
        f"  reasoning {arm:3}  first-try {pct(len(first), len(arm_runs)):>12}  "
        f"ready {pct(len(ready), len(arm_runs)):>12}  "
        f"slides {slides if slides else '-'}  "
        f"p50 {p50:5.1f}s  "
        f"bullets n={len(bullets)} max={max(bullets) if bullets else '-'}"
    )


async def main_async(n: int, arms: str = "both") -> int:
    from app.auth.deps import current_user
    from app.main import app

    prior_reasoning = settings.generation_reasoning
    print("=" * 78)
    print("deck_rate_check -- what a model-written deck actually is")
    print("=" * 78)
    print(f"model:      {settings.generation_model}")
    print(f"brief:      {BRIEF}")
    print(
        f"arms:       {'reasoning off / on, pair order alternated (R8)' if arms == 'both' else f'reasoning {arms} only'}"
        f", {n} run(s) per arm"
    )
    print()

    async with SessionLocal() as db:
        agent = (await db.scalars(select(Agent).where(Agent.name == AGENT_NAME))).first()
        if agent is None:
            print(f"FIXTURE MISSING: no agent named {AGENT_NAME!r}.")
            print("Run: backend/.venv/Scripts/python.exe scripts/agentic_check.py --setup")
            return 1
        user = await db.get(User, agent.owner_user_id)
        agent_id, user_id = agent.id, agent.owner_user_id
        print(f"agent:      {agent_id}")
        print(f"budget:     retrieve_k={agent.retrieve_k} rerank_top_n={agent.rerank_top_n} "
              f"chunk_size={agent.chunk_size}   <- HOSTILE fixture; see the docstring")
        print()

    async def _fake_user():
        async with SessionLocal() as d:
            return await d.get(User, user_id)

    app.dependency_overrides[current_user] = _fake_user
    base = f"/api/agents/{agent_id}/handouts"

    # (off,on),(on,off),(off,on),... every arm spends equal time in each slot.
    sequence: list[str] = []
    if arms == "both":
        for pair in range(n):
            sequence += ["off", "on"] if pair % 2 == 0 else ["on", "off"]
    else:
        # One arm: there is no order to balance, so R8 does not apply.
        sequence = [arms] * n

    runs = [Run(index=i + 1, arm=arm, position=i + 1) for i, arm in enumerate(sequence)]
    try:
        # See the equivalent block in `agentic_check.py`: the download route
        # redirects to object storage, and an ASGI transport passed as
        # `transport=` would serve the redirected https url too. `analyse()`
        # below reads `.content` off that route, so both lines are load-bearing.
        transport = httpx.ASGITransport(app=app)
        async with httpx.AsyncClient(
            transport=transport,
            mounts={"https://": httpx.AsyncHTTPTransport()},
            follow_redirects=True,
            base_url="http://test",
            timeout=60.0,
        ) as client:
            async with SessionLocal() as db:
                for run in runs:
                    await one_run(client, db, base, agent_id, run)
                    flag = "[ok]  " if run.first_try else (
                        "[retry]" if run.status == "ready" else "[FAIL]"
                    )
                    print(
                        f"{flag} {run.index:>2}. reasoning={run.arm:3} "
                        f"status={run.status:7} attempts={run.attempts} "
                        f"slides={run.slides} longest_bullet={run.longest_run} "
                        f"cited={run.cited}/{run.bullets} ({run.cite_pct:.0f}%) "
                        f"{run.wall_s:5.1f}s {run.saved_to}"
                    )
                    if run.error:
                        print(f"        error: {run.error}")
    finally:
        settings.generation_reasoning = prior_reasoning
        app.dependency_overrides.pop(current_user, None)
        await engine.dispose()

    print()
    print("=" * 78)
    print("summary")
    print("=" * 78)
    for arm in ("off", "on"):
        print(summarise(runs, arm))
    print()

    ready = [r for r in runs if r.status == "ready"]
    slides = sorted(r.slides for r in ready if r.slides is not None)
    bullets = sorted(b for r in ready for b in r.bullet_lengths)
    untitled = [r.index for r in ready if any(not t for t in r.titles)]
    print(f"  slide counts, all arms : {slides}")
    if slides:
        print(f"  min slides observed    : {min(slides)}   <- the floor must sit AT OR BELOW this")
    if bullets:
        k = int(len(bullets) * 0.95)
        print(f"  bullet chars n={len(bullets):<4}     : min={bullets[0]} p50={bullets[len(bullets)//2]} "
              f"p95={bullets[min(k, len(bullets)-1)]} max={bullets[-1]}")
        print(f"  max-bullet threshold   : must sit ABOVE {bullets[-1]} to not fire on these")
    print(f"  decks with an untitled slide: {untitled or 'none'}")

    # GROUNDING. Report the per-deck spread, never just the pooled share -- the
    # 2026-08-17 baseline was bimodal and the mean hid it completely.
    cites = sorted(r.cite_pct for r in ready)
    pooled_b = sum(r.bullets for r in ready)
    pooled_c = sum(r.cited for r in ready)
    sources: set[str] = set()
    for r in ready:
        sources |= r.sources
    if cites:
        print()
        print(f"  cited bullets, pooled  : {pooled_c}/{pooled_b} "
              f"({100 * pooled_c / pooled_b:.1f}%)" if pooled_b else "  cited bullets: none")
        print(f"  per-deck cite %        : {[round(c) for c in cites]}")
        print(f"  poorly grounded (<50%) : {sum(1 for c in cites if c < 50)}/{len(cites)} decks")
        print(f"  distinct sources cited : {sorted(sources) or 'NONE'}")
    print()
    print(f"  decks written to: {OUT_DIR}")
    print("  OPEN ONE IN POWERPOINT. No harness replaces that step.")
    return 0


def main() -> int:
    p = argparse.ArgumentParser(description="Measure the deck recipe. Not an assertion harness.")
    p.add_argument("--n", type=int, default=6, help="runs per arm (default 6)")
    p.add_argument(
        "--arms",
        choices=("both", "off", "on"),
        default="both",
        help="which reasoning arms to run. `off` isolates the budget variable "
             "at the shipped default and costs ~3x less wall clock.",
    )
    args = p.parse_args()
    return asyncio.run(main_async(args.n, args.arms))


if __name__ == "__main__":
    sys.exit(main())
