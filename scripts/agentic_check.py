"""End-to-end check for the agent loop, the tools, and Handouts.

Layer 2 of the three in `new features/06-test-plan.md`. Layer 1
(`scripts/sandbox_check.py`) needs nothing but the venv; layer 3 is Playwright in
a real browser. This one sits between them: real database, real OpenRouter, real
Pinecone, no browser.

    python scripts/agentic_check.py --setup     # throwaway agent + fixture corpus
    python scripts/agentic_check.py --run       # the scenarios
    python scripts/agentic_check.py --cleanup   # namespace, rows, handouts

`--only` is a substring match on a scenario's FUNCTION NAME and, since
2026-08-17, it reaches the handout-route scenarios too. It used to not: the
whole HTTP block was gated on `if not only:`, so every assertion about a real
handout job was all-or-nothing with the twenty-minute suite, and the five
criteria owed by `new features/12-robust-handouts/` could not be executed
without re-asking twenty questions of a live model. A check that is expensive
to run is a check that does not get run.

**The fixture corpus is two files on purpose.** CLAUDE.md records that context
precision and recall both scoring exactly 1.0 on the existing single-file corpus
does not mean retrieval is excellent -- it means retrieval cannot fail, because
there is only one chunk to return. A multi-hop test against that corpus would
pass without exercising anything. `scripts/fixtures/` holds a power briefing and
a comms briefing that overlap in exactly one place (the communications power
allocation appears in both, with different numbers -- an allocation in one and a
measured average in the other), which is the smallest corpus that can tell a
one-search answer from a two-search one.

**Scenarios S1, S7, S20 and S25 are the regression tests and matter most.**
Everything else checks that a new feature works; those four check that it did not
eat the old one. S1 asserts an agent with tools off produces exactly the six
pre-existing trace event types -- and, since 2026-08-16, that REWRITE is among
them, because the subset check alone passed whether the rewriter ran or had been
deleted. S7 asserts that giving a model tools did not turn "the corpus does not
cover this" into an invention. S20 is S1 for the orchestrator, and S25 is S7 for
the self-check: a persona designed to answer with a question must not be graded
for citing nothing.

Layer 1 for the orchestrator is `scripts/route_specialist_check.py` -- mention
parsing, the alias table, the self-check trigger and the critic's carve-out, with
no database and no model, in seconds. Read a red row there before a red row here.

HTTP routes are exercised through an ASGI transport rather than a running server,
with `current_user` overridden. `owned_agent` still runs for real, so the tenancy
wiring is genuinely under test -- only the identity assertion is stubbed, which
is the same split `POST /api/auth/dev-login` makes.
"""

from __future__ import annotations

import argparse
import asyncio
import re
import sys
import time
import uuid
from dataclasses import dataclass, field
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "backend"))

from sqlalchemy import delete, func, select  # noqa: E402

from app import storage as storage_mod  # noqa: E402
from app.config import settings  # noqa: E402
from app.db.models import (  # noqa: E402
    Agent,
    Chunk,
    Conversation,
    Document,
    Handout,
    IngestionRun,
    Query,
    QueryChunk,
    TraceEvent,
    User,
)
from app.db.session import SessionLocal, engine  # noqa: E402
from app.rag.ingest import ingest_file  # noqa: E402
from app.rag.retriever import aretrieve, get_vector_store  # noqa: E402

SEED_SUB = "agentic-check-local"
AGENT_NAME = "Agentic Check"
FIXTURES = ROOT / "scripts" / "fixtures"

# The six event types that existed before the tool loop. S1 asserts a
# tools-off turn produces a subset of exactly these and nothing else.
CLASSIC_EVENTS = {"RETRIEVE", "SCORE_CHECK", "REWRITE", "RERANK", "GENERATE", "REFUSE"}

# The three the orchestrator adds, and the nine that existed before it. S20 is
# to orchestration what S1 is to tools: with the feature off the trace must hold
# nothing from this set, and the assertion is written as a set difference so a
# FOURTH new event type is caught too. A subset check against a set that already
# contained the new names would pass whether the feature was off or merely
# broken -- the hole S1 carried until 2026-08-16.
ORCHESTRATOR_EVENTS = {"ROUTE", "DELEGATE", "SELF_CHECK"}
PRE_ORCHESTRATOR_EVENTS = CLASSIC_EVENTS | {"TOOL_CALL", "TOOL_RESULT", "TOOL_ERROR"}


# --------------------------------------------------------------------------
# Console
# --------------------------------------------------------------------------

def ascii_safe(text: str) -> str:
    """Force text to ASCII for the Windows console.

    The rule is usually read as being about string literals; the corpus and the
    model output are the bigger hazard. An em-dash in a generated answer raises
    UnicodeEncodeError under cp1252, several layers from anything this file
    wrote, and it has killed three throwaway scripts in this repo already.
    """
    return text.encode("ascii", "replace").decode("ascii")


def rule(title: str) -> None:
    print(f"\n{'=' * 74}\n{title}\n{'=' * 74}")


def preview(text: str, width: int = 160) -> str:
    return ascii_safe(" ".join(text.split())[:width])


@dataclass
class Outcome:
    name: str
    ok: bool
    detail: str = ""
    notes: list[str] = field(default_factory=list)
    rate_limited: bool = False
    #: Neither green nor red: the assertion did not run, because its input never
    #: arrived. Ported from `ui_check.py`'s `Results.unmeasured` on 2026-08-17,
    #: after an audit found five checks in this file that report green when their
    #: precondition is absent -- S8b vanishing inside `if ready:`, S11's
    #: `not leaked` over zero captured statements, S5's disjunction collapsing
    #: when no tool was called, S10 comparing two empty lists, S23's `all([])`.
    #:
    #: A two-state harness has to call those PASS, and `build.md` 5 is blunt
    #: about what that costs: "a check that cannot fail reports success, and it
    #: reports it in green, forever."
    unmeasured: bool = False


def unmeasured(name: str, detail: str) -> "Outcome":
    """The assertion could not run. Not a pass, not a failure."""
    return Outcome(name, False, detail, unmeasured=True)


#: The floor a deck must clear here. `DECK_PROMPT` asks for five to eight slides
#: and also tells the model to use only what the material supports, so a thin
#: fixture corpus SHOULD produce a short deck -- firing on that is the
#: `refusal_pass = 0/2` defect, a measurement punishing the behaviour the prompt
#: exists to produce. Three is "somebody made a deck", not "the deck is good".
DECK_MIN_SLIDES = 3


def artifact_problem(recipe: str, body: bytes) -> str | None:
    """Open the downloaded handout. Returns a reason it is unusable, or None.

    These are the four content checks `new features/06-test-plan.md:183` has
    promised since it was written -- "the .pptx opens with python-pptx; the .png
    opens with PIL; the .csv parses; the .md is non-empty" -- and which S8 did
    not perform. It asserted `status == "ready" and byte_size > 0`, which is
    satisfied by a zero-slide deck (27,387 bytes, measured) and by 28 bytes of
    PK junk.

    Deliberately independent of `app/handouts/validate.py`, which does not exist
    yet and which this must be able to outlive: a regression in the application's
    validator must not be able to make this scenario pass. Same double-entry
    reasoning as `sandbox_check.deck_problem`.

    Never raises. A harness that dies on a malformed artefact reports nothing
    about the artefact.
    """
    if not body:
        return "the download was empty"
    try:
        if recipe == "deck":
            import io as _io

            from pptx import Presentation

            prs = Presentation(_io.BytesIO(body))
            slides = list(prs.slides)
            if len(slides) < DECK_MIN_SLIDES:
                return f"deck has {len(slides)} slide(s), fewer than {DECK_MIN_SLIDES}"
            untitled = []
            for n, slide in enumerate(slides, 1):
                holder = getattr(slide.shapes, "title", None)
                text = (
                    holder.text_frame.text.strip()
                    if holder is not None and getattr(holder, "has_text_frame", False)
                    else ""
                )
                if not text:
                    untitled.append(n)
            if untitled:
                return f"slide(s) {untitled} have no title"
            return None

        if recipe == "chart":
            import io as _io

            from PIL import Image

            if not body.startswith(b"\x89PNG"):
                return "no PNG signature"
            image = Image.open(_io.BytesIO(body))
            image.verify()
            if min(image.size) <= 1:
                return f"image is {image.size}, which is not a chart"
            return None

        if recipe == "table":
            import csv as _csv
            import io as _io

            rows = list(_csv.reader(_io.StringIO(body.decode("utf-8", "replace"))))
            rows = [r for r in rows if any(cell.strip() for cell in r)]
            if len(rows) < 2:
                return f"csv has {len(rows)} non-empty row(s); a header alone is not a table"
            return None

        if recipe == "sheet":
            text = body.decode("utf-8", "replace").strip()
            if not text:
                return "the study sheet is empty"
            if "#" not in text:
                return "the study sheet has no markdown heading"
            return None
    except Exception as exc:  # noqa: BLE001 - see the docstring
        return f"could not open the artefact: {type(exc).__name__}: {exc}"
    return None


def slide_count(body: bytes) -> int | None:
    """How many slides the downloaded deck HAS, or None if it will not open.

    Deliberately separate from `artifact_problem`, which answers "is this
    usable" against a floor. A count and a verdict are different questions and
    the scenarios below need both: S28 has to distinguish "the validator fired"
    from "the model happened to clear a raised floor", and S8's feature-05 half
    has to compare the stored `preview_text` against the real number. A verdict
    cannot tell either pair apart.

    Never raises, for `artifact_problem`'s reason: a harness that dies on a
    malformed artefact reports nothing about the artefact.
    """
    try:
        import io as _io

        from pptx import Presentation

        return len(list(Presentation(_io.BytesIO(body)).slides))
    except Exception:  # noqa: BLE001 - see the docstring
        return None


# `validate._numbered` opens a deck outline with `"6 slides"` (or `"1 slide"`,
# because `_plural` exists precisely so the card does not read "1 slides").
# Anchored at the start, so a slide TITLE containing "12 slides" cannot be read
# as the count.
_PREVIEW_SLIDES = re.compile(r"^\s*(\d+)\s+slides?\b")


def preview_slide_count(preview_text: str | None) -> int | None:
    """The slide count the STORED preview claims, or None if it claims none.

    Feature 05's whole product claim is that a user can tell an empty deck from
    a real one without opening PowerPoint. That claim is only true if the number
    on the card is the number in the file, and nothing else in this repo
    compares the two -- `deck_check.py` asserts `outline()` against fixtures it
    built itself, which cannot see a job that wrote the wrong field, wrote the
    model's stdout caption instead, or wrote a preview from attempt 1 beside an
    artefact from attempt 2.

    **The anchor is what makes it able to fail.** Measured 2026-08-17:

        "7 slides\\n1. Three Downlink Paths"                    ->    7
        "1 slide\\n1. Only one"                                 ->    1
        "0 slides"                                             ->    0
        "Saved deck.pptx with 7 slides on the three downlink"   -> None
        "deck written with 3 slides"                           -> None

    The last two are the model's own `print()` caption -- exactly what
    `preview_text` held before feature 05, and exactly what `_preview_for`
    falls back to. So a regression that dropped the outline and kept the
    caption reads as `preview_claims=None` and S8c goes red, rather than
    matching a number out of the middle of a sentence and going green.
    """
    if not preview_text:
        return None
    match = _PREVIEW_SLIDES.match(preview_text)
    return int(match.group(1)) if match else None


# Phrases that mean "the upstream provider refused us", not "the code is wrong".
#
# This suite makes roughly twenty reranking calls in a couple of minutes -- every
# scenario retrieves at least once, several retrieve twice, and each of the four
# recipes retrieves to gather its material. **Cohere's trial key allows ten API
# calls per minute**, so running the suite twice in quick succession trips it,
# and the failure arrives as a handout stuck at `failed` and a scenario throwing.
#
# Reporting that as `[FAIL]` would be the same defect this project has already
# recorded twice: CLAUDE.md's note that `METRIC_TIMEOUT_S` silently doubled as a
# quota-retry ceiling, so "a rate limit and a hang are indistinguishable on the
# card, and they need opposite fixes". A red row that means "wait sixty seconds"
# sends the reader to debug code that is working.
# Both spellings of every phrase, because the matched text is
# `f"{type(exc).__name__}: {exc}"` and an SDK's exception CLASS is camel-cased
# with the spaces removed. Cohere raises `TooManyRequestsError`, which the
# spaced "too many requests" does not match -- so the first version of this list
# printed `[FAIL]` for a rate limit anyway, which is the exact bug it was
# written to fix.
#
# That is the same gap as the refusal marker list in `app/rag/refusal.py`,
# arrived at independently ten minutes later: a substring test matched one
# spelling of a phrase and missed the variant the model (or the SDK) actually
# used. The lesson from there applies here -- add the shape, not the string you
# just saw.
RATE_LIMIT_PHRASES = (
    "trial key",
    "rate limit",
    "ratelimit",
    "429",
    "too many requests",
    "toomanyrequests",
    "resource_exhausted",
    "resourceexhausted",
    "quota",
    "overloaded",
    "service unavailable",
    "serviceunavailable",
    "503",
    "502",
)


def _flag(outcome: "Outcome") -> str:
    """FOUR states, not two.

    `[rate]` says the provider refused, so re-run in a minute rather than opening
    an editor. `[warn]` says the assertion never ran at all, which is the state
    that was previously spelled "green".

    Neither fails the suite -- see `_is_failure`. A red row that means "wait
    sixty seconds" or "the fixture produced no input" sends its reader to debug
    working code, and teaches them to ignore red.
    """
    if outcome.ok:
        return "[ok]  "
    if outcome.unmeasured:
        return "[warn]"
    return "[rate]" if outcome.rate_limited else "[FAIL]"


def _is_failure(outcome: "Outcome") -> bool:
    """A genuine defect, as opposed to an environment refusal or a missing input.

    Until 2026-08-17 the exit code was `not o.ok`, so a `[rate]` row exited
    non-zero -- which is the opposite of what this file's own comment and
    CLAUDE.md both describe, and it meant a Cohere 429 looked exactly like a
    broken handout job. Treat both non-defect states as UNMEASURED, never as
    passing, and never as failing.
    """
    return not outcome.ok and not outcome.rate_limited and not outcome.unmeasured


def is_rate_limited(text: str) -> bool:
    lowered = text.lower()
    return any(phrase in lowered for phrase in RATE_LIMIT_PHRASES)


def record(
    outcome: "Outcome", outcomes: list["Outcome"], took: float | None = None
) -> None:
    """Re-flag, print, append. ONE implementation, used by both scenario loops.

    The scenario loop and the handout-route loop each did this inline and the
    two had already drifted once -- the summary spelled `_flag` out by hand and
    printed `[FAIL]` for a row the run had printed `[rate]`, which is the
    defect recorded at `_flag`. Two more copies is how that comes back.

    Printing here rather than after the block is the other half: a handout
    scenario is minutes of silence, and a reader watching a blank console
    cannot tell a slow deck from a hung one.
    """
    if not outcome.ok and is_rate_limited(outcome.detail):
        outcome.rate_limited = True
    print(
        f"  {_flag(outcome)} {outcome.name}"
        + (f"  ({took:.1f}s)" if took is not None else "")
    )
    print(f"         {ascii_safe(outcome.detail)}")
    outcomes.append(outcome)


# --------------------------------------------------------------------------
# Fixture agent
# --------------------------------------------------------------------------

async def get_or_create_agent(db) -> tuple[User, Agent]:
    user = await db.scalar(select(User).where(User.google_sub == SEED_SUB))
    if user is None:
        user = User(
            id=uuid.uuid4(),
            google_sub=SEED_SUB,
            email="agentic-check@localhost",
            name="Agentic Check",
            role="user",
        )
        db.add(user)
        await db.flush()

    agent = await db.scalar(
        select(Agent).where(Agent.owner_user_id == user.id, Agent.name == AGENT_NAME)
    )
    if agent is None:
        agent = Agent(
            id=uuid.uuid4(),
            owner_user_id=user.id,
            name=AGENT_NAME,
            description="Throwaway agent for the agentic end-to-end check.",
            tools_enabled=True,
            max_tool_steps=3,
            # 250/40 rather than the 800/120 default, and this is the difference
            # between a suite that tests something and one that reports success
            # for free.
            #
            # The two fixtures are ~550 tokens each. At the default chunk size
            # they produce ONE CHUNK EACH -- so a corpus of two chunks, against
            # `retrieve_k=20`, means a single retrieval returns the entire corpus
            # and no question can possibly require a second search. S3 would
            # pass without the tool ever being needed.
            #
            # This is the same trap CLAUDE.md records for context precision:
            # 1.000 on a single-chunk corpus is not excellent retrieval, it is
            # retrieval that cannot fail. 250/40 splits each briefing into
            # roughly three chunks, and `rerank_top_n=3` then means one
            # retrieval structurally cannot hold both topics.
            chunk_size=250,
            chunk_overlap=40,
            # And retrieval is deliberately STARVED: 3 candidates reranked to 2,
            # against the 20/3 default.
            #
            # The first attempt at this suite used the defaults and S3 failed
            # with `searches=0` while already citing both files -- because
            # `retrieve_k=20` over a 7-chunk corpus makes every chunk a
            # candidate, and a reranker picking 3 of 7 will happily span both
            # documents. The model was right not to search again; there was
            # nothing left to find. A passing S3 there would have measured
            # nothing.
            #
            # 3/2 recreates the condition a second search exists FOR: a two-part
            # question whose halves are semantically distant gets one half back
            # and has to go looking for the other. That is a real property of a
            # large corpus reproduced on a small one, which is the only honest
            # way to test it without shipping a large fixture.
            retrieve_k=3,
            rerank_top_n=2,
        )
        db.add(agent)
        await db.flush()

    await db.commit()
    return user, agent


async def setup(db) -> int:
    user, agent = await get_or_create_agent(db)
    rule("Setup")
    print(f"  agent      : {agent.id}")
    print(f"  namespace  : {agent.namespace}")

    existing = await db.scalar(
        select(Document).where(Document.agent_id == agent.id).limit(1)
    )
    if existing is not None:
        print("  corpus already ingested; skipping. Use --cleanup first to rebuild.")
        return 0

    for path in sorted(FIXTURES.glob("*.md")):
        started = time.perf_counter()
        run = await ingest_file(db, agent, path, uploaded_by_user_id=user.id)
        took = int((time.perf_counter() - started) * 1000)
        print(f"  ingested {path.name}: {run.chunk_count} chunks in {took} ms")

    return 0


async def cleanup(db) -> int:
    user = await db.scalar(select(User).where(User.google_sub == SEED_SUB))
    if user is None:
        print("Nothing to clean up.")
        return 0

    rule("Cleanup")
    agents = (await db.scalars(select(Agent).where(Agent.owner_user_id == user.id))).all()
    for agent in agents:
        # The namespace first. A leaked Pinecone namespace is a real cost -- the
        # Builder plan's 1,000-namespace cap IS the maximum number of agents this
        # deployment can ever hold, and it binds long before storage does.
        try:
            get_vector_store(agent)._index.delete(delete_all=True, namespace=agent.namespace)
            print(f"  deleted namespace {agent.namespace}")
        except Exception as exc:  # namespace may never have been written
            print(f"  namespace {agent.namespace}: {type(exc).__name__} (already absent?)")

        # The object prefix, for exactly the reason the namespace above goes
        # first -- and this helper is a SECOND leak site beyond the one the API
        # route handles, because it bypasses the route entirely. The Core DELETE
        # below removes every handout row without any Python seeing it, so
        # without this the suite quietly fills a bucket with the artefacts of
        # every run that was ever cleaned up. `--cleanup` looks like it left
        # nothing behind, which is what makes it worth doing here rather than
        # trusting the route.
        try:
            removed = storage_mod.delete_prefix(storage_mod.agent_prefix(agent.id))
            if removed:
                print(f"  deleted {removed} stored file(s) for {agent.id}")
        except Exception as exc:  # noqa: BLE001
            print(f"  stored files for {agent.id}: {type(exc).__name__} (leaked)")

        await db.execute(delete(Handout).where(Handout.agent_id == agent.id))
        docs = (await db.scalars(select(Document).where(Document.agent_id == agent.id))).all()
        for doc in docs:
            await db.execute(delete(Chunk).where(Chunk.document_id == doc.id))
            await db.execute(delete(IngestionRun).where(IngestionRun.document_id == doc.id))
        await db.execute(delete(Document).where(Document.agent_id == agent.id))
        await db.execute(delete(Conversation).where(Conversation.agent_id == agent.id))
        await db.execute(delete(Query).where(Query.agent_id == agent.id))
        await db.delete(agent)

    await db.delete(user)
    await db.commit()
    print("  deleted seed user, agents, documents, chunks, conversations, handouts")
    return 0


# --------------------------------------------------------------------------
# Turn helper
# --------------------------------------------------------------------------

async def ask(db, agent: Agent, user: User, question: str,
              conversation: Conversation | None = None,
              emit=None):
    """One turn through the real engine, returning (AskOut, events, conversation).

    Goes through `ask.run_turn` rather than `pipeline.answer_question`, because
    the trace rows are half of what is being asserted and only `run_turn` writes
    them.

    `emit` is the streaming seam, and only S26 uses it. It is here rather than in
    a second helper because `answer_reset` is a WIRE frame with no trace row --
    `app/rag/events.py` says in as many words that SSE names must never become
    trace event types -- so a scenario asserting that a draft was not discarded
    has nowhere else to look. With `emit=None` every branch below `run_turn` is
    the line it already was, which is what S1 depends on.
    """
    from app.api.ask import run_turn

    if conversation is None:
        conversation = Conversation(
            id=uuid.uuid4(), agent_id=agent.id, user_id=user.id, title=None
        )
        db.add(conversation)
        await db.flush()

    out = await run_turn(
        db, agent=agent, user=user, session=None,
        conversation=conversation, question=question, emit=emit,
    )

    events = (
        await db.scalars(
            select(TraceEvent)
            .where(TraceEvent.query_id == out.query_id)
            .order_by(TraceEvent.step_index.asc())
        )
    ).all()
    return out, list(events), conversation


def event_types(events) -> list[str]:
    return [e.event_type for e in events]


def payloads_of(events, kind: str) -> list[dict]:
    return [e.payload or {} for e in events if e.event_type == kind]


def frame_collector():
    """(frames, emit) -- an `Emit` that records instead of sending.

    Frames are `(name, payload)` in order. Used only by S26; see `ask()`.
    """
    frames: list[tuple[str, dict]] = []

    async def emit(name: str, payload: dict) -> None:
        frames.append((name, payload))

    return frames, emit


# --------------------------------------------------------------------------
# Scenarios
# --------------------------------------------------------------------------

async def s1_classic_path(db, agent, user) -> Outcome:
    """Tools OFF must produce exactly the pre-existing trace shape.

    The regression test. Everything else here checks the new feature works; this
    one checks it did not change the old one.
    """
    agent.tools_enabled = False
    await db.commit()
    try:
        out, events, _ = await ask(
            db, agent, user, "How much power do the solar arrays generate?"
        )
        seen = set(event_types(events))
        stray = seen - CLASSIC_EVENTS
        # **`REWRITE in seen` was added 2026-08-16, and its absence was a real
        # hole.** `stray` is a SUBSET check against a set that already contains
        # REWRITE, so this scenario passed identically whether the rewriter ran
        # or had been deleted -- it could see the feature neither appear nor
        # disappear. It runs on every turn now, first ones included, so its row
        # is an invariant of a turn rather than a property of a thread.
        ok = not stray and "RETRIEVE" in seen and "GENERATE" in seen and "REWRITE" in seen
        return Outcome(
            "S1 classic path unchanged", ok,
            f"events={sorted(seen)}" + (f" STRAY={sorted(stray)}" if stray else ""),
        )
    finally:
        agent.tools_enabled = True
        await db.commit()


async def s2_no_reflex_tools(db, agent, user) -> Outcome:
    """Tools ON, single-topic question: the model must not call tools reflexively."""
    out, events, _ = await ask(
        db, agent, user, "What is the usable capacity of one battery module?"
    )
    calls = payloads_of(events, "TOOL_CALL")
    return Outcome(
        "S2 no reflex tool use", len(calls) == 0,
        f"tool_calls={len(calls)} answer={preview(out.answer, 90)}",
    )


async def s3_multi_hop(db, agent, user) -> Outcome:
    """A question spanning both fixture files should provoke a second search.

    The halves are chosen to be semantically DISTANT, which the first version of
    this scenario got wrong. It asked to compare the communications power
    allocation with the measured communications draw -- two facts in two
    different files, but both of them plainly about communications power, so a
    single embedding retrieved both and the model correctly did not search again.

    Battery modules and science-data storage share no vocabulary and no topic --
    except that both fixtures use the word "storage" as a heading, which is
    exactly how the SECOND version of this scenario also failed with
    `searches=0`: reranking to 2 returned one chunk from each file and answered
    the whole question in one pass.

    So this scenario now starves retrieval ITSELF rather than relying on the
    agent's configuration, the same way S1 owns `tools_enabled` and S6 owns
    `max_tool_steps`. `retrieve_k=1` is the smallest configuration in which the
    property is testable at all: exactly one chunk comes back, so a two-part
    question CANNOT be answered from the first retrieval, and searching again is
    the only way through. Anything larger, on a seven-chunk corpus, leaves the
    model correctly deciding it already has what it needs -- which is good
    behaviour and a worthless test.

    The general lesson, which cost three attempts: **a test that the feature is
    needed has to make the feature necessary.** Twice here the scenario passed
    the model a question it could already answer and then asserted it had worked
    harder.
    """
    prior_k, prior_n = agent.retrieve_k, agent.rerank_top_n
    agent.retrieve_k, agent.rerank_top_n = 1, 1
    await db.commit()
    try:
        out, events, _ = await ask(
            db, agent, user,
            "How many battery modules does the platform carry, and separately, "
            "how much onboard storage do the science instruments have?",
        )
    finally:
        agent.retrieve_k, agent.rerank_top_n = prior_k, prior_n
        await db.commit()
    calls = [p for p in payloads_of(events, "TOOL_CALL") if p.get("tool") == "search_corpus"]
    results = payloads_of(events, "TOOL_RESULT")
    new_chunks = sum(int(p.get("new_chunks") or 0) for p in results)
    files = {c.filename for c in out.citations}
    ok = len(calls) >= 1 and new_chunks > 0 and len(files) >= 2
    return Outcome(
        "S3 multi-hop search", ok,
        f"searches={len(calls)} new_chunks={new_chunks} files={sorted(files)}",
    )


async def s4_chart_from_chat(db, agent, user) -> Outcome:
    """Asking for a chart mid-conversation should produce a handout."""
    out, events, convo = await ask(
        db, agent, user,
        "List the power allocation for each subsystem, then draw me a bar chart "
        "of those figures as a PNG.",
    )
    calls = [p for p in payloads_of(events, "TOOL_CALL") if p.get("tool") == "run_python"]
    rows = (
        await db.scalars(
            select(Handout).where(Handout.conversation_id == convo.id)
        )
    ).all()
    ready = [r for r in rows if r.status == "ready"]
    has_png = any(r.mime_type == "image/png" for r in ready)
    has_code = all(bool(r.source_code) for r in ready)
    ok = len(calls) >= 1 and has_png and has_code
    return Outcome(
        "S4 chart handout from chat", ok,
        f"run_python={len(calls)} handouts={len(rows)} ready={len(ready)} "
        f"png={has_png} source_code={has_code}",
    )


async def s5_self_correction(db, agent, user) -> Outcome:
    """A tool failure must be recoverable, not fatal.

    Provoked rather than mocked: asking for a library that is not on the import
    allowlist makes the static check refuse, the model reads the refusal (which
    NAMES what is allowed) and retries. That round trip is the single most
    valuable behaviour a code interpreter has and it is worth a scenario.
    """
    out, events, _ = await ask(
        db, agent, user,
        "Use the seaborn library to plot the three downlink data rates, and if "
        "seaborn is unavailable use whatever plotting library you do have.",
    )
    errors = payloads_of(events, "TOOL_ERROR")
    calls = payloads_of(events, "TOOL_CALL")
    detail = f"calls={len(calls)} errors={len(errors)} answered={bool(out.answer)}"

    # `len(errors) == 0 or len(calls) > len(errors)` short-circuits when the
    # model called no tool at all, collapsing the whole assertion to
    # `bool(out.answer)` -- the error-shaped test T2 forbids, sitting in the
    # scenario whose docstring calls self-correction the most valuable behaviour
    # a code interpreter has. Recovery cannot be observed if nothing failed.
    if not errors:
        return unmeasured(
            "S5 tool failure is recoverable",
            f"{detail} -- no tool error was provoked, so recovery was not exercised",
        )
    ok = bool(out.answer) and len(calls) > len(errors)
    return Outcome("S5 tool failure is recoverable", ok, detail)


async def s6_step_budget(db, agent, user) -> Outcome:
    """max_tool_steps=1 on a two-search question still returns a CLEAN answer.

    **This scenario shipped a user-visible bug while green, and the second
    assertion is the fix.** It asserted `bool(out.answer)` -- did a turn come
    back -- which is the error-shaped test `new features/loop.md` T2 exists to
    forbid, and it is exactly the shape that misses.

    What it missed, found in a browser on 2026-08-16: when the budget runs out
    the loop re-invokes with `tool_choice="none"`, and
    `deepseek/deepseek-v4-flash-0731` still wanted to search, so it emitted its
    own tool-call syntax into the CONTENT channel --

        Let me try searching for the thermal rejection budget document
        <|DSML|tool_calls> <|DSML|invoke name="search_corpus"> ...

    -- and the user read it. Nothing raised, `out.answer` was non-empty and long,
    and this scenario stayed green through it. This is the ONLY scenario that
    exercises the forced-final-answer path, so it is the only one that could have
    caught it.

    The assertion is now "is the answer free of machinery", which is a statement
    about the OUTCOME. `_LEAKED_TOOL_MARKUP` is imported rather than retyped, so a
    change to the sentinel cannot leave this scenario testing the old one.
    """
    from app.rag.agent_loop import _LEAKED_TOOL_MARKUP

    agent.max_tool_steps = 1
    await db.commit()
    try:
        out, events, _ = await ask(
            db, agent, user,
            "Compare the Ka-band availability gaps with the battery replacement "
            "schedule, and explain how the two interact.",
        )
        gen = payloads_of(events, "GENERATE")
        stopped = gen[0].get("stopped_reason") if gen else None
        leaked = _LEAKED_TOOL_MARKUP.search(out.answer or "")
        ok = bool(out.answer) and leaked is None
        return Outcome(
            "S6 step budget forces a clean answer", ok,
            f"stopped_reason={stopped} answer_chars={len(out.answer)} "
            f"leaked_markup={'YES at char ' + str(leaked.start()) if leaked else 'no'}",
        )
    finally:
        agent.max_tool_steps = 3
        await db.commit()


async def s7_refusal_survives(db, agent, user) -> Outcome:
    """Tools must not turn "not covered" into invention.

    The probe is chosen the way CLAUDE.md says a good refusal probe is chosen:
    both fixtures RAISE the modulation and coding schemes and explicitly say they
    are held elsewhere. A weak probe asks about something the corpus never
    mentions; a tight one asks about a detail the corpus starts and does not
    finish, because that is what a model is actually tempted to complete.
    """
    out, events, _ = await ask(
        db, agent, user,
        "What modulation and coding scheme does the Ka-band downlink use?",
    )
    return Outcome(
        "S7 refusal survives tools", out.refused,
        f"refused={out.refused} answer={preview(out.answer, 120)}",
    )


async def s9_timing_adds_up(db, agent, user) -> Outcome:
    """The phases must sum to the turn instead of overlapping."""
    out, events, _ = await ask(
        db, agent, user, "How many battery modules are there and how are they arranged?"
    )
    by_type = {e.event_type: (e.duration_ms or 0) for e in events}
    parts = sum(by_type.get(k, 0) for k in ("REWRITE", "RETRIEVE", "GENERATE"))
    parts += sum(e.duration_ms or 0 for e in events if e.event_type in ("TOOL_RESULT", "TOOL_ERROR"))
    total = out.latency_ms or 1
    ratio = parts / total
    # Generous: the turn also carries the history read and the citation join,
    # and RERANK deliberately records no duration.
    ok = 0.5 <= ratio <= 1.15
    return Outcome(
        "S9 phase timings add up", ok,
        f"parts={parts} ms total={total} ms ratio={ratio:.2f}",
    )


async def s10_citation_integrity(db, agent, user) -> Outcome:
    """After any number of searches, markers must be contiguous and resolvable."""
    import re

    out, events, _ = await ask(
        db, agent, user,
        "What are the three downlink paths, and what does each one carry?",
    )
    markers = [c.marker for c in out.citations]
    contiguous = markers == list(range(1, len(markers) + 1))
    used = {int(m) for m in re.findall(r"\[(\d+)\]", out.answer)}
    unresolved = sorted(used - set(markers))
    detail = f"markers={markers} unresolved={unresolved}"

    # With zero citations both halves are vacuously true: `[] == list(range(1,1))`
    # and `set() - set()` is empty. A turn that retrieved nothing and cited
    # nothing used to pass a check named "citation integrity".
    if not markers:
        return unmeasured(
            "S10 citation integrity",
            f"{detail} -- the turn cited nothing, so contiguity was not tested",
        )
    ok = contiguous and not unresolved
    return Outcome("S10 citation integrity", ok, detail)


# --------------------------------------------------------------------------
# S13-S16 -- the DeepSeek swap. See `new features/09-deepseek-agentic.md`.
#
# The swap inverted this suite's founding assumption. S3 and the gap trigger were
# both built on `new features/loop.md` T1: gemma-4-31b-it would not initiate a
# search on its own judgement, 0 tool calls in every prompt configuration tried.
# `deepseek/deepseek-v4-flash-0731` self-initiates 6/6 on the same probe.
#
# That does not retire the trigger, it makes it CONDITIONAL, and conditional
# machinery is what rots. These four pin both halves: the trigger still works
# where it is still needed (S13), it no longer fires where it would now be waste
# (S14), the model really is choosing to search (S15), and the two redundant
# mechanisms holding S15 up have not BOTH been removed (S16).
# --------------------------------------------------------------------------

# The generation model the gap trigger was designed against. S13 owns this the
# way S1 owns `tools_enabled` -- the scenario must not depend on the fixture
# happening to be configured hostilely.
GAP_TRIGGER_MODEL = "google/gemma-4-31b-it"


def _searches(events) -> list[dict]:
    return [p for p in payloads_of(events, "TOOL_CALL") if p.get("tool") == "search_corpus"]


def _gap_forced(events) -> list[dict]:
    """Searches the GAP TRIGGER forced, not ones the model chose.

    `agent_loop` stamps `trigger="gap_detected"` into the invocation args, and
    `ask.run_turn` records `args` verbatim, so the durable trace distinguishes the
    two. A model-chosen call has no `trigger` key at all.
    """
    return [p for p in _searches(events) if (p.get("args") or {}).get("trigger") == "gap_detected"]


async def s13_gap_trigger_still_fires(db, agent, user) -> Outcome:
    """The gap trigger must still work on a model that does not self-initiate.

    **This is the scenario that stops the trigger becoming dead code.** With the
    default model searching on its own, every other scenario here would stay green
    if the whole gap branch were deleted -- so its only remaining proof is a model
    that behaves the way gemma does, which `agents.generation_model` can still
    select and which CLAUDE.md documents an operator being able to type in.

    Owns TWO preconditions, because either one alone makes the test vacuous:
    `generation_model`, so the model genuinely will not search unprompted; and
    `retrieve_k=1`, per `loop.md` section 5 -- on a seven-chunk corpus a wider
    retrieval answers both halves and there is no gap to detect. Both restored.

    Read as `loop.md` T2: the assertion is "did a gap-triggered search occur",
    never "did the turn succeed". A turn that answers half and stops succeeds.
    """
    prior_model = agent.generation_model
    prior_k, prior_n = agent.retrieve_k, agent.rerank_top_n
    agent.generation_model = GAP_TRIGGER_MODEL
    agent.retrieve_k, agent.rerank_top_n = 1, 1
    await db.commit()
    try:
        out, events, _ = await ask(
            db, agent, user,
            "How many battery modules does the platform carry, and separately, "
            "how much onboard storage do the science instruments have?",
        )
    finally:
        agent.generation_model = prior_model
        agent.retrieve_k, agent.rerank_top_n = prior_k, prior_n
        await db.commit()
    forced = _gap_forced(events)
    total = _searches(events)
    # Either the trigger fired, or the model searched by itself -- on a model
    # measured at 0/N that second branch would be news, so it is reported rather
    # than failed. What must NOT happen is neither.
    ok = bool(total)
    return Outcome(
        "S13 gap trigger still fires", ok,
        f"model={GAP_TRIGGER_MODEL} searches={len(total)} gap_forced={len(forced)}"
        + ("" if forced else "  (no gap-forced search: the model self-initiated)"),
    )


async def s14_no_redundant_gap_search(db, agent, user) -> Outcome:
    """A correct refusal that already searched must not be made to search again.

    The saving the `not corpus_searched` gate buys. Before it, the ordinary shape
    of a correct refusal on a self-initiating model was: search, find nothing, say
    so -- and then have `detect_gap` force the SAME search a second time, on every
    single refusal.

    The probe is S7's, chosen the way CLAUDE.md says a refusal probe should be:
    both fixtures RAISE the modulation scheme and say it is held elsewhere, so the
    model is tempted to complete it and a search genuinely returns nothing useful.

    Necessity: asserting only "no gap-forced search" would pass on a turn that
    never searched at all, which is the failure this scenario is supposed to
    detect the absence of. So it requires a real search first.
    """
    out, events, _ = await ask(
        db, agent, user,
        "What modulation and coding scheme does the Ka-band downlink use?",
    )
    total = _searches(events)
    forced = _gap_forced(events)
    ok = bool(total) and not forced
    return Outcome(
        "S14 no redundant gap search", ok,
        f"searches={len(total)} gap_forced={len(forced)} refused={out.refused} "
        f"answer={preview(out.answer, 90)}",
    )


async def s15_model_initiates_search(db, agent, user) -> Outcome:
    """The premise of the swap: the model chooses to search, unprompted.

    `loop.md` T1 says to assume it will not, and that assumption is why the gap
    trigger exists. This asserts the measured inversion on the shipped
    configuration -- a search whose `trigger` is absent, i.e. one the MODEL
    decided on rather than one the loop forced.

    Two trials, pass on either. Not superstition: measured 2026-08-16 the shipped
    configuration self-initiates 6/6, while the configuration this suite is meant
    to catch -- reasoning off AND the guidance paragraph gone -- scores 2/6. Two
    trials separate those two populations (about 89 percent of broken runs go red)
    without spending four full turns on a property S16 also pins deterministically.

    Owns `retrieve_k`, like S3 and S13: the question must be unanswerable from the
    first retrieval or choosing not to search is the correct behaviour.
    """
    prior_k, prior_n = agent.retrieve_k, agent.rerank_top_n
    agent.retrieve_k, agent.rerank_top_n = 1, 1
    await db.commit()
    try:
        seen = []
        for _ in range(2):
            _out, events, _convo = await ask(
                db, agent, user,
                "How many battery modules does the platform carry, and separately, "
                "how much onboard storage do the science instruments have?",
            )
            chosen = [p for p in _searches(events) if not (p.get("args") or {}).get("trigger")]
            seen.append(len(chosen))
            if chosen:
                break
    finally:
        agent.retrieve_k, agent.rerank_top_n = prior_k, prior_n
        await db.commit()
    ok = any(seen)
    return Outcome(
        "S15 model initiates search", ok,
        f"model={agent.generation_model or settings.generation_model} "
        f"self_initiated_per_trial={seen}",
    )


async def s16_tool_use_has_a_belt_and_braces(db, agent, user) -> Outcome:
    """At least one of the two mechanisms that produce tool use must survive.

    No model call, no database read -- this is a structural assertion, and it is
    deliberately the only one of the four that cannot be flaky.

    Measured 2026-08-16, 6 trials per cell, "did it search unprompted":

                             guidance paragraph      no guidance paragraph
        reasoning on              6/6                      6/6
        reasoning off             6/6                      2/6

    The two are REDUNDANT WITH EACH OTHER. Either alone holds the behaviour, which
    is what makes `generation_reasoning=False` affordable -- and it is also what
    makes this dangerous, because it means the guidance paragraph now looks like
    dead weight from a superseded model. Deleting it costs nothing measurable
    until reasoning is also off, which it already is. Nothing raises; tool use
    drops to a third; S15 goes intermittent and gets marked flaky.

    So this asserts the DISJUNCTION rather than either half, because either half
    alone is a legitimate configuration and only losing both is the bug. That is
    `loop.md` T2 aimed at a config invariant: the failure has no error to test for,
    so the test has to name the outcome.
    """
    from app.rag.agent_loop import TOOL_GUIDANCE

    # The sequencing sentence, not the whole paragraph -- reworded prose should
    # not fail this, only losing the instruction should.
    sequencing = "after searching" in TOOL_GUIDANCE.lower()
    reasoning_on = bool(settings.generation_reasoning)
    ok = sequencing or reasoning_on
    return Outcome(
        "S16 tool use has a belt and braces", ok,
        f"guidance_sequencing={sequencing} generation_reasoning={reasoning_on}"
        + ("" if ok else "  BOTH REMOVED -- tool use measured 2/6 in this state"),
    )


# --------------------------------------------------------------------------
# S18-S19 -- the rewriter on every turn. See `new features/10-*.md`.
#
# The rewriter is a plain code path, not a tool and not a trigger (`loop.md`
# section 6 item 1), so there is no trigger to design -- which moves the whole
# design risk onto proving the thing is NEEDED. These two are the halves of that:
# S18 says the rewrite reaches a chunk the raw question cannot, S19 says it does
# not maul the questions that were already fine.
# --------------------------------------------------------------------------

S18_RAW = "how much dos the uhf lnk drw whl actv"
# comms-subsystem.md chunk 1, the power-draw paragraph. Unique in the corpus:
# no other chunk contains the string, checked against all seven.
S18_TARGET = "0.5 kW"

# How far BEHIND the winning chunk the target has to sit before "the raw form
# cannot reach it" is a finding rather than a coin flip.
#
# **This constant is the fix, and the number that forced it was 0.0002.** The
# probe this scenario shipped with -- "how meny km is the prox lnk gd for" ->
# "200 km" -- was adopted on a bare `missed_raw=True`, which is a boolean over a
# continuous quantity. An independent reading put the winner at 0.5422 against
# the target chunk's 0.5420; this machine measures the same probe at +0.0066
# with a DIFFERENT chunk winning. Two readings, two winners, and the boolean
# said "necessary" both times.
#
# 0.02 comes from the measured spread, not from taste. The same string embedded
# three times in one process moves 0.000000 -- the endpoint is bit-deterministic
# -- so the noise is not in the model. It is in the batch: the same probe scored
# against the LIVE namespace and against locally re-embedded chunks differs by
# up to 0.005, because those vectors were written by different `embed_documents`
# calls. 0.02 is four times that drift, and the shipped probe clears 0.02 by
# 1.7x. It is deliberately not tuned down to admit a near-miss: `loop.md`
# section 5 is that a scenario which can pass without exercising the feature is
# worse than no scenario, and a red row that means "re-run me" teaches its
# reader to ignore red.
S18_MIN_MISS_MARGIN = 0.02


def target_margin(retrieval, target: str) -> tuple[float | None, int | None]:
    """(margin, rank) of the best-scoring retrieved chunk containing `target`.

    The margin is `top score - best target score` in Pinecone cosine: 0.0 when
    the target IS the winner, and positive by exactly how much the query missed
    by. Returns `(None, None)` when nothing retrieved contains the target at
    all, which is a louder failure than a large margin and must not be read as
    one -- hence None rather than a sentinel float that would compare.

    Sorted here rather than trusted: Pinecone returns descending today, and an
    assertion resting on that staying true forever is one more thing a provider
    change can invert without raising.
    """
    scored = [(doc, float(score)) for doc, score in retrieval.scored]
    if not scored:
        return None, None
    scored.sort(key=lambda pair: -pair[1])
    hits = [i for i, (doc, _) in enumerate(scored) if target in doc.page_content]
    if not hits:
        return None, None
    return scored[0][1] - scored[hits[0]][1], hits[0] + 1


async def s18_rewrite_is_necessary(db, agent, user) -> Outcome:
    """A question the RAW form cannot retrieve, and the rewritten form can.

    Both halves matter. `new features/loop.md` section 5 records S3 passing twice
    while proving nothing, because the scenario never checked that the un-helped
    path FAILED. A rewriter that has completely stopped working still produces: a
    successful turn, a WARNING nobody reads, no exception, and marginally worse
    retrieval. Every error-shaped check passes over it -- "did it return", "did
    the turn answer", "is there a REWRITE row" (a dead rewriter that returns the
    input unchanged still writes one), and even "after != before" (a rewriter
    that MANGLES the question also passes that).

    **The assertion is a MARGIN and not a boolean, and that is the 2026-08-16
    fix.** The version before this one asserted `target not in the k=1 result`,
    which is section 5's trap one level down: the boolean was True on a gap of
    two ten-thousandths of a cosine, so the single assertion carrying the whole
    rewriter change was a coin flip that had come up heads. `target_margin`
    above reports how far behind the target actually is, `S18_MIN_MISS_MARGIN`
    is what that has to clear, and the number is PRINTED in the detail so a
    near-miss is visible in the run that produced it rather than in the run
    after it flips.

    **The probe was chosen by measurement over 81 candidates across both fixture
    files** -- misspellings, shorthand, and vocabulary the corpus never writes --
    each re-split with the shipped splitter at the fixture agent's own 250/40,
    embedded through `retriever.get_embeddings()` and ranked by cosine.
    Fourteen missed by >= 0.02; twelve of those went through the REAL
    `contextualize_question`, five to seven trials each; four survived both
    halves. Measured live against the ingested namespace, 2026-08-16 --
    raw top-1 chunk, best chunk holding the target, and the gap:

        how much dos the uhf lnk drw whl actv   -> "0.5 kW"      <- SHIPPED
            comms#2 0.6218  comms#1 0.5881  +0.0337  rank 3 of 7
            rewrite 7/7 byte-identical, target at rank 1 (0.7222, next 0.7045)
        wat dos the uhf lnk drw whn its actv    -> "0.5 kW"
            comms#2 0.6581  comms#1 0.6252  +0.0329  rank 3 of 7
            rewrite 7/7 to rank 1, in two spellings
        how much does the uhf lnk drw whn actv  -> "0.5 kW"
            comms#2 0.6309  comms#1 0.6073  +0.0236  rank 3 of 7
            rewrite 7/7 byte-identical to rank 1

    **Both probes this docstring used to carry were wrong, and the second was
    wrong in the dangerous direction.** `"how meny km is the prox lnk gd for"`
    -> `"200 km"` measures +0.0066, a coin flip rather than the stable miss it
    was adopted as. `"wats the bcklg drain hrzn in days"` -> `"eighteen days"`
    was recorded here as "5/5 and 5/5": it measures +0.0077 on this machine,
    while the reviewer measured the RAW form retrieving that target at rank 1
    (comms#1, 0.6023) -- so swapping to it would have made S18 permanently red.
    Two readings, opposite winners, eight thousandths apart. **A pass count
    cannot tell a stable miss from a coin flip**, which is why every line above
    carries a margin instead of a tally.

    **A large raw miss and a working rewrite pull AGAINST each other**, which is
    the transferable half if this corpus is ever edited. The biggest miss
    measured was +0.0785 -- `"how much for cmd data hndlng n avncs"` ->
    `"2.4 kW"`, target at rank 6 of 7 -- and its rewrite fails 0/5: shorthand
    mangled enough to destroy the embedding is usually mangled enough that the
    rewriter guesses too. "avncs" came back as "advances", and "drw" in
    `"wat does the ka bnd xmtr n its pntng mech drw"` (+0.0573) came back as
    "drawing show" in 4 of 5 trials. Conversely, shorthand the rewriter repairs
    cleanly is usually shorthand the embedding already survived. The band that
    satisfies both is narrow, and "uhf lnk drw whl actv" sits in it because
    every token repairs exactly one way while the mangled form lands on comms#2
    -- the least topically specific chunk in the corpus (Ka-band availability
    plus "what this briefing does not cover"), and the attractor for most of the
    measured misses.

    **No reranking in the two measurement retrievals, deliberately.** The turn
    is still starved to `retrieve_k=1, rerank_top_n=1` so the model's own
    context is the un-helped one, and at k=1 reranking a single document cannot
    reorder anything -- the vector ranking IS the outcome. Measuring it directly
    reports the distance rather than the boolean a reranked k=1 call could only
    have returned, and costs two fewer Cohere calls in a suite CLAUDE.md records
    tripping a rate limit.

    `k` is the corpus chunk count read from the database rather than a literal
    7. A window smaller than the corpus can leave the target outside it, and
    `target_margin` would then report "not retrieved at all" for a reason that
    has nothing to do with the rewriter.
    """
    prior_k, prior_n = agent.retrieve_k, agent.rerank_top_n
    # Owned, not inherited: on a seven-chunk corpus at the fixture's own k=3 the
    # candidate set is nearly half the corpus and the raw form cannot fail. Same
    # reasoning as S3 and S13, restored in the `finally` the same way.
    agent.retrieve_k, agent.rerank_top_n = 1, 1
    await db.commit()
    try:
        # Through the real turn, so the string under test is the one the PIPELINE
        # produced and recorded -- not one this scenario asked for separately.
        _out, events, _ = await ask(db, agent, user, S18_RAW)
        payload = (payloads_of(events, "REWRITE") or [{}])[0]
        after = payload.get("after")
        corpus = await db.scalar(
            select(func.count(Chunk.id))
            .join(Document, Chunk.document_id == Document.id)
            .where(Document.agent_id == agent.id)
        ) or 0
        raw = await aretrieve(agent, S18_RAW, rerank=False, k=corpus)
        new = await aretrieve(agent, after or S18_RAW, rerank=False, k=corpus)
    finally:
        agent.retrieve_k, agent.rerank_top_n = prior_k, prior_n
        await db.commit()

    raw_margin, raw_rank = target_margin(raw, S18_TARGET)
    _new_margin, new_rank = target_margin(new, S18_TARGET)
    ok = (
        bool(after)
        and raw_margin is not None
        and raw_margin >= S18_MIN_MISS_MARGIN
        and new_rank == 1
    )
    shown = "--" if raw_margin is None else format(raw_margin, "+.4f")
    return Outcome(
        "S18 the rewrite is necessary", ok,
        f"raw_margin={shown} (needs >= +{S18_MIN_MISS_MARGIN:.4f}) "
        f"raw_rank={raw_rank}/{corpus} rewritten_rank={new_rank} "
        f"trigger={payload.get('trigger')} after={preview(after or '', 70)}",
    )


async def s19_no_over_firing(db, agent, user) -> Outcome:
    """A clean question must survive the rewriter byte-identical.

    **The over-firing guard, and it is the half that protects the user's
    meaning.** The asymmetry is stated above `CONTEXTUALIZE_SYSTEM_PROMPT`: a
    false positive rewrites a question that was already fine and can change what
    was asked, while a false negative leaves a typo and costs slightly worse
    retrieval. Every number in EVAL.md was also measured on unrewritten text.

    One scenario catches four different failures at once, which is why it asserts
    three keys rather than one: the rewriter dead (no row at all), the rewriter
    mangling clean questions (`changed=True`), the rewriter silently failing
    (`failed=True` -- `contextualize_question` swallows every exception, so this
    key is the only place that shows), and the trigger label being wrong
    (`first_turn`, which was hardcoded to `conversation_history` until
    2026-08-16 and mislabelled every first turn).

    The question is the one `scripts/rewrite_check.py` case 8 measured at 5/5
    byte-identical, so a red here is a change in behaviour rather than a sample.
    """
    question = "How much power do the solar arrays generate?"
    _out, events, _ = await ask(db, agent, user, question)
    rows = payloads_of(events, "REWRITE")
    payload = rows[0] if rows else {}
    ok = (
        len(rows) == 1
        and payload.get("trigger") == "first_turn"
        and payload.get("failed") is False
        and payload.get("changed") is False
    )
    return Outcome(
        "S19 clean question is left alone", ok,
        f"rows={len(rows)} trigger={payload.get('trigger')} "
        f"failed={payload.get('failed')} changed={payload.get('changed')} "
        f"after={preview(str(payload.get('after')), 80)}",
    )


# --------------------------------------------------------------------------
# S20-S27 -- the orchestrator, @mentions and the self-check.
# See `new features/11-orchestrator-and-self-check.md`, section 1 Q5.
#
# Layer 1 for this feature is `scripts/route_specialist_check.py`: mention
# parsing, the alias table, the self-check trigger and the critic's carve-out,
# with no database and no model. Read a red row there before a red row here --
# the S16-before-S15 ordering, for the same reason.
#
# **Every one of these owns the columns it needs.** `agents.specialists` being
# NULL is the entire off switch, so a scenario that inherited the fixture's
# configuration would be asserting against whatever the previous scenario left
# behind. `loop.md` section 5 records S3 passing twice while proving nothing for
# exactly that reason, and S21, S23, S25, S26 and S27 would each be vacuous in a
# different way: S21 with no roster records no ROUTE row at all, and S25 with
# `self_check_enabled` false finds zero SELF_CHECK rows because the feature never
# ran, which is the answer it is looking for arrived at by the wrong route.
# --------------------------------------------------------------------------

from app.db.specialists import DEFAULT_ROSTER  # noqa: E402
from app.rag import events as rag_events  # noqa: E402
from app.rag.selfcheck import (  # noqa: E402
    SIGNAL_PHANTOM,
    VERDICT_UNGROUNDED,
    markers_in,
)

# What S21 measured, read by S22. See S22's docstring: a mention that agrees
# with the router proves nothing, so S22 needs to know what the router actually
# chose rather than assuming.
_ROUTED: dict[str, str] = {}

# The question S21 uses for the explanatory arm, reused verbatim by S22 so that
# the only variable between them is the mention.
S21_EXPLAIN = "Explain simply how the Ka-band downlink budget works."
S21_QUIZ = "Write me five practice questions on the power allocation by subsystem."


def _orchestration(agent, *, specialists, self_check):
    """Set the two orchestrator columns, returning what was there before.

    `getattr` with a default because these columns arrive with the feature's own
    migration: before it lands the attributes do not exist, and a scenario that
    raised `AttributeError` while reading them would report the feature as broken
    rather than as absent.
    """
    prior = (
        getattr(agent, "specialists", None),
        getattr(agent, "self_check_enabled", False),
    )
    agent.specialists = specialists
    agent.self_check_enabled = self_check
    return prior


def _route_payload(events) -> dict:
    rows = payloads_of(events, "ROUTE")
    return rows[0] if rows else {}


async def s20_orchestration_off(db, agent, user) -> Outcome:
    """Orchestration OFF must produce exactly the pre-existing trace shape.

    The regression test, and the S1 of this feature. Everything else here checks
    that routing and the self-check work; this checks they did not change the
    turn of an agent that has neither.

    `loop.md` S4 is the standard being held to: not "similar", identical. NULL
    and false are the classic path by construction -- no backfill was needed,
    unlike `bc307f5fc31f` -- so a stray ROUTE row here means the branch is keyed
    on something other than the column.
    """
    prior = _orchestration(agent, specialists=None, self_check=False)
    await db.commit()
    try:
        out, events, _ = await ask(
            db, agent, user, "How much power do the solar arrays generate?"
        )
        seen = set(event_types(events))
        stray = seen - PRE_ORCHESTRATOR_EVENTS
        orchestrated = sorted(seen & ORCHESTRATOR_EVENTS)
        ok = bool(out.answer) and not stray
        return Outcome(
            "S20 orchestration off is unchanged", ok,
            f"events={sorted(seen)}"
            + (f" STRAY={sorted(stray)}" if stray else "")
            + (f" ORCHESTRATOR_ROWS={orchestrated}" if orchestrated else ""),
        )
    finally:
        agent.specialists, agent.self_check_enabled = prior
        await db.commit()


async def s21_router_discriminates(db, agent, user) -> Outcome:
    """Two questions on one corpus must route to two DIFFERENT specialists.

    **One question cannot test a router.** A router hardwired to return
    `feynman-explainer`, a router whose structured output silently failed and
    fell back, and a router that read the question all produce the same passing
    row on a single "explain this" probe. The finding is the DIFFERENCE, so the
    scenario has to ask twice and compare.

    The pair is chosen so the topic is constant and only the ASK varies -- "the
    Ka-band budget" against "the power allocation" would confound the router's
    signal with the corpus. What separates them is "explain simply" versus "write
    me five practice questions", which is the distinction `route.py`'s prompt
    names in its own first bullet.

    Separate conversations, deliberately: the router sees history, so asking both
    in one thread would let the first answer influence the second and the
    difference would no longer be a property of the questions.

    Also asserts `trigger == "router"` on both. Without it a fallback -- which
    `route_specialist` returns on a failed structured call, by design, so that a
    dead router degrades to a working answer -- would be indistinguishable from a
    decision, and this scenario would be measuring the fallback path.
    """
    prior = _orchestration(agent, specialists=list(DEFAULT_ROSTER), self_check=False)
    await db.commit()
    try:
        _out_a, events_a, _ = await ask(db, agent, user, S21_EXPLAIN)
        _out_b, events_b, _ = await ask(db, agent, user, S21_QUIZ)
    finally:
        agent.specialists, agent.self_check_enabled = prior
        await db.commit()

    a, b = _route_payload(events_a), _route_payload(events_b)
    chosen_a, chosen_b = a.get("specialist"), b.get("specialist")
    if chosen_a:
        _ROUTED["explain"] = chosen_a
    if chosen_b:
        _ROUTED["quiz"] = chosen_b

    ok = (
        bool(chosen_a)
        and bool(chosen_b)
        and chosen_a != chosen_b
        and a.get("trigger") == "router"
        and b.get("trigger") == "router"
    )
    return Outcome(
        "S21 the router discriminates", ok,
        f"explain->{chosen_a} ({a.get('trigger')}) quiz->{chosen_b} "
        f"({b.get('trigger')}) why={preview(str(a.get('why')), 60)} / "
        f"{preview(str(b.get('why')), 60)}",
    )


async def s22_mention_overrides_router(db, agent, user) -> Outcome:
    """A typed mention must beat the router's own choice.

    **Without S21 having run first this scenario proves nothing**, and that is
    not a caution, it is the reason the pair exists. Assert `@quiz ... -> quiz
    writer` on its own and a router that would have picked the Quiz Writer anyway
    passes it -- the mention could be agreeing with the router by luck, and a
    parser that had been deleted entirely would look identical. So the specialist
    mentioned here is chosen to be one S21 MEASURED the router not choosing for
    this exact question, and the scenario says so in its detail line either way.

    Run with `--only s22` and the fallback below keeps it runnable, but the
    detail prints `s21_not_run` and the finding is weaker by exactly that much.
    """
    prior = _orchestration(agent, specialists=list(DEFAULT_ROSTER), self_check=False)
    await db.commit()
    try:
        router_pick = _ROUTED.get("explain")
        # Anything the router demonstrably did not pick for S21_EXPLAIN.
        alternatives = [s for s in DEFAULT_ROSTER if s != (router_pick or "feynman-explainer")]
        target = "quiz-generator" if "quiz-generator" in alternatives else alternatives[0]
        # An ALIAS rather than the slug, because that is what a user types --
        # and because it exercises the alias table on the live path rather than
        # only in `route_specialist_check.py`.
        token = "@quiz" if target == "quiz-generator" else f"@{target}"
        out, events, _ = await ask(db, agent, user, f"{token} {S21_EXPLAIN}")
    finally:
        agent.specialists, agent.self_check_enabled = prior
        await db.commit()

    payload = _route_payload(events)
    ok = (
        payload.get("trigger") == "mention"
        and payload.get("specialist") == target
        # The mention must be STRIPPED before the embedder. `@quiz` in vector
        # space is noise, and the rewriter is documented to mangle terms it does
        # not recognise -- so a mention that survives into the search query is a
        # silent retrieval regression, not a cosmetic one.
        and token not in (out.rewritten_question or "")
    )
    return Outcome(
        "S22 a mention overrides the router", ok,
        f"router_picked={router_pick or 's21_not_run'} mentioned={token} "
        f"got={payload.get('specialist')} trigger={payload.get('trigger')} "
        f"search_query={preview(out.rewritten_question or '-', 60)}",
    )


async def s23_two_mentions_two_sections(db, agent, user) -> Outcome:
    """Two mentions produce two sections over ONE shared ledger.

    Two DELEGATE events and two `## ` headings are the visible half. **The real
    assertion is the ledger**, because concatenating two specialists' text is
    only safe if `[2]` means the same chunk in both sections -- and a
    two-ledger implementation produces an answer that looks exactly like a
    one-ledger implementation while attributing half its claims to the wrong
    source. There is no error to test for; this is `loop.md` T2 aimed at a
    citation.

    It is asserted through `query_chunks` rather than by comparing the two
    sections' text, because `query_chunks.rank` IS the marker (see
    `ask.CitationOut`: "there is no `marker` column"). Two ledgers written to one
    turn therefore collide on rank, and a turn that wrote only the first
    ledger's chunks leaves the second section's markers unresolvable. Those are
    the two ways the feature can be built wrongly, and this catches both.

    A marker shared BETWEEN the sections is asserted when there is one, and its
    absence is reported rather than failed: whether two specialists happen to
    cite the same passage is a property of the corpus and the question, while
    rank-uniqueness is a property of the design and is always available.
    """
    import re

    prior = _orchestration(agent, specialists=list(DEFAULT_ROSTER), self_check=False)
    await db.commit()
    try:
        out, events, _ = await ask(
            db, agent, user,
            "@feynman @quiz the onboard storage margin for science data",
        )
        rows = (await db.scalars(
            select(QueryChunk).where(QueryChunk.query_id == out.query_id)
        )).all()
    finally:
        agent.specialists, agent.self_check_enabled = prior
        await db.commit()

    delegates = payloads_of(events, "DELEGATE")
    bodies = re.split(r"^##\s+", out.answer or "", flags=re.M)[1:]

    by_rank: dict[int, uuid.UUID] = {}
    duplicate_ranks: list[int] = []
    for row in rows:
        if row.rank in by_rank:
            duplicate_ranks.append(row.rank)
        by_rank[row.rank] = row.chunk_id

    used = markers_in(out.answer or "")
    unresolved = sorted(used - set(by_rank))

    shared = set.intersection(*(markers_in(b) for b in bodies)) if len(bodies) >= 2 else set()
    shared_chunks = {m: by_rank.get(m) for m in sorted(shared)}
    # `all([])` is True, so this half is vacuous whenever the two specialists
    # happened not to cite the same passage. That is deliberate (see the
    # docstring) and it is NOT the whole scenario -- the four assertions above it
    # are real and always available, so flagging the run `[warn]` would hide four
    # passing checks behind one that could not run. It is spelled out in the
    # detail line instead, which is the smallest honest thing.
    shared_ok = all(cid is not None for cid in shared_chunks.values())
    shared_note = (
        f"shared_markers={sorted(shared)}"
        if shared
        else "shared_markers=none <- that half NOT MEASURED"
    )

    ok = (
        len(delegates) == 2
        and len(bodies) >= 2
        and not duplicate_ranks
        and not unresolved
        and shared_ok
    )
    return Outcome(
        "S23 two mentions, two sections, one ledger", ok,
        f"delegates={len(delegates)} sections={len(bodies)} "
        f"query_chunks={len(rows)} duplicate_ranks={duplicate_ranks} "
        f"unresolved_markers={unresolved} {shared_note}",
    )


async def s24_phantom_marker_is_caught(db, agent, user) -> Outcome:
    """A citation outside the ledger must fire, and no verdict may excuse it.

    **This one does NOT go through a model, and that is the design rather than a
    shortcut.** The end-to-end version needs the generation model to invent a
    citation on demand, which is not reproducible: it would pass on the runs
    where the model happened to hallucinate and go red on the runs where it
    behaved, and CLAUDE.md already records what an intermittent red does -- the
    fifth refusal-marker miss "surfaced as an INTERMITTENT red ... a flaky
    refusal scenario reads as model variance and gets re-run". A scenario nobody
    believes is worse than no scenario.

    So the draft is constructed and the detector is called directly. What earns
    this a place in layer 2 rather than beside the identical case in
    `route_specialist_check.py` is the SECOND half, which layer 1 cannot reach:
    with the critic stubbed to say `grounded=True`, the result must STILL be
    ungrounded. The signal overrides the verdict in exactly one direction -- no
    model opinion can make a citation to a passage that does not exist true --
    and that asymmetry is a line of code with nothing else asserting it.
    """
    from app.rag import selfcheck

    draft = (
        "The Ka-band high-rate downlink is budgeted at 220 Mbps [1] and the "
        "onboard storage holds 34 TB [2]. The relay constellation is visible "
        "for 74 percent of each orbit [99]."
    )
    signal = selfcheck.self_check_signal(draft, ledger_size=3, expects_citations=True)

    class _AlwaysGrounded:
        async def ainvoke(self, _payload):
            return selfcheck.GroundingVerdict(
                grounded=True, unsupported=[], suggested_query=None
            )

    prior_critic = selfcheck.get_critic
    selfcheck.get_critic = lambda: _AlwaysGrounded()
    try:
        result = await selfcheck.run_grounding_critic(
            answer=draft, context="[1] ... [2] ... [3] ...", signal=signal
        )
    finally:
        selfcheck.get_critic = prior_critic

    ok = (
        signal is not None
        and signal.name == SIGNAL_PHANTOM
        and signal.phantom_markers == (99,)
        and result.verdict == VERDICT_UNGROUNDED
        and result.phantom_markers == (99,)
    )
    return Outcome(
        "S24 a phantom marker is caught, and no verdict excuses it", ok,
        f"signal={getattr(signal, 'name', None)} "
        f"phantom={getattr(signal, 'phantom_markers', None)} "
        f"verdict_with_a_grounded_critic={result.verdict}",
    )


async def s25_no_self_check_on_a_socratic_turn(db, agent, user) -> Outcome:
    """The self-check must not fire on a persona designed to cite nothing.

    **This is the S7 of this feature** -- the one that checks the addition did
    not eat the existing behaviour. A Socratic turn is a question put back to the
    learner; it asserts nothing, so it cites nothing, and it is CORRECT when it
    does. Firing a critic on it is the same defect as `refusal_pass = 0/2`: a
    measurement penalising the behaviour the persona exists to produce, which
    CLAUDE.md records costing three separate investigations before anyone read
    the answers.

    **It owns `self_check_enabled = True`, and that is the whole scenario.**
    Zero SELF_CHECK rows on an agent with the feature switched off is not a
    finding, it is the feature being absent -- the right answer reached by the
    wrong route, and precisely the trap `loop.md` section 5 describes. It also
    asserts the turn genuinely routed to `socratic-tutor`, because a failed route
    would produce zero SELF_CHECK rows too, for a third wrong reason.
    """
    prior = _orchestration(agent, specialists=list(DEFAULT_ROSTER), self_check=True)
    await db.commit()
    try:
        out, events, _ = await ask(
            db, agent, user,
            "@socratic I think the Ka-band link clears all the science data "
            "generated each day, because it is the high-rate path. Is that right?",
        )
    finally:
        agent.specialists, agent.self_check_enabled = prior
        await db.commit()

    payload = _route_payload(events)
    checks = payloads_of(events, "SELF_CHECK")
    routed = payload.get("specialist") == "socratic-tutor"
    ok = routed and not checks and bool(out.answer)
    return Outcome(
        "S25 no self-check on a Socratic turn", ok,
        f"self_check_enabled=True routed_to={payload.get('specialist')} "
        f"trigger={payload.get('trigger')} self_check_rows={len(checks)} "
        f"markers={sorted(markers_in(out.answer or '')) or 'none'} "
        f"answer={preview(out.answer, 90)}",
    )


async def s26_critic_exempts_pedagogy(db, agent, user) -> Outcome:
    """A teaching answer must not have its draft discarded for teaching.

    **PRD open item 20, and the one scenario here that can fail by making the
    product quietly worse rather than visibly broken.** EVAL run 3 scored an
    answer 0.571 on faithfulness whose sentences 1-4 were four correct figures
    straight from the context; the deductions were a labelled analogy and
    "restate this in your own words". The scorecard then named faithfulness as
    the weakest metric and advised reducing persona verbosity -- that is,
    deleting the pedagogy. A groundedness critic with no carve-out rebuilds that
    instrument inside the live turn, where it does not merely recommend deleting
    the teaching: it discards the draft the user watched stream and tells the
    model to write a duller one. Nothing raises.

    The layer-1 half is `route_specialist_check.py` cases 25-31 -- the free
    pre-check does not fire on that answer, and the critic prompt still carries
    its exemptions. This is the end-to-end half, and it goes through the
    streaming seam because `answer_reset` is a wire frame with no trace row (see
    `ask()`).

    Two ways the draft can be thrown away and both are asserted: the frame the
    browser would render, and `SELF_CHECK.acted`, which section 7.1 exists to
    separate "we checked and it was fine" from "we checked and acted". A turn
    where the pre-check never fired passes; a turn where it fired and the critic
    said grounded passes; a turn where the draft was discarded does not.
    """
    prior = _orchestration(agent, specialists=list(DEFAULT_ROSTER), self_check=True)
    await db.commit()
    frames, emit = frame_collector()
    try:
        out, events, _ = await ask(
            db, agent, user,
            "@feynman Explain simply how the battery store covers the power "
            "deficit, with an analogy.",
            emit=emit,
        )
    finally:
        agent.specialists, agent.self_check_enabled = prior
        await db.commit()

    resets = [
        p for name, p in frames
        if name == rag_events.ANSWER_RESET and p.get("reason") == "self_check"
    ]
    checks = payloads_of(events, "SELF_CHECK")
    acted = [c for c in checks if c.get("acted")]
    ungrounded = [c for c in checks if c.get("verdict") == VERDICT_UNGROUNDED]

    # NOT MEASURED, not passing, when the critic never ran. The assertion below
    # is `not resets and not acted and not ungrounded` -- three negatives, every
    # one of them trivially true over an EMPTY list. And empty is the expected
    # shape here: the pre-check returns None for any draft carrying a citation
    # (`selfcheck.py:124-125`, `if used: return None`), and a cited answer is
    # exactly what an @feynman turn produces. So the green this printed was the
    # green of a check that could not fail -- `Outcome.unmeasured`'s own comment,
    # ten lines up from its definition, is about this.
    if not checks:
        return unmeasured(
            "S26 the critic exempts pedagogy",
            "the pre-check never fired (the draft carried a citation), so the "
            "carve-out was not exercised -- see route_specialist_check case 41, "
            "which calls the critic directly",
        )

    ok = bool(out.answer) and not resets and not acted and not ungrounded
    return Outcome(
        "S26 the critic exempts pedagogy", ok,
        f"self_check_rows={len(checks)} "
        f"verdicts={[c.get('verdict') for c in checks] or 'none'} "
        f"acted={len(acted)} answer_reset_self_check={len(resets)} "
        f"answer_chars={len(out.answer or '')}",
    )


# The unsupported claim S27 puts in the critic's mouth. Deliberately a string no
# model would ever write: "the answer text is unmodified" is otherwise not
# provable after the fact, because the draft the check saw is recorded nowhere
# this scenario can read. A sentinel makes it exact -- if the turn splices the
# critic's findings into the answer, this appears verbatim.
S27_SENTINEL = "UNSUPPORTED-SENTINEL-Q7"

# A second, weaker guard for a caveat the SYSTEM composed rather than quoted.
# Phrases, not a shape, and it is the half that can miss -- kept because its
# false negatives cost nothing and the sentinel above is what actually bites.
S27_CAVEAT_PHRASES = (
    "may not be fully grounded",
    "could not be verified",
    "not supported by the sources",
    "self-check",
    "this answer was flagged",
)


async def s27_budget_exhaustion_still_answers(db, agent, user) -> Outcome:
    """An ungrounded verdict with no budget left keeps the draft, unmodified.

    Section 4.4's last row, and its asymmetry is the one worth reading twice:
    **editing a model's answer to add a caveat it did not write is the one
    outcome worse than shipping the draft**, because it makes the system's voice
    unreliable in a way no trace event records. So the turn is flagged and the
    text is left exactly as the model wrote it.

    **The failure is forced, not hoped for.** Getting a live model to produce an
    ungrounded draft on demand is the same non-reproducibility S24 refuses, so
    two module globals in `selfcheck` are swapped for the duration: `markers_in`
    makes every draft look as though it cited a passage that does not exist, and
    `get_critic` returns an ungrounded verdict with no model call. Both are
    resolved from `selfcheck`'s own globals at call time, so the patch holds
    whatever import style `pipeline.py` uses -- and both are restored in the
    `finally`, along with `max_tool_steps`.

    `suggested_query` is null on purpose: a search would be the other branch of
    that table, and this scenario is about the branch with nowhere left to go.

    **It owns retrieval as well as the step budget, and the reason is arithmetic
    rather than caution.** `pipeline` computes `steps_left = max_steps -
    loop.steps`, so `max_tool_steps = 1` only exhausts the budget if the loop
    actually spends its step. Starving retrieval to `k=1` on a seven-chunk corpus
    is what makes a search necessary -- the same condition S3, S13 and S15 own
    for the same reason. If the model answers without searching anyway,
    `steps_left` is 1 and the turn takes the REDRAFT branch instead: the three
    assertions below still hold there, because "never edit the model's text"
    binds on both branches, but the budget branch was not the one exercised. The
    detail line prints `acted` and `tool_steps` so that is visible in the run
    rather than inferred from a green row.
    """
    from app.rag import selfcheck

    class _Ungrounded:
        async def ainvoke(self, _payload):
            return selfcheck.GroundingVerdict(
                grounded=False, unsupported=[S27_SENTINEL], suggested_query=None
            )

    prior = _orchestration(agent, specialists=list(DEFAULT_ROSTER), self_check=True)
    prior_steps = agent.max_tool_steps
    prior_k, prior_n = agent.retrieve_k, agent.rerank_top_n
    agent.max_tool_steps = 1
    agent.retrieve_k, agent.rerank_top_n = 1, 1
    await db.commit()

    prior_markers = selfcheck.markers_in
    prior_critic = selfcheck.get_critic
    selfcheck.markers_in = lambda text: {999}
    selfcheck.get_critic = lambda: _Ungrounded()
    try:
        out, events, _ = await ask(
            db, agent, user,
            "Compare the Ka-band availability gaps with the battery replacement "
            "schedule, and explain how the two interact.",
        )
    finally:
        selfcheck.markers_in = prior_markers
        selfcheck.get_critic = prior_critic
        agent.specialists, agent.self_check_enabled = prior
        agent.max_tool_steps = prior_steps
        agent.retrieve_k, agent.rerank_top_n = prior_k, prior_n
        await db.commit()

    checks = payloads_of(events, "SELF_CHECK")
    verdicts = [c.get("verdict") for c in checks]
    acted = [c.get("acted") for c in checks]
    answer = out.answer or ""
    tail = answer[-400:].lower()
    caveats = [p for p in S27_CAVEAT_PHRASES if p in tail]
    ok = (
        bool(answer)
        and VERDICT_UNGROUNDED in verdicts
        and S27_SENTINEL not in answer
        and not caveats
    )
    return Outcome(
        "S27 budget exhaustion still answers, unmodified", ok,
        f"answer_chars={len(answer)} verdicts={verdicts or 'none'} "
        f"acted={acted or 'none'} tool_steps={out.tool_steps} "
        + ("(redraft branch, not the budget branch)" if any(acted) else "")
        + f" critic_text_spliced_in={S27_SENTINEL in answer} "
        f"appended_caveat={caveats or 'no'}",
    )


# --------------------------------------------------------------------------
# S32-S33 -- the OTHER door. See `new features/12-robust-handouts/06-*.md`.
#
# A deck can be asked for in two places and until this change set they shared no
# prompt, no grounding rules and no validation: the panel button ran
# `run_handout_job` with a fifty-line `DECK_PROMPT`, and the chat turn ran
# `run_python` with one bullet of guidance and nothing that opened the file. So
# feature 02 landing alone would have shipped the same defect through the door
# the workshop actually demonstrates.
#
# S32 is the necessity case and S33 is its regression twin, which is the S3/S1
# pairing one layer over.
# --------------------------------------------------------------------------

# One question, asked by both, so the only variable between them is
# `tools_enabled`. It names the file type explicitly: "make me some slides" is a
# request a model can honour in prose, and a scenario about artefacts must not
# be satisfiable without one.
S32_DECK_ASK = (
    "Make me a PowerPoint slide deck as a .pptx file summarising the three "
    "downlink paths, one slide per path."
)


def _pptx_names(payloads: list[dict], key: str) -> list[str]:
    """Filenames ending `.pptx` under `key` in a list of tool payloads."""
    return [
        str(entry.get("filename") or "")
        for payload in payloads
        for entry in (payload.get(key) or [])
        if str(entry.get("filename") or "").lower().endswith(".pptx")
    ]


async def s32_tool_deck_is_rejected(db, agent, user) -> Outcome:
    """A chat-made deck that fails validation: TOOL_ERROR, and NO handout row.

    Feature 06, A7. Two halves, and the second is the one that matters: the tool
    path may not persist a file it has already told the model is broken.
    `PLAN.md` 1.8's all-or-nothing principle -- "a deck missing half its slides
    is worse than a deck that failed" -- applies to the panel as much as to the
    prompt, because a `ready` handout is a download button.

    **It owns its conditions twice over, and the second one is the point.**
    Feature 06 A7 says to starve `retrieve_k` "so a chat-requested deck cannot
    fill `handout_deck_min_slides`", and that alone would be the S3 trap
    verbatim: `PLAN.md` 8.3 MEASURED the model producing five to nine slides at
    the most starved budget in this repo, so a scenario waiting for starvation
    to shrink a deck below three is waiting for something that does not happen.
    An intermittent scenario is re-run until it passes -- CLAUDE.md records
    exactly that, on the fifth refusal-marker miss. So the deck is made
    un-usable by `IMPOSSIBLE_BULLET_LIMIT` instead, for the reasons written
    there, and retrieval is starved as well so the turn is a thin-material turn
    rather than only a moved-threshold one. Both restored in the `finally`.

    Three ways this is honestly unmeasurable, and each reports as such rather
    than as green: the model may not call `run_python` at all, it may call it
    and write something that is not a deck, or the run may crash before saving.
    None of those is evidence about validation.

    **Seen failing before it was believed** (`build.md` section 5). 2026-08-17,
    same ask, same lever, with `handout_validate_artifacts=False` for one run:

        validation on   [ok]   run_python=2 tool_errors=2 error_kinds=[invalid,
                               invalid] rejected=[downlink_paths.pptx x2]
                               persisted_pptx=0
        validation off  [FAIL] run_python=1 tool_errors=0 error_kinds=none
                               wrote_pptx=[downlink_paths.pptx] persisted_pptx=1

    Read the left column twice. With validation on the model called
    `run_python` a SECOND time after being told what was wrong -- the whole
    point of returning a `ToolMessage` instead of raising (`loop.md` section 4)
    -- and neither deck reached the panel.
    """
    prior_limit = settings.handout_deck_max_bullet_chars
    prior_k, prior_n = agent.retrieve_k, agent.rerank_top_n
    settings.handout_deck_max_bullet_chars = IMPOSSIBLE_BULLET_LIMIT
    agent.retrieve_k, agent.rerank_top_n = 1, 1
    await db.commit()
    try:
        out, events, _convo = await ask(db, agent, user, S32_DECK_ASK)
        rows = (await db.scalars(
            select(Handout).where(Handout.query_id == out.query_id)
        )).all()
    finally:
        settings.handout_deck_max_bullet_chars = prior_limit
        agent.retrieve_k, agent.rerank_top_n = prior_k, prior_n
        await db.commit()

    name = "S32 an invalid chat deck is rejected, not stored"
    calls = [p for p in payloads_of(events, "TOOL_CALL") if p.get("tool") == "run_python"]
    errors = payloads_of(events, "TOOL_ERROR")
    results = payloads_of(events, "TOOL_RESULT")

    rejected = _pptx_names(errors, "rejected")
    # Anything the run wrote, from either payload. On a crash `artifacts` lists
    # what the program managed to save and those are NOT persisted, so this is
    # only ever used to decide whether a deck was ATTEMPTED -- the database
    # below is what decides whether one was kept.
    written = _pptx_names(results, "artifacts") + _pptx_names(errors, "artifacts")
    persisted = [r for r in rows if "presentationml" in (r.mime_type or "")]
    kinds = [p.get("error_kind") for p in errors]
    detail = (
        f"bullet_limit={IMPOSSIBLE_BULLET_LIMIT} run_python={len(calls)} "
        f"tool_errors={len(errors)} "
        f"error_kinds={kinds or 'none'} rejected_pptx={rejected or 'none'} "
        f"wrote_pptx={written or 'none'} persisted_handouts={len(rows)} "
        f"persisted_pptx={len(persisted)}"
    )

    if not calls:
        return unmeasured(name, f"{detail} -- the model never called run_python")
    if not rejected and not written and not persisted:
        return unmeasured(
            name, f"{detail} -- no .pptx was produced, so nothing was validated"
        )

    ok = (
        bool(rejected)
        and ARTEFACT_ERROR_KIND in kinds
        and not persisted
    )
    return Outcome(name, ok, detail)


async def s33_tools_off_makes_no_handout(db, agent, user) -> Outcome:
    """The same deck request with tools OFF is the classic path, and makes nothing.

    Feature 06, A8, and the standing form from `loop.md` S4: not "similar",
    identical. S1 asserts it for a plain question; this asserts it for the one
    request most likely to tempt a code path into running -- because the
    interesting regression is not "tools off produces no TOOL_CALL", it is
    "tools off produces no HANDOUT", and a `run_python` reached by any route
    other than `_tools_active(agent)` would show up here and nowhere else.

    Owns `tools_enabled` and restores it, the way S1 does.

    **Seen failing before it was believed** (`build.md` section 5). Same ask,
    2026-08-17, with the flag flipped for one run:

        tools off  [ok]   events=[GENERATE, RERANK, RETRIEVE, REWRITE,
                          SCORE_CHECK] handouts=0
        tools on   [FAIL] STRAY=[TOOL_CALL, TOOL_RESULT] handouts_by_query=1

    Both halves moved, which is what says the assertion is reading the flag and
    not the weather.
    """
    agent.tools_enabled = False
    await db.commit()
    try:
        out, events, convo = await ask(db, agent, user, S32_DECK_ASK)
        by_query = (await db.scalars(
            select(Handout).where(Handout.query_id == out.query_id)
        )).all()
        by_convo = (await db.scalars(
            select(Handout).where(Handout.conversation_id == convo.id)
        )).all()
    finally:
        agent.tools_enabled = True
        await db.commit()

    seen = set(event_types(events))
    stray = seen - CLASSIC_EVENTS
    ok = (
        not stray
        and "RETRIEVE" in seen
        and "GENERATE" in seen
        and "REWRITE" in seen
        and bool(out.answer)
        and not by_query
        and not by_convo
    )
    return Outcome(
        "S33 tools off makes no handout", ok,
        f"events={sorted(seen)}" + (f" STRAY={sorted(stray)}" if stray else "")
        + f" handouts_by_query={len(by_query)} handouts_by_conversation={len(by_convo)} "
        f"answer_chars={len(out.answer or '')}",
    )


# --------------------------------------------------------------------------
# HTTP scenarios -- the handout routes
#
# **These are ordinary scenarios now, and `--only` reaches them.** Until
# 2026-08-17 the whole block was one function gated on `if not only:` in
# `run_scenarios`, so every handout assertion was all-or-nothing with the
# twenty-minute suite: there was no way to run S8 without also re-asking twenty
# questions of a live model. `01-deck-harness-floor.md` section F names that as
# the reason it pushed everything it could down to `deck_check.py` -- and it is
# also why the criteria that CANNOT go to layer 1, the ones about a real job,
# kept not being executed. A test that is expensive to run is a test that is not
# run.
#
# Each function takes one `Http` and returns a LIST of outcomes, because some of
# them genuinely produce several: S8 makes four recipes and its download check
# reads what they left behind, and splitting those apart would mean making the
# four recipes twice.
# --------------------------------------------------------------------------


@dataclass
class Http:
    """What one handout-route scenario is handed.

    `agent` and `db` are the OUTER session's objects -- the same pair every
    scenario above owns and restores -- so a handout scenario starves retrieval
    exactly the way S3 and S13 do rather than inventing a second way. The
    background job re-loads the agent in a session of its own, so a committed
    write here is what that job reads.
    """

    client: object
    base: str
    agent_url: str
    db: object
    agent: Agent
    user: User


# THE LEVER S28 AND S32 USE TO FORCE A DECK TO BE INVALID, and it is not the
# one `02-artefact-validation.md` A10 and `06-tool-path-parity.md` A7 name.
#
# Both criteria say to raise `handout_deck_min_slides` "above what the fixture
# corpus can support". **Measured 2026-08-17, that lever does not work, and the
# reason is a property of the retry rather than of this fixture: the repair turn
# is TOLD the threshold.** `validate._check_deck` writes "has 6 slides, and the
# handout needs at least N. Add slides with `prs.slides.add_slide(...)`", and
# `_repair_message` hands that to the model verbatim. So:
#
#   floor = 12  ->  attempt 2 complies, the row ends `ready`, nothing measured
#   floor = 40  ->  attempt 2 tries, and the program is TRUNCATED at
#                   CODE_MAX_TOKENS: "Python syntax error on line 322:
#                   unterminated string literal". Observed, twice.
#
# The second row is `PLAN.md` R7 and `04-failure-legibility.md` D happening
# live -- nothing in the backend reads `finish_reason`, so a cut-off program
# fails as a syntax error and is reported as `error_kind="import"`. A slide
# floor is therefore either satisfiable (no failure to observe) or inflationary
# (a failure that is not the one under test). There is no third value.
#
# A bullet-length limit of ZERO has neither problem. It is un-satisfiable BY
# CONSTRUCTION rather than by luck -- `_check_deck` measures every paragraph
# including the title, and a slide with a title has a paragraph longer than
# zero characters -- and complying makes the program SHORTER, so it cannot
# truncate. Its message still opens "The file deck.pptx opened with 6 slides,
# but ...", so a failure is still legible as a statement about a deck.
#
# It costs nothing in discriminating power, which is the property that actually
# matters: a validator that had been deleted, or one that only ever fired on
# zero slides, ships the six-slide deck `ready` and S28 goes red.
IMPOSSIBLE_BULLET_LIMIT = 0

# What S31 shrinks `sandbox_timeout_s` to, so that both attempts time out and
# the row fails with `error_kind="timeout"` -- a real sandbox classification
# rather than a mocked one.
#
# **It was 1.0 and 1.0 DID NOT WORK, which is the reason this comment carries a
# number.** Measured on this box, 2026-08-17: a matplotlib bar chart through
# `sandbox.run` completes in **528 ms**, so a one-second ceiling let the run
# succeed and S31 reported `[warn] ... the run succeeded inside 1.0s`. The
# three-state result did its job -- a two-state harness would have called that
# green or red and either would have been a lie -- but the lever was
# machine-dependent, which is the thing to fix rather than to re-run.
#
# 50 ms is under Python's own interpreter start, so it cannot depend on how fast
# matplotlib imports here. That is the property wanted: this scenario is about
# `error_kind` surviving to the API, not about a realistic timeout.
S31_TIMEOUT_S = 0.05

# The five the sandbox mints, plus the one feature 02 adds. `error_kind` is a
# plain string with no enum behind it (PLAN.md 3.4), so this list is the only
# thing that would notice a sixth appearing unannounced.
SANDBOX_ERROR_KINDS = {"import", "syntax", "timeout", "runtime", "output"}
ARTEFACT_ERROR_KIND = "invalid"

# Words that make a failed deck's `error` a statement about MATERIAL or about
# the artefact, as opposed to about the sandbox falling over. S29 accepts a
# failure only when it is one of these; anything else means the run died for a
# reason that has nothing to do with what S29 is asking, and the honest report
# is `unmeasured`.
MATERIAL_SHAPED = ("slide", "material", "corpus", ".pptx", "no file", "upload a document")


async def make_handout(
    http: Http, recipe: str, brief: str, *, budget_s: float = 420
) -> tuple[str | None, dict | None, str]:
    """POST one handout and wait for a terminal row. Returns `(id, row, note)`.

    `row` is the LIST route's view of it -- the one the panel polls, and the one
    `error` and `error_kind` have to survive to (PLAN.md 3.4). `row` is None
    when the POST was refused or the job never settled, and `note` says which;
    every caller turns that into `unmeasured` rather than into a failure,
    because "the job did not finish" is not an answer to any question asked
    below.

    Polls rather than trusting the POST. Under `httpx.ASGITransport` a Starlette
    `BackgroundTask` runs inside the same ASGI call, so the 202 in fact does not
    return until the job is done -- which is convenient and is an implementation
    detail of the transport, not a property of the route. A scenario that
    assumed it would break silently against a real server.
    """
    r = await http.client.post(http.base, json={"recipe": recipe, "brief": brief})
    if r.status_code != 202:
        return None, None, f"POST returned {r.status_code}: {preview(r.text, 160)}"
    hid = r.json()["id"]

    deadline = time.monotonic() + budget_s
    while True:
        listing = await http.client.get(http.base, params={"limit": 200})
        if listing.status_code != 200:
            return hid, None, f"GET list returned {listing.status_code}"
        row = next((x for x in listing.json() if x["id"] == hid), None)
        if row is not None and row["status"] in ("ready", "failed"):
            return hid, row, ""
        if time.monotonic() >= deadline:
            status_now = row["status"] if row else "no row"
            return hid, None, f"still {status_now} after {budget_s:.0f} s"
        await asyncio.sleep(3)


async def handout_meta(handout_id: str) -> dict:
    """`handouts.meta` for one row, read in a session of its own.

    **`meta` is deliberately absent from every response model** -- putting it on
    the wire would hand a client the vocabulary `extra="forbid"` withholds
    (`api/handouts.py`) -- so the database is the only place `chunk_ids` can be
    read, and S30 needs it. Two facts are lifted out of `meta` onto the response
    and those two are read through the API, where they belong.

    Empty dict on anything unexpected: this feeds a detail line, never a
    verdict on its own.
    """
    async with SessionLocal() as db:
        row = await db.get(Handout, uuid.UUID(str(handout_id)))
        meta = getattr(row, "meta", None)
        return dict(meta) if isinstance(meta, dict) else {}


async def corpus_chunks(db, agent: Agent) -> int:
    """How many chunks the fixture corpus holds, read rather than assumed.

    S18 makes the same call for the same reason: a literal 7 here would turn a
    re-ingest at a different chunk size into a red row that names retrieval.
    """
    return await db.scalar(
        select(func.count(Chunk.id))
        .join(Document, Chunk.document_id == Document.id)
        .where(Document.agent_id == agent.id)
    ) or 0


async def s8_recipes(http: Http) -> list[Outcome]:
    """The four recipes, opened rather than weighed -- plus S8b and S8c.

    Until 2026-08-17 this was `status == "ready" and byte_size > 0`, which is
    satisfied by a zero-slide 27 KB deck and by 28 bytes of PK junk (both
    measured, `PLAN.md` 1.1). `06-test-plan.md:183` promised these four content
    checks since it was written and none of them had ever run.
    """
    outcomes: list[Outcome] = []
    created: dict[str, str] = {}
    for recipe, brief in (
        ("sheet", "A one-page study sheet on the power subsystem."),
        ("chart", "A bar chart of the power allocation by subsystem."),
        ("table", "A CSV of each subsystem and its power allocation in kW."),
        ("deck", "A short deck on the three downlink paths."),
    ):
        r = await http.client.post(http.base, json={"recipe": recipe, "brief": brief})
        if r.status_code == 202:
            created[recipe] = r.json()["id"]
        else:
            outcomes.append(Outcome(
                f"S8 recipe {recipe}", False,
                f"POST returned {r.status_code}: {preview(r.text, 200)}",
            ))

    # Poll. Recipes are LLM calls plus a subprocess; sheet is fastest.
    deadline = time.monotonic() + 480
    final: dict[str, dict] = {}
    while created and time.monotonic() < deadline:
        await asyncio.sleep(4)
        r = await http.client.get(http.base, params={"limit": 200})
        if r.status_code != 200:
            break
        rows = {row["id"]: row for row in r.json()}
        for recipe, hid in list(created.items()):
            row = rows.get(hid)
            if row and row["status"] in ("ready", "failed"):
                final[recipe] = row
                created.pop(recipe)

    deck_body: bytes | None = None
    deck_row: dict | None = None
    for recipe in ("sheet", "chart", "table", "deck"):
        row = final.get(recipe)
        if row is None:
            outcomes.append(Outcome(f"S8 recipe {recipe}", False, "did not finish in 480 s"))
            continue
        detail = (
            f"status={row['status']} bytes={row['byte_size']} "
            f"mime={row['mime_type']} err={preview(row.get('error') or '-', 90)}"
        )
        if row["status"] != "ready":
            outcomes.append(Outcome(f"S8 recipe {recipe}", False, detail))
            continue

        # THE FILE IS OPENED.
        body_r = await http.client.get(f"{http.base}/{row['id']}/download")
        if body_r.status_code != 200:
            outcomes.append(unmeasured(
                f"S8 recipe {recipe}",
                f"{detail} download={body_r.status_code}, so the bytes were never inspected",
            ))
            continue
        if recipe == "deck":
            deck_body, deck_row = body_r.content, row
        problem = artifact_problem(recipe, body_r.content)
        outcomes.append(Outcome(
            f"S8 recipe {recipe}", problem is None,
            detail if problem is None else f"{detail} -- {problem}",
        ))

    # Download, including the header-safety check.
    ready = [r for r in final.values() if r["status"] == "ready"]
    if not ready:
        # Was `if ready:` with no else, so when nothing reached `ready` this
        # appended NO Outcome at all -- not a pass, not a failure, not a
        # warning. It simply vanished from the summary count.
        outcomes.append(unmeasured(
            "S8b download + safe filename",
            "no recipe reached `ready`, so no download was attempted",
        ))
    else:
        hid = ready[0]["id"]
        r = await http.client.get(f"{http.base}/{hid}/download")
        disp = r.headers.get("content-disposition", "")
        safe = "\n" not in disp and "\r" not in disp and '"' in disp
        outcomes.append(Outcome(
            "S8b download + safe filename",
            r.status_code == 200 and len(r.content) > 0 and safe,
            f"status={r.status_code} bytes={len(r.content)} disp={ascii_safe(disp)}",
        ))

    # S8c -- feature 05's A9. The number on the CARD against the number in the
    # FILE.
    #
    # `deck_check.py` cases 50-56 assert `outline()` against fixtures it built
    # itself, which is the function working. This is the only assertion anywhere
    # that the right string reached the right column of a real row: a job that
    # wrote the model's stdout caption instead of the outline, or wrote attempt
    # 1's preview beside attempt 2's artefact, passes every layer-1 case and
    # puts a lie on the card -- and the measured failure this change set exists
    # for is precisely a caption reading "deck written with 6 slides" over a
    # file holding none.
    if deck_row is None or deck_body is None:
        outcomes.append(unmeasured(
            "S8c deck preview names the real slide count",
            "the deck recipe never reached `ready` and downloaded, so no preview was read",
        ))
    else:
        detail_r = await http.client.get(f"{http.base}/{deck_row['id']}")
        if detail_r.status_code != 200:
            outcomes.append(unmeasured(
                "S8c deck preview names the real slide count",
                f"GET detail returned {detail_r.status_code}",
            ))
        else:
            stored = detail_r.json().get("preview_text")
            claimed = preview_slide_count(stored)
            actual = slide_count(deck_body)
            outcomes.append(Outcome(
                "S8c deck preview names the real slide count",
                claimed is not None and claimed == actual,
                f"preview_claims={claimed} file_has={actual} "
                f"preview={preview(stored or '-', 90)}",
            ))

    return outcomes


async def s11_list_does_not_load_bytea(http: Http) -> list[Outcome]:
    """Listing must never select the bytea column."""
    import logging

    statements: list[str] = []

    class _Capture(logging.Handler):
        # `entry`, not `record` -- there is a module-level `record()` helper now
        # and a parameter of that name inside a handler is a trap for whoever
        # next wants to report something from in here.
        def emit(self, entry):
            statements.append(entry.getMessage())

    sql_log = logging.getLogger("sqlalchemy.engine.Engine")
    handler = _Capture()
    prior_level = sql_log.level
    sql_log.addHandler(handler)
    sql_log.setLevel(logging.INFO)
    try:
        await http.client.get(http.base, params={"limit": 200})
    finally:
        sql_log.removeHandler(handler)
        sql_log.setLevel(prior_level)

    leaked = [s for s in statements if "handouts.content" in s]
    if not statements:
        # `ok = not leaked` over an empty capture is green forever. If the
        # `sqlalchemy.engine.Engine` logger ever stops emitting -- a logging
        # config change, an echo default, a library rename -- this asserted
        # nothing and said so in green.
        return [unmeasured(
            "S11 list does not load bytea",
            "the SQL log captured 0 statements, so nothing was inspected",
        )]

    outcomes = [Outcome(
        "S11 list does not load bytea", not leaked,
        f"statements={len(statements)} leaked={len(leaked)}",
    )]

    # S11b. THE SAME QUESTION ASKED SO THAT DELETION CANNOT ANSWER IT.
    #
    # S11 above greps SQL for `handouts.content`. That is correct today and it is
    # unfalsifiable the moment the column is dropped: the string stops being
    # emittable, `leaked` is empty forever, and the row prints green while
    # measuring nothing. It would be the seventh entry in build.md section 7's
    # table and a NEW mechanism -- not an assertion that was too weak, but one
    # whose subject was deleted underneath it.
    #
    # It also has a nearer failure. Once the download route presigns, the
    # obvious-looking optimisation is to presign every row in the LIST response
    # so the panel need not call again. That emits no `handouts.content` SQL
    # whatsoever, so S11 passes -- on a change that hands a client up to 200 live
    # bearer capabilities in one body. S11 would go green on the defect it exists
    # to catch, inside this change set rather than a later one.
    #
    # So this asserts the positive property instead: listing performs no
    # object-storage work at all. Zero is itself a suspicious number, so the
    # control below proves the counter can move.
    calls: list[str] = []
    real_presign = storage_mod.presigned_get_url

    def _counting(*args, **kwargs):
        calls.append(args[0] if args else kwargs.get("key", "?"))
        return real_presign(*args, **kwargs)

    storage_mod.presigned_get_url = _counting
    try:
        await http.client.get(http.base, params={"limit": 200})
        during_list = len(calls)

        # The control. Signing is a local HMAC with no round trip, so this needs
        # no `ready` handout and no network -- which is what lets the control be
        # mandatory rather than conditional.
        try:
            _counting("agents/x/handouts/y.pptx", filename="y.pptx", mime_type=None)
            control_works = len(calls) == during_list + 1
        except Exception:  # noqa: BLE001
            control_works = False
    finally:
        storage_mod.presigned_get_url = real_presign

    if not storage_mod.enabled():
        outcomes.append(unmeasured(
            "S11b listing performs no object-storage work",
            "storage_route is not 'r2', so there is no object storage to call",
        ))
    elif not control_works:
        outcomes.append(unmeasured(
            "S11b listing performs no object-storage work",
            "the call counter did not register the control, so a zero here "
            "would mean nothing",
        ))
    else:
        outcomes.append(Outcome(
            "S11b listing performs no object-storage work",
            during_list == 0,
            f"presign calls during list={during_list} (control proved the "
            f"counter moves)",
        ))

    return outcomes


async def s17_generation_model_switchable(http: Http) -> list[Outcome]:
    """The model is switchable through the API, and a typo is not.

    S13 proves that CHANGING the model changes behaviour, but it writes the
    column directly, which was the only way to set it until 2026-08-16. This
    asserts the route: the field round-trips, null clears it, and a bare id is
    refused HERE rather than 404ing on every answer the agent later gives.

    The refusal is the half worth having. Accepting `"gemma-4-31b-it"` is easy
    and stores a value that makes the agent fail with `404 No endpoints
    found...`, which CLAUDE.md records as reading like an outage -- so the user
    changed a setting and the agent broke, with nothing connecting the two.
    """
    prior_model = (await http.client.get(http.agent_url)).json().get("generation_model")
    try:
        set_r = await http.client.patch(
            http.agent_url, json={"generation_model": "google/gemma-4-31b-it"}
        )
        bad_r = await http.client.patch(
            http.agent_url, json={"generation_model": "gemma-4-31b-it-typo"}
        )
        after_bad = (await http.client.get(http.agent_url)).json().get("generation_model")
        clear_r = await http.client.patch(http.agent_url, json={"generation_model": None})
        cleared = (await http.client.get(http.agent_url)).json().get("generation_model")
        ok = (
            set_r.status_code == 200
            and set_r.json().get("generation_model") == "google/gemma-4-31b-it"
            and bad_r.status_code == 422
            # The rejected write must not have landed.
            and after_bad == "google/gemma-4-31b-it"
            and clear_r.status_code == 200
            and cleared is None
        )
        return [Outcome(
            "S17 generation_model is switchable, typos are not", ok,
            f"set={set_r.status_code} bad={bad_r.status_code} "
            f"after_bad={after_bad!r} cleared={cleared!r}",
        )]
    finally:
        await http.client.patch(http.agent_url, json={"generation_model": prior_model})


async def s12_quota_refuses(http: Http) -> list[Outcome]:
    """Quota refuses rather than evicting."""
    agent_id = http.agent.id
    async with SessionLocal() as db:
        before = len((await db.scalars(
            select(Handout).where(Handout.agent_id == agent_id)
        )).all())
    original = settings.handout_max_per_agent
    settings.handout_max_per_agent = max(1, before)
    try:
        r = await http.client.post(http.base, json={"recipe": "sheet", "brief": "over quota"})
        async with SessionLocal() as db:
            after = len((await db.scalars(
                select(Handout).where(Handout.agent_id == agent_id)
            )).all())
        return [Outcome(
            "S12 quota refuses, evicts nothing",
            r.status_code == 409 and after == before,
            f"status={r.status_code} before={before} after={after}",
        )]
    finally:
        settings.handout_max_per_agent = original


async def s28_invalid_deck_fails_the_job(http: Http) -> list[Outcome]:
    """An unusable deck must FAIL, in a real job. Feature 02, A10.

    **This is the case that makes `validate.py` necessary, and without it every
    other assertion about the validator is over a function nothing is proven to
    call.** `deck_check.py` cases 20-28 pass `check()` a fixture and read what
    comes back; case 24 goes one step further and calls `_problem`. None of them
    can see the branch being reached from `run_handout_job` -- a
    `handout_validate_artifacts` read that got inverted, an exception swallowed
    between `_primary_artifact` and `_problem`, a `_problem` whose result the
    job stopped acting on -- and each of those leaves every layer-1 case green
    while shipping the defect this whole change set exists to fix. That is
    exactly how S3 went green twice: it tested a function and not a path.

    **The scenario OWNS its condition** (`loop.md` section 5). Hoping the
    fixture corpus is thin enough to produce a two-slide deck is the trap S3
    fell into twice, and `PLAN.md` 8.3 has already disproved it by measurement:
    at the most starved budget in this repo the model still produces five to
    nine slides. So a threshold moves instead -- see `IMPOSSIBLE_BULLET_LIMIT`
    for which one, and for the measurement that ruled out the slide floor A10
    actually names. Restored in a `finally`, and the setting is process-global:
    the job runs in this process, which is why writing it works and why leaking
    it would poison every row after.

    Four-state on purpose:

      failed + error_kind="invalid"  -> the branch fired in the job. PASS
      ready                          -> the branch did NOT fire. FAIL
      failed + another error_kind    -> the run died elsewhere. unmeasured
      no terminal row                -> the job never settled. unmeasured

    The third state is not hypothetical padding: measured on this fixture, a
    handout generation comes back with a mangled or truncated first program
    often enough that a two-state version of this scenario would be
    intermittently red for a reason that has nothing to do with validation --
    and CLAUDE.md records what an intermittent red does, which is get re-run
    until it passes.

    **Seen failing before it was believed** (`build.md` section 5, applied to
    the scenario rather than to the code). 2026-08-17, same brief, same lever,
    with `handout_validate_artifacts=False` added for one run:

        validation on   [ok]   status=failed error_kind=invalid attempts=2
        validation off  [FAIL] status=ready  error_kind=None    attempts=1 slides=7

    A seven-slide deck full of 107-character bullets, stored `ready`, is what
    ships when the third branch is not reached. That is the row this scenario
    exists to make impossible to miss.
    """
    prior = settings.handout_deck_max_bullet_chars
    settings.handout_deck_max_bullet_chars = IMPOSSIBLE_BULLET_LIMIT
    try:
        hid, row, note = await make_handout(
            http, "deck",
            "A deck covering the downlink paths, the power allocation and the "
            "battery store.",
        )
    finally:
        settings.handout_deck_max_bullet_chars = prior

    name = "S28 an invalid deck fails the job"
    if row is None:
        return [unmeasured(name, f"no terminal handout row: {note}")]

    error = row.get("error") or ""
    kind = row.get("error_kind")
    detail = (
        f"bullet_limit={IMPOSSIBLE_BULLET_LIMIT} status={row['status']} "
        f"error_kind={kind} attempts={row.get('attempts')} "
        f"names_the_deck={'slide' in error.lower()} err={preview(error, 110)}"
    )

    if row["status"] == "failed":
        if kind == ARTEFACT_ERROR_KIND:
            # The message is asserted as well as the kind. `error_kind` alone
            # would pass on a validator that returned a bare "invalid" -- and a
            # refusal the model cannot act on wastes the retry it triggers
            # (`loop.md` section 4), which is the whole reason `_check_deck`
            # writes a finding and a fix rather than a verdict.
            return [Outcome(name, "slide" in error.lower(), detail)]
        return [unmeasured(
            name,
            f"{detail} -- failed with a PROCESS error rather than an artefact "
            "one, so the validation branch was not the one exercised",
        )]

    body_r = await http.client.get(f"{http.base}/{hid}/download")
    slides = slide_count(body_r.content) if body_r.status_code == 200 else None
    return [Outcome(
        name, False,
        f"{detail} slides={slides} -- a deck violating the limit was stored "
        "`ready`; `_problem`'s third branch did not fire in the job",
    )]


async def s29_deck_survives_a_starved_agent(http: Http) -> list[Outcome]:
    """A starved AGENT must not yield a silently-thin `ready` deck. Feature 03, A4.

    **The scenario owns the starvation**, `loop.md` section 5: `retrieve_k=1,
    rerank_top_n=1` is the smallest budget the agent row can express, and it is
    written here rather than inherited from the fixture, exactly the way S3, S13,
    S15, S18 and S27 own the same two columns.

    What makes this interesting is that the agent's budget is the one thing the
    deck job does NOT use: `RECIPES["deck"]` carries `retrieve_k=40,
    rerank_top_n=10` and `gather_material` passes them through as per-call
    overrides. So the honest reading of a pass here is "the recipe's budget
    survived an agent configured to starve it", and the honest reading of a
    silently-thin `ready` deck is that the override was dropped somewhere
    between the recipe and `aretrieve`.

    **Read the pairing with S30 rather than this row alone.** A deck can clear
    the floor because the corpus is easy, and `PLAN.md` 8.3 measured five to
    nine slides at the most starved budget in the repo -- so this assertion
    alone is satisfiable by an override that never took effect. S30 is what
    measures whether it did. The two are a pair for the same reason
    `route_specialist_check.py` cases 25 and 26 are.

    The failure branch is accepted only when the error is about material or the
    artefact. A timeout or an import refusal is the run dying for an unrelated
    reason, and calling that a pass would make this row green on a broken
    sandbox.
    """
    agent = http.agent
    prior_k, prior_n = agent.retrieve_k, agent.rerank_top_n
    agent.retrieve_k, agent.rerank_top_n = 1, 1
    await http.db.commit()
    try:
        hid, row, note = await make_handout(
            http, "deck", "A deck on the three downlink paths and what each carries."
        )
    finally:
        agent.retrieve_k, agent.rerank_top_n = prior_k, prior_n
        await http.db.commit()

    name = "S29 a starved agent yields no thin deck"
    if row is None:
        return [unmeasured(name, f"no terminal handout row: {note}")]

    floor = settings.handout_deck_min_slides
    meta = await handout_meta(hid)
    chunks = len(meta.get("chunk_ids") or [])
    error = row.get("error") or ""
    detail = (
        f"agent_budget=1/1 recipe_chunks={chunks or '-'} floor={floor} "
        f"status={row['status']} error_kind={row.get('error_kind')} "
        f"err={preview(error, 100)}"
    )

    if row["status"] == "failed":
        if any(word in error.lower() for word in MATERIAL_SHAPED):
            return [Outcome(name, True, f"{detail} -- failed, and said why")]
        return [unmeasured(
            name,
            f"{detail} -- failed for a reason unrelated to material, so the "
            "property was not exercised",
        )]

    body_r = await http.client.get(f"{http.base}/{hid}/download")
    if body_r.status_code != 200:
        return [unmeasured(name, f"{detail} download={body_r.status_code}")]
    slides = slide_count(body_r.content)
    if slides is None:
        return [Outcome(
            name, False, f"{detail} -- stored `ready` and does not open at all"
        )]
    return [Outcome(name, slides >= floor, f"{detail} slides={slides}")]


async def s30_deck_uses_the_recipe_budget(http: Http) -> list[Outcome]:
    """The deck retrieved at the RECIPE's width, not the agent's. Feature 03, A5.

    **A4 alone is satisfiable by an override that never took effect** -- 03's
    own acceptance table says so -- so this is the row that measures whether
    `retrieve_k=40, rerank_top_n=10` actually reached `aretrieve`.

    **HOW IT LOOKS, and what that CANNOT see.** A5 as written asks for "the
    RETRIEVE trace payload for a deck job". There is no such payload and there
    cannot be one: `jobs.py` never imports `TraceRecorder`, a recipe handout
    carries `query_id = NULL` by construction, and `PLAN.md` 3.5 rules a trace
    for background jobs out of this change set as needing a nullable anchor on
    `trace_events` plus a new view. So the observation used here is
    `meta["chunk_ids"]`, which `gather_material` writes from
    `retrieval.documents` -- the set that actually reached the prompt.

    What it sees: the WIDTH of the material, end to end, through the real job.
    On this fixture, with the agent owned down to `retrieve_k=1,
    rerank_top_n=1` and a seven-chunk corpus, the three possibilities separate
    cleanly:

        both overrides honoured   -> min(40, 7) reranked to min(10, 7) = 7
        only `top_n` honoured     -> the agent's k=1 caps it at            1
        neither honoured          ->                                       1

    What it CANNOT see: `retrieve_k` apart from `rerank_top_n` above the corpus
    size. Forty and eight are indistinguishable against seven chunks, and no
    fixture this size can tell them apart -- that is a property of the corpus,
    not a weakness of the observation, and it is why the assertion is written
    against the corpus count rather than against 40.

    **The `table` control is not decoration.** Without it, seven chunk ids would
    also be produced by an agent whose starvation never landed -- a write that
    did not commit, a job reading a cached row. `table` carries `retrieve_k =
    None`, the identity case, so under the identical agent it must come back at
    1. One variable between the two POSTs: the recipe.

    **Seen failing before it was believed.** `gather_material` called directly
    against this fixture at `agent = 1/1`, 2026-08-17, with no model call:

        shipped recipe (40/10)         chunk_ids = 7
        override REMOVED (None/None)   chunk_ids = 1

    So the assertion below separates the shipped feature from its absence by
    six chunks, not by a rounding error.
    """
    from app.handouts.recipes import RECIPES

    deck = RECIPES["deck"]
    agent = http.agent
    corpus = await corpus_chunks(http.db, agent)
    prior_k, prior_n = agent.retrieve_k, agent.rerank_top_n
    agent.retrieve_k, agent.rerank_top_n = 1, 1
    await http.db.commit()
    try:
        deck_id, deck_row, deck_note = await make_handout(
            http, "deck", "A deck on the battery store and the power deficit it covers."
        )
        table_id, table_row, table_note = await make_handout(
            http, "table", "A CSV of each downlink path and its data rate."
        )
    finally:
        agent.retrieve_k, agent.rerank_top_n = prior_k, prior_n
        await http.db.commit()

    name = "S30 the deck retrieves at the recipe's budget"
    deck_meta = await handout_meta(deck_id) if deck_id else {}
    table_meta = await handout_meta(table_id) if table_id else {}
    deck_chunks = deck_meta.get("chunk_ids")
    table_chunks = table_meta.get("chunk_ids")

    # `chunk_ids` is written by the SUCCESS path only; `_settle` records just
    # `error_kind` and `attempts` on a failed row. So a failed job leaves this
    # question genuinely unanswered rather than answered "zero".
    if not isinstance(deck_chunks, list) or not isinstance(table_chunks, list):
        return [unmeasured(
            name,
            f"deck={deck_row['status'] if deck_row else deck_note} "
            f"table={table_row['status'] if table_row else table_note} -- "
            "`meta[\"chunk_ids\"]` is written on success only, so the widths "
            "were never recorded",
        )]

    expected_deck = min(corpus, deck.rerank_top_n or corpus)
    detail = (
        f"corpus={corpus} agent_budget=1/1 recipe_budget="
        f"{deck.retrieve_k}/{deck.rerank_top_n} deck_chunks={len(deck_chunks)} "
        f"(expected {expected_deck}) table_chunks={len(table_chunks)} "
        f"(expected 1) rerank_enabled={agent.rerank_enabled}"
    )
    ok = (
        len(deck_chunks) == expected_deck
        and len(table_chunks) == 1
        and len(deck_chunks) > len(table_chunks)
    )
    return [Outcome(name, ok, detail)]


async def s31_failure_kind_reaches_the_api(http: Http) -> list[Outcome]:
    """A failed handout carries `error` AND `error_kind` on the list route. Feature 04, A8.

    `HandoutOut.error_kind` is not an attribute of `Handout` -- it is lifted out
    of the `meta` JSONB by a wrap validator -- and `PLAN.md` 3.4 records that
    the plan very nearly shipped it as a bare field, which would have serialised
    `None` forever with no error, no warning and valid JSON. `deck_check.py`
    cannot see that: the defect is in Pydantic serialisation of an ORM row, and
    only a real failed row fetched through the real route proves the lift
    happened.

    **The lever is NOT the one feature 04 names, and that is a defect in the
    criterion rather than a shortcut here.** A8 says to reach the
    `Material.is_empty` refusal by POSTing against an agent with no corpus. That
    path raises a bare `ValueError`, and `run_handout_job` records a kind only
    for `HandoutFailure` -- deliberately, with a comment saying that a job which
    never reached the sandbox has no kind to record. Measured 2026-08-17 by
    driving `run_handout_job` into that same branch directly:

        status         failed
        error          "Unknown recipe 'no-such-recipe'"
        meta           {}
        API error_kind None      API attempts None

    So written as specified this scenario would assert a value the shipped code
    is designed not to write, and would be red forever. A red row that means
    "the criterion was wrong" teaches its reader to ignore red.

    So the failure is forced through the sandbox instead, by shrinking
    `sandbox_timeout_s` for the duration -- owned and restored, and process-
    global for `IMPOSSIBLE_BULLET_LIMIT`'s reason. Both attempts time out, the run
    raises `HandoutFailure(error_kind="timeout")`, and the assertion is that
    both strings survive to the panel's own polling endpoint. `chart` rather
    than `deck` because it is the shortest of the three sandbox programs and
    this scenario is paying for two generations it intends to throw away.
    """
    prior = settings.sandbox_timeout_s
    settings.sandbox_timeout_s = S31_TIMEOUT_S
    try:
        _hid, row, note = await make_handout(
            http, "chart", "A bar chart of the three downlink data rates."
        )
    finally:
        settings.sandbox_timeout_s = prior

    name = "S31 a failed handout carries error and error_kind"
    if row is None:
        return [unmeasured(name, f"no terminal handout row: {note}")]

    error = row.get("error") or ""
    kind = row.get("error_kind")
    detail = (
        f"timeout={S31_TIMEOUT_S}s status={row['status']} error_kind={kind} "
        f"attempts={row.get('attempts')} err={preview(error, 100)}"
    )
    if row["status"] != "failed":
        return [unmeasured(
            name,
            f"{detail} -- the run succeeded inside {S31_TIMEOUT_S}s, so no "
            "failure was produced to describe",
        )]

    known = SANDBOX_ERROR_KINDS | {ARTEFACT_ERROR_KIND}
    ok = bool(error.strip()) and isinstance(kind, str) and kind in known
    if ok and kind != "timeout":
        # Still a pass -- both fields arrived, which is what A8 asks -- but the
        # lever missed, and a reader should see that rather than infer it.
        detail = f"{detail} <- expected `timeout`; the lever hit a different failure"
    return [Outcome(name, ok, detail)]


HTTP_SCENARIOS = [
    s8_recipes,
    s11_list_does_not_load_bytea,
    s17_generation_model_switchable,
    s12_quota_refuses,
    # The four owed by change set 12. Appended rather than interleaved, so a
    # default run prints the first four in the order it always has.
    #
    # S28 before S29/S30 for the S16-before-S15 reason: it is the one that says
    # whether the validator is reached at all, and a red S28 explains a thin
    # deck under S29 without anyone re-deriving a retrieval budget.
    s28_invalid_deck_fails_the_job,
    s29_deck_survives_a_starved_agent,
    s30_deck_uses_the_recipe_budget,
    s31_failure_kind_reaches_the_api,
]


async def http_scenarios(db, agent: Agent, user: User, only: str | None) -> list[Outcome]:
    """Run the selected handout-route scenarios through an ASGI transport.

    `current_user` is overridden; `owned_agent` is NOT. That split is deliberate:
    the identity assertion is the only thing a script cannot perform (it needs a
    human at a Google consent screen), while the ownership hop is exactly the
    thing worth testing, so it runs for real against the database.

    Returns `[]` without starting the app when `--only` selects none of them, so
    a targeted run of a chat scenario pays nothing for this block.
    """
    selected = [fn for fn in HTTP_SCENARIOS if not only or only in fn.__name__]
    if not selected:
        return []

    import httpx

    from app.auth.deps import current_user
    from app.db.session import SessionLocal as SL
    from app.main import app

    user_id = user.id

    async def _fake_user():
        async with SL() as db_:
            return await db_.get(User, user_id)

    app.dependency_overrides[current_user] = _fake_user
    outcomes: list[Outcome] = []

    rule("Handout routes")
    try:
        # TWO changes, and neither works without the other.
        #
        # `follow_redirects=True`, because the download route answers 302 to a
        # presigned URL on the R2 road. Nine assertions in this file read
        # `.content` off that route -- S8 x4, S8b, S8c, S28, S29 -- and without
        # this every one of them goes red simultaneously, reading as "the
        # migration broke downloads" rather than "the client does not follow
        # redirects".
        #
        # `mounts`, because a client's explicit `transport=` serves EVERY url,
        # not only those under `base_url`. Follow the redirect with the ASGI
        # transport still in charge and the R2 url is re-issued against the
        # FastAPI app, which has no such route -- so the nine assertions go from
        # reading an empty redirect body to reading a 404 body. Same red, and the
        # second cause reads as an application routing bug.
        transport = httpx.ASGITransport(app=app)
        async with httpx.AsyncClient(
            transport=transport,
            mounts={"https://": httpx.AsyncHTTPTransport()},
            follow_redirects=True,
            base_url="http://test",
            timeout=900,
        ) as client:
            http = Http(
                client=client,
                base=f"/api/agents/{agent.id}/handouts",
                agent_url=f"/api/agents/{agent.id}",
                db=db,
                agent=agent,
                user=user,
            )
            for fn in selected:
                started = time.perf_counter()
                try:
                    produced = await fn(http)
                except Exception as exc:  # one scenario must never abort the block
                    text = f"{type(exc).__name__}: {str(exc)[:300]}"
                    produced = [Outcome(
                        fn.__name__, False, ascii_safe(text)[:220],
                        rate_limited=is_rate_limited(text),
                    )]
                took = time.perf_counter() - started
                # The elapsed time is on the LAST row of the group, because it
                # is the group that was timed: `s8_recipes` makes four handouts
                # and stamping the same number on each would read as four jobs
                # of that length.
                for index, outcome in enumerate(produced, start=1):
                    record(
                        outcome, outcomes,
                        took if index == len(produced) else None,
                    )
    finally:
        app.dependency_overrides.pop(current_user, None)

    return outcomes


# --------------------------------------------------------------------------
# Runner
# --------------------------------------------------------------------------

SCENARIOS = [
    s1_classic_path,
    s2_no_reflex_tools,
    s3_multi_hop,
    s4_chart_from_chat,
    s5_self_correction,
    s6_step_budget,
    s7_refusal_survives,
    s9_timing_adds_up,
    s10_citation_integrity,
    # S16 first of the swap group: it is the only structural one, costs no model
    # call, and names the configuration the three behavioural ones depend on. A
    # red S16 explains a red S15; the reverse ordering makes the reader debug the
    # model instead of the config.
    s16_tool_use_has_a_belt_and_braces,
    s15_model_initiates_search,
    s14_no_redundant_gap_search,
    s13_gap_trigger_still_fires,
    # S19 before S18, for the S16-before-S15 reason: S19 is the cheaper one and
    # it names the failure -- a red S19 says the rewriter is dead or mangling,
    # which explains a red S18 without anyone re-deriving a retrieval probe.
    s19_no_over_firing,
    s18_rewrite_is_necessary,
    # S20-S27, the orchestrator. Ordered cheapest-and-most-explanatory first,
    # the S16-before-S15 convention:
    #
    #   S20 is the regression test and costs one plain turn. A red S20 means the
    #       classic path moved, which explains every row under it.
    #   S24 makes no model call at all, so a red S24 is the detector rather than
    #       the model and is worth knowing before reading S25-S27.
    #   S21 must precede S22, which is not a preference -- S22 reads what S21
    #       measured, and without it can only report `s21_not_run`.
    #   S26 is the one to write first and the one to read first among the
    #       behavioural rows: it is the only scenario here that fails by making
    #       the product quietly worse rather than visibly broken.
    s20_orchestration_off,
    s24_phantom_marker_is_caught,
    s21_router_discriminates,
    s22_mention_overrides_router,
    s23_two_mentions_two_sections,
    s25_no_self_check_on_a_socratic_turn,
    s26_critic_exempts_pedagogy,
    s27_budget_exhaustion_still_answers,
    # S32-S33, the tool path's half of change set 12. S33 first, for the
    # S16-before-S15 reason: it is the regression twin and it costs one plain
    # turn, so a red S33 means the classic path moved and explains S32 without
    # anyone reading a sandbox trace.
    s33_tools_off_makes_no_handout,
    s32_tool_deck_is_rejected,
]


async def run_scenarios(db, only: str | None) -> list[Outcome]:
    user, agent = await get_or_create_agent(db)

    doc_count = len((await db.scalars(
        select(Document).where(Document.agent_id == agent.id)
    )).all())
    if doc_count == 0:
        print("ERROR: no corpus. Run with --setup first.")
        return [Outcome("setup", False, "no documents ingested")]

    rule("Configuration")
    print(f"  generation   : {settings.generation_model}")
    print(f"  tools        : global={settings.agent_tools_enabled} "
          f"agent={agent.tools_enabled} max_steps={agent.max_tool_steps}")
    print(f"  sandbox      : timeout={settings.sandbox_timeout_s}s "
          f"mem={settings.sandbox_memory_mb}MB")
    print(f"  namespace    : {agent.namespace}  ({doc_count} documents)")

    outcomes: list[Outcome] = []
    selected = [fn for fn in SCENARIOS if not only or only in fn.__name__]
    if selected:
        rule("Scenarios")
    for fn in selected:
        name = fn.__name__
        started = time.perf_counter()
        try:
            outcome = await fn(db, agent, user)
        except Exception as exc:  # a scenario must never abort the suite
            text = f"{type(exc).__name__}: {str(exc)[:300]}"
            outcome = Outcome(
                name, False, ascii_safe(text)[:220], rate_limited=is_rate_limited(text)
            )
        record(outcome, outcomes, time.perf_counter() - started)

    # **NOT gated on `if not only` any more, and that gate is why five handout
    # criteria kept not being executed.** `--only` could never reach this block,
    # so every assertion about a real handout job was all-or-nothing with the
    # twenty-minute suite -- and a check that is expensive to run is a check
    # that does not get run. `http_scenarios` applies the same substring match
    # to its own function names and returns `[]` without starting the app when
    # none of them is selected, so a targeted chat run costs exactly what it did
    # before and a default run is unchanged.
    outcomes.extend(await http_scenarios(db, agent, user, only))

    if only and not outcomes:
        print(f"\n  --only {only!r} matched no scenario.")

    return outcomes


async def main_async(args) -> int:
    try:
        async with SessionLocal() as db:
            if args.cleanup:
                return await cleanup(db)
            if args.setup:
                return await setup(db)

            outcomes = await run_scenarios(db, args.only)

        rule("Summary")
        failed = [o for o in outcomes if _is_failure(o)]
        rated = [o for o in outcomes if o.rate_limited and not o.ok]
        unrun = [o for o in outcomes if o.unmeasured]
        passed = [o for o in outcomes if o.ok]
        for o in outcomes:
            # `_flag`, not a second copy of it. The summary used to spell the
            # two-state version out inline, so a `[rate]` row printed as [FAIL]
            # here after printing as [rate] above -- the one line a reader
            # actually screenshots disagreed with the run.
            print(f"  {_flag(o)} {o.name}")
        print(f"\n  {len(passed)} / {len(outcomes)} passed")

        # Printed even on a green run, and that is the whole point: an unmeasured
        # check is invisible precisely when everything looks fine.
        if unrun:
            print(f"\n  {len(unrun)} NOT MEASURED -- treat as unknown, never as passing:")
            for o in unrun:
                print(f"    [warn] {o.name}  {o.detail}")
        if rated:
            print(f"\n  {len(rated)} refused upstream -- re-run in a minute:")
            for o in rated:
                print(f"    [rate] {o.name}  {o.detail}")
        return 1 if failed else 0
    finally:
        await engine.dispose()


def main() -> int:
    p = argparse.ArgumentParser(description="End-to-end check for the agent loop and Handouts.")
    p.add_argument("--setup", action="store_true", help="Create the agent and ingest fixtures")
    p.add_argument(
        "--run",
        action="store_true",
        help="Run the scenarios (the default when no other mode is given)",
    )
    p.add_argument("--cleanup", action="store_true", help="Delete namespace, agent and rows")
    p.add_argument(
        "--only",
        default=None,
        help=(
            "Substring match on a scenario function name, INCLUDING the handout"
            " routes (e.g. --only s28, --only s8_recipes). Match on the trailing"
            " underscore -- `s3_` -- to separate S3 from S30-S33."
        ),
    )
    args = p.parse_args()
    return asyncio.run(main_async(args))


if __name__ == "__main__":
    raise SystemExit(main())
